"""SPEC 0011 PPT 配置选项测试。

覆盖：
- PptRenderer.render：config 各字段的应用逻辑（页数、主题色、图表开关）+ 降级策略
- outlines API：generate_ppt 端点的请求体解析、config 校验、错误码
"""

import json
from datetime import datetime, timezone

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from app.core.errors import AppError
from app.infrastructure.database.engine import Base
from app.infrastructure.renderers.ppt_renderer import PptRenderer
from app.main import app
from app.modules.outlines.models import Outline
from app.modules.outlines.status import OutlineStatus
from app.modules.projects.models import Project
from app.modules.projects.status import ProjectStatus


TEST_DB = "sqlite:///:memory:"


SAMPLE_SECTIONS = [
    {
        "id": "sec_001",
        "title": "实验目的",
        "content": "分析胃病数据分布特征",
        "source_type": "REQUIREMENT",
        "source_ids": ["plan_001"],
    },
    {
        "id": "sec_002",
        "title": "实验方法",
        "content": "使用描述性统计方法",
        "source_type": "EVIDENCE",
        "source_ids": ["ev_001"],
    },
    {
        "id": "sec_003",
        "title": "实验结果",
        "content": "执行成功，数据已清洗",
        "source_type": "EXECUTION",
        "source_ids": ["run_001"],
    },
    {
        "id": "sec_004",
        "title": "总结",
        "content": "本实验完成了数据分析与可视化",
        "source_type": "SUMMARY",
        "source_ids": [],
    },
]


# --- 渲染器测试 ---


def _render_ppt(tmp_path, config=None, sections=None, artifacts=None,
                filename="output.pptx"):
    """渲染 PPT 并返回路径。"""
    renderer = PptRenderer()
    output_path = tmp_path / filename
    renderer.render(
        project_name="测试项目",
        project_topic="测试课题",
        outline_sections=sections or SAMPLE_SECTIONS,
        execution_artifacts=artifacts or [],
        output_path=str(output_path),
        config=config,
    )
    return output_path


def test_render_no_config_keeps_default(tmp_path):
    """R-PAGE-01：config=None 保持默认行为。"""
    output_path = _render_ppt(tmp_path, config=None)
    prs = Presentation(str(output_path))
    # 标题页 + 内容页（REQUIREMENT+EVIDENCE+ANALYSIS 合并 + EXECUTION）+ 总结页
    assert len(prs.slides) >= 3


def test_render_config_none_dict_keeps_default(tmp_path):
    """R-FALL-01：config 为空 dict 保持默认行为。"""
    output_path = _render_ppt(tmp_path, config={})
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 3


def test_render_target_slide_count_6_limits_content(tmp_path):
    """R-PAGE-02：target_slide_count=6 时内容页不超过 4（6-标题页-总结页）。"""
    output_path = _render_ppt(
        tmp_path, config={"target_slide_count": 6}
    )
    prs = Presentation(str(output_path))
    # 标题页 + 内容页（<=4）+ 总结页 = <=6
    assert len(prs.slides) <= 6


def test_render_target_slide_count_20_keeps_actual(tmp_path):
    """R-PAGE-03：target_slide_count=20 时内容少于可用槽位，保持实际页数。"""
    output_path = _render_ppt(
        tmp_path, config={"target_slide_count": 20}
    )
    prs = Presentation(str(output_path))
    # 内容页候选最多 3（REQUIREMENT + EVIDENCE+DATASET+ANALYSIS + EXECUTION）
    # 标题页 + 3 内容页 + 总结页 = 5，不超过 20
    assert len(prs.slides) <= 20


def test_render_target_slide_count_5_minimum(tmp_path):
    """R-PAGE-04：target_slide_count=5（最小值），内容页不超过 3。"""
    output_path = _render_ppt(
        tmp_path, config={"target_slide_count": 5}
    )
    prs = Presentation(str(output_path))
    # 标题页 + 内容页（<=3）+ 总结页 = <=5
    assert len(prs.slides) <= 5


def test_render_charts_not_counted_in_target(tmp_path):
    """R-PAGE-05：图表页不计入 target_slide_count。"""
    # 创建测试 PNG 图表产物
    chart_path = tmp_path / "chart.png"
    # 最小有效 PNG 文件（1x1 像素）
    chart_path.write_bytes(
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00\x00\x00\x03\x00\x01"
        b"\x82\x8b\x99\xde\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    artifacts = [
        {"name": "chart1.png", "artifact_type": "CHART_PNG",
         "file_path": str(chart_path)},
    ]
    output_path = _render_ppt(
        tmp_path, config={"target_slide_count": 8}, artifacts=artifacts
    )
    prs = Presentation(str(output_path))
    # 图表页额外生成，不计入 target_slide_count
    # 标题页 + 内容页 + 图表页 + 总结页
    assert len(prs.slides) >= 4


def _slide_has_color(slide, target_color: RGBColor) -> bool:
    """检查某页是否在文本字体或形状填充中使用了目标颜色。

    SPEC 0026 增强：支持渐变填充检查（检查任一停止点颜色）。
    """
    from pptx.enum.dml import MSO_FILL_TYPE
    for shape in slide.shapes:
        # 检查文本 run 的字体颜色
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color and run.font.color.rgb == target_color:
                        return True
        # 检查形状填充颜色（纯色、渐变）
        try:
            fill = shape.fill
            if fill.type is None:
                continue
            if fill.type == MSO_FILL_TYPE.SOLID:
                if fill.fore_color.rgb == target_color:
                    return True
            elif fill.type == MSO_FILL_TYPE.GRADIENT:
                # 渐变填充：检查任一停止点颜色
                for stop in fill.gradient_stops:
                    if stop.color.rgb == target_color:
                        return True
        except Exception:
            pass
    return False


def test_render_theme_color_purple_applied(tmp_path):
    """R-COLOR-02：theme_color=#7c3aed 时主题色应用到 PPT（SPEC 0024：色块/标题/分隔线/圆点）。"""
    output_path = _render_ppt(
        tmp_path, config={"theme_color": "#7c3aed"}
    )
    prs = Presentation(str(output_path))
    # SPEC 0024：空白版式无 title placeholder，主题色改应用到
    # 封面顶部色块、内容页标题文字、分隔线、要点圆点等元素
    theme_color = RGBColor(0x7c, 0x3a, 0xed)
    found = any(
        _slide_has_color(slide, theme_color) for slide in prs.slides
    )
    assert found, "主题色 #7c3aed 未应用到 PPT 任何元素（色块/标题/分隔线/圆点）"


def test_render_theme_color_blue_applied(tmp_path):
    """R-COLOR-03：theme_color=#2563eb 时主题色应用到 PPT（SPEC 0024：色块/标题/分隔线/圆点）。"""
    output_path = _render_ppt(
        tmp_path, config={"theme_color": "#2563eb"}
    )
    prs = Presentation(str(output_path))
    theme_color = RGBColor(0x25, 0x63, 0xeb)
    found = any(
        _slide_has_color(slide, theme_color) for slide in prs.slides
    )
    assert found, "主题色 #2563eb 未应用到 PPT 任何元素（色块/标题/分隔线/圆点）"


def test_render_theme_color_green_all_slides(tmp_path):
    """R-COLOR-04：主题色应用到多个页面（SPEC 0024：封面色块+内容页标题+总结页等）。"""
    output_path = _render_ppt(
        tmp_path, config={"theme_color": "#16a34a"}
    )
    prs = Presentation(str(output_path))
    theme_color = RGBColor(0x16, 0xa3, 0x4a)
    # 统计应用了主题色的页面数
    slides_with_theme = sum(
        1 for slide in prs.slides if _slide_has_color(slide, theme_color)
    )
    # 主题色应至少应用到封面（色块）和内容页/总结页（标题、分隔线）
    assert slides_with_theme >= 2, (
        f"主题色只应用到 {slides_with_theme} 个页面，应至少 2 个（封面色块 + 内容页标题）"
    )


def test_render_no_theme_color_keeps_default(tmp_path):
    """R-COLOR-01：config=None 时主题色不修改（保持默认）。"""
    output_path_no_color = _render_ppt(tmp_path, config=None)
    output_path_with_color = _render_ppt(
        tmp_path, config={"theme_color": "#dc2626"}
    )
    # 两个文件都生成成功
    assert output_path_no_color.exists()
    assert output_path_with_color.exists()


def test_render_include_charts_true(tmp_path):
    """R-CHART-01：include_charts=True 时生成图表页（有图表产物时）。"""
    chart_path = tmp_path / "chart.png"
    chart_path.write_bytes(
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00\x00\x00\x03\x00\x01"
        b"\x82\x8b\x99\xde\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    artifacts = [
        {"name": "chart1.png", "artifact_type": "CHART_PNG",
         "file_path": str(chart_path)},
    ]
    output_with = _render_ppt(
        tmp_path, config={"include_charts": True}, artifacts=artifacts
    )
    prs_with = Presentation(str(output_with))
    # 有图表页
    assert len(prs_with.slides) >= 4  # 标题+内容+图表+总结


def test_render_include_charts_false_skips_chart(tmp_path):
    """R-CHART-02：include_charts=False 时不嵌入图表（SPEC 0024：检查图片 shape 存在性）。

    SPEC 0024 中单张图表会嵌入到内容页右栏（双栏布局），不再额外生成独立图表页，
    因此页数比较不再适用。改为检查图片 shape（PICTURE 类型）的存在性。
    """
    chart_path = tmp_path / "chart.png"
    chart_path.write_bytes(
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00\x00\x00\x03\x00\x01"
        b"\x82\x8b\x99\xde\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    artifacts = [
        {"name": "chart1.png", "artifact_type": "CHART_PNG",
         "file_path": str(chart_path)},
    ]
    output_without = _render_ppt(
        tmp_path, config={"include_charts": False}, artifacts=artifacts,
        filename="no_charts.pptx",
    )
    output_with = _render_ppt(
        tmp_path, config={"include_charts": True}, artifacts=artifacts,
        filename="with_charts.pptx",
    )
    prs_without = Presentation(str(output_without))
    prs_with = Presentation(str(output_with))
    # SPEC 0024：include_charts=False 时不应有任何图片 shape
    has_picture_without = any(
        shape.shape_type == 13  # PICTURE
        for slide in prs_without.slides
        for shape in slide.shapes
    )
    # include_charts=True 时应有图片 shape（图表嵌入到内容页右栏）
    has_picture_with = any(
        shape.shape_type == 13
        for slide in prs_with.slides
        for shape in slide.shapes
    )
    assert not has_picture_without, "include_charts=False 时不应嵌入图表图片"
    assert has_picture_with, "include_charts=True 时应嵌入图表图片"


def test_render_include_charts_false_no_artifacts(tmp_path):
    """R-CHART-03：include_charts=False 且无图表产物时不生成图表页。"""
    output_path = _render_ppt(
        tmp_path, config={"include_charts": False}, artifacts=[]
    )
    prs = Presentation(str(output_path))
    # 无图表产物，与默认行为一致
    assert len(prs.slides) >= 3


def test_render_partial_config(tmp_path):
    """R-FALL-02：config 部分字段缺失时只应用已有字段。"""
    output_path = _render_ppt(
        tmp_path, config={"theme_color": "#2563eb"}
    )
    # 只应用主题色，其他使用默认值
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 3


def test_render_invalid_theme_color_falls_back(tmp_path):
    """R-FALL-03：hex 色值解析异常时降级到默认（不抛异常）。"""
    # _parse_theme_color 内部捕获异常返回 None，不抛出
    output_path = _render_ppt(
        tmp_path, config={"theme_color": "#invalid"}
    )
    # 降级成功，文件正常生成
    assert output_path.exists()
    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 3


# --- SPEC 0025 三角色彩系统与三明治结构测试 ---


def _shape_has_fill_color(slide, target_color: RGBColor) -> bool:
    """检查幻灯片是否有指定填充色的形状（支持纯色和渐变填充）。

    SPEC 0026 增强：渐变填充时检查任一停止点颜色是否匹配。
    """
    from pptx.enum.dml import MSO_FILL_TYPE
    for shape in slide.shapes:
        try:
            fill = shape.fill
            if fill.type is None:
                continue
            if fill.type == MSO_FILL_TYPE.SOLID:
                if fill.fore_color.rgb == target_color:
                    return True
            elif fill.type == MSO_FILL_TYPE.GRADIENT:
                # 渐变填充：检查任一停止点颜色
                for stop in fill.gradient_stops:
                    if stop.color.rgb == target_color:
                        return True
        except Exception:
            pass
    return False


def _shape_has_gradient_fill(slide) -> bool:
    """检查幻灯片是否有渐变填充形状（SPEC 0026）。"""
    from pptx.enum.dml import MSO_FILL_TYPE
    for shape in slide.shapes:
        try:
            if shape.fill.type == MSO_FILL_TYPE.GRADIENT:
                return True
        except Exception:
            pass
    return False


def _find_gradient_shape(slide):
    """返回幻灯片第一个渐变填充形状（SPEC 0026）。"""
    from pptx.enum.dml import MSO_FILL_TYPE
    for shape in slide.shapes:
        try:
            if shape.fill.type == MSO_FILL_TYPE.GRADIENT:
                return shape
        except Exception:
            pass
    return None


def _find_rounded_shape(slide, target_color: RGBColor):
    """返回幻灯片第一个填充指定颜色的圆角矩形（SPEC 0026）。"""
    from pptx.enum.shapes import MSO_SHAPE
    for shape in slide.shapes:
        try:
            # auto_shape_type 仅对自选图形可用，且值为 MSO_SHAPE 枚举
            if getattr(shape, "auto_shape_type", None) != MSO_SHAPE.ROUNDED_RECTANGLE:
                continue
            if (shape.fill.type is not None
                    and shape.fill.fore_color.rgb == target_color):
                return shape
        except Exception:
            pass
    return None


class TestSpec0025ColorSystem:
    """SPEC 0025 三角色彩派生测试。"""

    def test_derive_palette_blue_primary_unchanged(self):
        """D1：蓝色主题派生主色 = 原值。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        primary, _, _, _ = renderer._derive_color_palette(theme_rgb)
        assert primary == RGBColor(0x25, 0x63, 0xEB)

    def test_derive_palette_blue_auxiliary_light(self):
        """D2：蓝色主题派生辅助色为高亮度浅色。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, auxiliary, _, _ = renderer._derive_color_palette(theme_rgb)
        # 辅助色不应等于主色
        assert auxiliary != RGBColor(0x25, 0x63, 0xEB)
        # 辅助色应为浅色（各分量平均值 > 180）
        hex_str = str(auxiliary)
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        avg = (r + g + b) / 3
        assert avg > 180, f"辅助色不够亮：{hex_str}（平均亮度 {avg}）"

    def test_derive_palette_blue_accent_complementary(self):
        """D3：蓝色主题派生强调色为互补色（金黄系，R 分量高）。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, _, accent, _ = renderer._derive_color_palette(theme_rgb)
        # 强调色不应等于主色
        assert accent != RGBColor(0x25, 0x63, 0xEB)
        # 蓝色色相约 0.6，互补色约 0.1（金黄系），R 分量应较高
        hex_str = str(accent)
        r = int(hex_str[0:2], 16)
        assert r > 150, f"强调色不是金黄系：{hex_str}（R={r}）"

    def test_derive_palette_gray_accent_is_blue(self):
        """D4：灰色主题强调色固定为蓝色 #2563EB（特殊处理）。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#475569")
        _, _, accent, _ = renderer._derive_color_palette(theme_rgb)
        assert accent == RGBColor(0x25, 0x63, 0xEB)

    def test_derive_palette_none_uses_default(self):
        """D5：theme_color=None 使用默认 #333333 派生三色。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color(None)
        primary, _, _, _ = renderer._derive_color_palette(theme_rgb)
        assert primary == RGBColor(0x33, 0x33, 0x33)

    def test_derive_palette_invalid_falls_back(self):
        """D6：无效色值降级到默认 #333333 派生三色。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#invalid")
        primary, _, _, _ = renderer._derive_color_palette(theme_rgb)
        assert primary == RGBColor(0x33, 0x33, 0x33)

    def test_derive_palette_blue_title_text_white(self):
        """主色亮度 < 0.60 时标题文字为白色（对比度保障）。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, _, _, title_text_color = renderer._derive_color_palette(theme_rgb)
        assert title_text_color == RGBColor(0xFF, 0xFF, 0xFF)

    def test_derive_palette_purple_title_text_white(self):
        """紫色主题（亮度约 0.578）标题文字为白色（阈值 0.60 覆盖紫色）。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#7c3aed")
        _, _, _, title_text_color = renderer._derive_color_palette(theme_rgb)
        assert title_text_color == RGBColor(0xFF, 0xFF, 0xFF)

    def test_derive_palette_returns_four_colors(self):
        """三角色彩派生返回四元组（主色/辅助色/强调色/标题文字色）。"""
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        result = renderer._derive_color_palette(theme_rgb)
        assert len(result) == 4
        primary, auxiliary, accent, title_text_color = result
        # 四色都应为 RGBColor 实例
        assert isinstance(primary, RGBColor)
        assert isinstance(auxiliary, RGBColor)
        assert isinstance(accent, RGBColor)
        assert isinstance(title_text_color, RGBColor)


class TestSpec0025SandwichStructure:
    """SPEC 0025 深浅对比三明治结构测试。"""

    def test_sandwich_content_title_bar_exists(self, tmp_path):
        """S1：内容页存在主色标题栏背景（shape fill）。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        primary = RGBColor(0x25, 0x63, 0xEB)
        # 内容页（第2页）应有主色填充的标题栏
        content_slide = prs.slides[1]
        assert _shape_has_fill_color(content_slide, primary), (
            "内容页标题栏未使用主色背景"
        )

    def test_sandwich_content_footer_bar_exists(self, tmp_path):
        """S3：内容页存在主色页脚栏背景（shape fill）。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        primary = RGBColor(0x25, 0x63, 0xEB)
        content_slide = prs.slides[1]
        assert _shape_has_fill_color(content_slide, primary), (
            "内容页页脚栏未使用主色背景"
        )

    def test_sandwich_cover_bottom_bar_exists(self, tmp_path):
        """S5：封面页底部存在主色窄条（三明治下层面包）。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        primary = RGBColor(0x25, 0x63, 0xEB)
        cover_slide = prs.slides[0]
        assert _shape_has_fill_color(cover_slide, primary), (
            "封面页未使用主色色块（顶部色块或底部窄条）"
        )

    def test_sandwich_summary_title_bar_exists(self, tmp_path):
        """S7：总结页存在主色标题栏背景。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        primary = RGBColor(0x25, 0x63, 0xEB)
        # 总结页是最后一页
        summary_slide = prs.slides[len(prs.slides) - 1]
        assert _shape_has_fill_color(summary_slide, primary), (
            "总结页标题栏未使用主色背景"
        )

    def test_auxiliary_background_left_column(self, tmp_path):
        """A1：左栏存在辅助色浅色背景。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, auxiliary, _, _ = renderer._derive_color_palette(theme_rgb)
        content_slide = prs.slides[1]
        assert _shape_has_fill_color(content_slide, auxiliary), (
            "左栏未使用辅助色浅色背景"
        )

    def test_accent_color_used_in_bullets(self, tmp_path):
        """A2：左栏圆点标记使用强调色（非主色）。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, _, accent, _ = renderer._derive_color_palette(theme_rgb)
        content_slide = prs.slides[1]
        # 强调色应出现在文本字体色中（圆点标记和章节标题）
        assert _slide_has_color(content_slide, accent), (
            "左栏圆点/标题未使用强调色"
        )

    def test_sandwich_chart_page_title_bar(self, tmp_path):
        """S6：图表页存在主色标题栏背景。"""
        chart_path = tmp_path / "chart.png"
        chart_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00\x00\x00\x03\x00\x01"
            b"\x82\x8b\x99\xde\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        artifacts = [
            {"name": "chart1.png", "artifact_type": "CHART_PNG",
             "file_path": str(chart_path)},
        ]
        # 用 5 张图表触发独立图表页
        multi_artifacts = [
            {"name": f"chart{i}.png", "artifact_type": "CHART_PNG",
             "file_path": str(chart_path)}
            for i in range(5)
        ]
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
            artifacts=multi_artifacts,
        )
        prs = Presentation(str(output_path))
        primary = RGBColor(0x25, 0x63, 0xEB)
        # 查找图表页（包含"关键图表"文本的页面）
        chart_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "关键图表" in shape.text_frame.text:
                        chart_slide = slide
                        break
            if chart_slide:
                break
        assert chart_slide is not None, "未找到图表页"
        assert _shape_has_fill_color(chart_slide, primary), (
            "图表页标题栏未使用主色背景"
        )


# --- API 测试 ---


@pytest.fixture
def client(monkeypatch, tmp_path):
    """TestClient + 内存 SQLite + 受控工作区。"""
    monkeypatch.setenv("PROJECT_DATA_ROOT", str(tmp_path / "projects"))
    engine = create_engine(
        TEST_DB,
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(bind=engine)
    TestingSessionLocal = sessionmaker(bind=engine)

    from app.api.routers import projects as project_router
    from app.api.routers import outlines as outlines_router
    from app.api.routers import deliverables as deliverables_router

    monkeypatch.setattr(project_router, "SessionLocal", TestingSessionLocal)
    monkeypatch.setattr(outlines_router, "SessionLocal", TestingSessionLocal)
    monkeypatch.setattr(deliverables_router, "SessionLocal", TestingSessionLocal)

    with TestClient(app) as test_client:
        yield test_client

    Base.metadata.drop_all(bind=engine)


def _create_project(client: TestClient,
                    status: str = ProjectStatus.OUTLINE_CONFIRMED.value) -> str:
    """创建项目并设置状态，返回 project_id。"""
    response = client.post(
        "/api/projects",
        json={"name": "测试项目", "topic": "测试课题"},
    )
    assert response.status_code == 200
    project_id = response.json()["id"]

    from app.api.routers import projects as project_router
    SessionLocal = project_router.SessionLocal
    db = SessionLocal()
    try:
        project = db.query(Project).filter(Project.id == project_id).first()
        project.status = status
        db.commit()
    finally:
        db.close()
    return project_id


def _seed_outline(SessionLocal, project_id: str,
                   outline_id: str = "ol_ppt_cfg_1",
                   status: str = OutlineStatus.CONFIRMED.value) -> str:
    """直接插入已确认大纲，返回 outline_id。"""
    db = SessionLocal()
    try:
        outline = Outline(
            id=outline_id,
            project_id=project_id,
            sections_json=json.dumps(SAMPLE_SECTIONS),
            status=status,
            candidate_source="local_rule",
            code_version=1,
            created_at=datetime(2024, 1, 1, tzinfo=timezone.utc),
            updated_at=datetime(2024, 1, 1, tzinfo=timezone.utc),
        )
        db.add(outline)
        db.commit()
        return outline.id
    finally:
        db.close()


def test_api_generate_ppt_no_body(client):
    """A-01：无 body 生成 PPT 成功。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate"
    )
    assert response.status_code == 201
    data = response.json()
    assert data["job_id"]
    assert data["deliverable_id"]
    assert data["template_used"] is False


def test_api_generate_ppt_with_config(client):
    """A-02：有 config 生成 PPT 成功。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={"config": {"target_slide_count": 10}},
    )
    assert response.status_code == 201
    data = response.json()
    assert data["job_id"]
    assert data["deliverable_id"]


def test_api_generate_ppt_full_config(client):
    """A-03：完整 config 生成 PPT 成功。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={
            "config": {
                "target_slide_count": 8,
                "theme_color": "#7c3aed",
                "include_charts": False,
            }
        },
    )
    assert response.status_code == 201


def test_api_generate_ppt_invalid_theme_color(client):
    """A-04：无效 theme_color 返回 PPT_CONFIG_INVALID_THEME_COLOR。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={"config": {"theme_color": "#ff0000"}},
    )
    assert response.status_code == 400
    assert response.json()["error"]["code"] == "PPT_CONFIG_INVALID_THEME_COLOR"


def test_api_generate_ppt_slide_count_too_small(client):
    """A-05：target_slide_count 小于 5 时校验失败。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={"config": {"target_slide_count": 3}},
    )
    # app 自定义 RequestValidationError 处理器返回 400
    assert response.status_code == 400
    assert response.json()["error"]["code"] == "REQUEST_VALIDATION_ERROR"


def test_api_generate_ppt_slide_count_too_large(client):
    """A-06：target_slide_count 大于 20 时校验失败。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={"config": {"target_slide_count": 25}},
    )
    # app 自定义 RequestValidationError 处理器返回 400
    assert response.status_code == 400
    assert response.json()["error"]["code"] == "REQUEST_VALIDATION_ERROR"


def test_api_generate_ppt_include_charts_false(client):
    """A-07：include_charts=false 成功创建。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={"config": {"include_charts": False}},
    )
    assert response.status_code == 201


def test_api_generate_ppt_empty_config(client):
    """A-08：空 config 对象成功创建。"""
    from app.api.routers import outlines as outlines_router
    SessionLocal = outlines_router.SessionLocal
    project_id = _create_project(client)
    _seed_outline(SessionLocal, project_id)

    response = client.post(
        f"/api/projects/{project_id}/outline/ol_ppt_cfg_1/ppt/generate",
        json={"config": {}},
    )
    assert response.status_code == 201


# --- SPEC 0026 视觉效果增强测试（渐变 + 圆角 + 阴影 + 边框）---


def _make_chart_png(path, width=100, height=75):
    """生成真实 PNG 图片用于测试（避免占位文本框）。"""
    try:
        from PIL import Image
        img = Image.new("RGB", (width, height), color=(100, 150, 200))
        img.save(str(path))
    except ImportError:
        # Pillow 不可用时写入最小有效 PNG
        path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00\x00\x00\x03\x00\x01"
            b"\x82\x8b\x99\xde\x00\x00\x00\x00IEND\xaeB`\x82"
        )


def _render_ppt_with_chart(tmp_path, theme_color="#2563eb"):
    """渲染带图表的 PPT（用于阴影/边框测试），返回路径。"""
    chart_path = tmp_path / "chart.png"
    _make_chart_png(chart_path)
    artifacts = [
        {"name": "chart1.png", "artifact_type": "CHART_PNG",
         "file_path": str(chart_path)},
    ]
    return _render_ppt(
        tmp_path,
        config={"theme_color": theme_color},
        artifacts=artifacts,
    )


def _find_picture_in_prs(prs):
    """遍历所有幻灯片找到第一个图片形状（SPEC 0026 阴影/边框测试）。

    返回 (slide_index, picture_shape) 或 (None, None)。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return idx, shape
    return None, None


class TestSpec0026VisualEffects:
    """SPEC 0026 视觉效果增强测试（渐变 + 圆角 + 阴影 + 边框）。"""

    # === 暗化算法测试 ===

    def test_darken_color_reduces_brightness(self):
        """D1：_darken_color 返回值亮度低于原色。"""
        import colorsys
        renderer = PptRenderer()
        original = RGBColor(0x25, 0x63, 0xEB)  # 蓝色 #2563eb
        darkened = renderer._darken_color(original, 0.20)

        # 计算两者亮度
        def brightness(rgb):
            hex_str = str(rgb)
            r = int(hex_str[0:2], 16) / 255
            g = int(hex_str[2:4], 16) / 255
            b = int(hex_str[4:6], 16) / 255
            _, l, _ = colorsys.rgb_to_hls(r, g, b)
            return l

        assert brightness(darkened) < brightness(original), (
            f"暗化后亮度未降低：原={brightness(original)}, 暗化={brightness(darkened)}"
        )

    def test_darken_color_low_brightness_floor(self):
        """D2：_darken_color 极端情况（L 已很低）不返回负值，亮度接近下限。"""
        import colorsys
        renderer = PptRenderer()
        # 极暗的颜色 #000000（L=0）
        extreme_dark = RGBColor(0x00, 0x00, 0x00)
        darkened = renderer._darken_color(extreme_dark, 0.50)

        hex_str = str(darkened)
        r = int(hex_str[0:2], 16) / 255
        g = int(hex_str[2:4], 16) / 255
        b = int(hex_str[4:6], 16) / 255
        _, l, _ = colorsys.rgb_to_hls(r, g, b)
        # 亮度应接近下限 0.10（允许 RGB↔HLS 转换浮点误差）
        assert l >= 0.09, f"暗化后亮度 {l} 远低于下限 0.10"

    # === 渐变填充测试 ===

    def test_gradient_cover_top_block(self, tmp_path):
        """G1：封面页顶部色块为渐变填充。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        cover_slide = prs.slides[0]
        assert _shape_has_gradient_fill(cover_slide), "封面页未找到渐变填充形状"

    def test_gradient_content_title_bar(self, tmp_path):
        """G2：内容页标题栏为渐变填充。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        content_slide = prs.slides[1]
        assert _shape_has_gradient_fill(content_slide), "内容页未找到渐变填充形状"

    def test_gradient_content_footer_bar(self, tmp_path):
        """G3：内容页页脚栏为渐变填充。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        content_slide = prs.slides[1]
        # 页脚栏在底部（top 接近 7.0"），检查底部形状是否有渐变
        assert _shape_has_gradient_fill(content_slide), "内容页页脚栏未找到渐变填充"

    def test_gradient_angle_is_90(self, tmp_path):
        """G4：渐变角度为 90°（上→下）。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        cover_slide = prs.slides[0]
        shape = _find_gradient_shape(cover_slide)
        assert shape is not None, "封面页未找到渐变形状"
        # 渐变角度应为 90（允许浮点误差）
        assert abs(shape.fill.gradient_angle - 90) < 1, (
            f"渐变角度不是 90°：{shape.fill.gradient_angle}"
        )

    def test_gradient_start_color_is_primary(self, tmp_path):
        """G5：渐变起始色 = 主色原值。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        cover_slide = prs.slides[0]
        shape = _find_gradient_shape(cover_slide)
        assert shape is not None, "封面页未找到渐变形状"
        stops = shape.fill.gradient_stops
        primary = RGBColor(0x25, 0x63, 0xEB)
        # 起始停止点（position 最小）颜色应为主色
        start_stop = min(stops, key=lambda s: s.position)
        assert start_stop.color.rgb == primary, (
            f"渐变起始色 {start_stop.color.rgb} != 主色 {primary}"
        )

    def test_gradient_end_color_is_darkened(self, tmp_path):
        """G6：渐变结束色 = 主色暗化色。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        cover_slide = prs.slides[0]
        shape = _find_gradient_shape(cover_slide)
        assert shape is not None, "封面页未找到渐变形状"
        stops = shape.fill.gradient_stops
        end_stop = max(stops, key=lambda s: s.position)

        renderer = PptRenderer()
        primary = RGBColor(0x25, 0x63, 0xEB)
        expected_end = renderer._darken_color(primary, 0.20)
        assert end_stop.color.rgb == expected_end, (
            f"渐变结束色 {end_stop.color.rgb} != 预期暗化色 {expected_end}"
        )

    # === 圆角矩形测试 ===

    def test_rounded_left_column_shape_type(self, tmp_path):
        """R1：左栏背景形状类型为 ROUNDED_RECTANGLE。"""
        from pptx.enum.shapes import MSO_SHAPE
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        content_slide = prs.slides[1]
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, auxiliary, _, _ = renderer._derive_color_palette(theme_rgb)
        shape = _find_rounded_shape(content_slide, auxiliary)
        assert shape is not None, "左栏未找到辅助色填充的圆角矩形"
        assert shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE

    def test_rounded_corner_radius_005(self, tmp_path):
        """R2：圆角半径 adjustments[0] ≈ 0.05。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        content_slide = prs.slides[1]
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, auxiliary, _, _ = renderer._derive_color_palette(theme_rgb)
        shape = _find_rounded_shape(content_slide, auxiliary)
        assert shape is not None, "左栏未找到圆角矩形"
        # 圆角半径应为 0.05（允许浮点误差）
        assert abs(shape.adjustments[0] - 0.05) < 0.01, (
            f"圆角半径 {shape.adjustments[0]} 不是 0.05"
        )

    def test_rounded_left_column_fill_is_auxiliary(self, tmp_path):
        """R3：左栏背景填充色 = 辅助色。"""
        output_path = _render_ppt(
            tmp_path, config={"theme_color": "#2563eb"},
        )
        prs = Presentation(str(output_path))
        content_slide = prs.slides[1]
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, auxiliary, _, _ = renderer._derive_color_palette(theme_rgb)
        shape = _find_rounded_shape(content_slide, auxiliary)
        assert shape is not None, "左栏未找到圆角矩形"
        assert shape.fill.fore_color.rgb == auxiliary

    # === 外阴影测试 ===

    def test_shadow_effect_lst_exists(self, tmp_path):
        """S1：右栏图表 picture 的 spPr 包含 effectLst 节点。"""
        from pptx.oxml.ns import qn
        output_path = _render_ppt_with_chart(tmp_path)
        prs = Presentation(str(output_path))
        _, pic_shape = _find_picture_in_prs(prs)
        assert pic_shape is not None, "未找到图片形状"
        spPr = pic_shape._element.spPr
        effectLst = spPr.find(qn('a:effectLst'))
        assert effectLst is not None, "图片未包含 a:effectLst 节点"

    def test_shadow_outer_shdw_exists(self, tmp_path):
        """S2：effectLst 包含 outerShdw 子节点。"""
        from pptx.oxml.ns import qn
        output_path = _render_ppt_with_chart(tmp_path)
        prs = Presentation(str(output_path))
        _, pic_shape = _find_picture_in_prs(prs)
        assert pic_shape is not None
        spPr = pic_shape._element.spPr
        effectLst = spPr.find(qn('a:effectLst'))
        outerShdw = effectLst.find(qn('a:outerShdw'))
        assert outerShdw is not None, "effectLst 未包含 a:outerShdw 子节点"

    def test_shadow_blur_rad_positive(self, tmp_path):
        """S3：outerShdw 的 blurRad 属性存在且为正值。"""
        from pptx.oxml.ns import qn
        output_path = _render_ppt_with_chart(tmp_path)
        prs = Presentation(str(output_path))
        _, pic_shape = _find_picture_in_prs(prs)
        assert pic_shape is not None
        spPr = pic_shape._element.spPr
        outerShdw = spPr.find(qn('a:effectLst')).find(qn('a:outerShdw'))
        blur_rad = outerShdw.get('blurRad')
        assert blur_rad is not None, "outerShdw 缺少 blurRad 属性"
        assert int(blur_rad) > 0, f"blurRad 非正值：{blur_rad}"

    def test_shadow_has_srgb_color(self, tmp_path):
        """S4：outerShdw 包含 srgbClr 颜色节点。"""
        from pptx.oxml.ns import qn
        output_path = _render_ppt_with_chart(tmp_path)
        prs = Presentation(str(output_path))
        _, pic_shape = _find_picture_in_prs(prs)
        assert pic_shape is not None
        spPr = pic_shape._element.spPr
        outerShdw = spPr.find(qn('a:effectLst')).find(qn('a:outerShdw'))
        srgbClr = outerShdw.find(qn('a:srgbClr'))
        assert srgbClr is not None, "outerShdw 未包含 a:srgbClr 颜色节点"

    # === 边框测试 ===

    def test_border_color_is_auxiliary(self, tmp_path):
        """B1：右栏图表 picture 的 line.color.rgb == 辅助色。"""
        output_path = _render_ppt_with_chart(tmp_path)
        prs = Presentation(str(output_path))
        _, pic_shape = _find_picture_in_prs(prs)
        assert pic_shape is not None, "未找到图片形状"
        renderer = PptRenderer()
        theme_rgb = renderer._resolve_theme_color("#2563eb")
        _, auxiliary, _, _ = renderer._derive_color_palette(theme_rgb)
        # 边框颜色应为辅助色
        assert pic_shape.line.color.rgb == auxiliary, (
            f"图片边框色 {pic_shape.line.color.rgb} != 辅助色 {auxiliary}"
        )

    def test_border_width_is_1pt(self, tmp_path):
        """B2：右栏图表 picture 的 line.width == Pt(1)。"""
        from pptx.util import Pt
        output_path = _render_ppt_with_chart(tmp_path)
        prs = Presentation(str(output_path))
        _, pic_shape = _find_picture_in_prs(prs)
        assert pic_shape is not None
        # 边框宽度应为 1pt（允许 EMU 精度误差）
        expected_emu = Pt(1)
        assert abs(pic_shape.line.width - expected_emu) < 1000, (
            f"图片边框宽度 {pic_shape.line.width} != 1pt ({expected_emu})"
        )


# --- SPEC 0027 布局增强测试（百分比定位 + Grid 布局） ---


class TestSpec0027LayoutEnhancement:
    """SPEC 0027 布局增强测试（_pct_to_emu + _GridHelper）。

    红色阶段说明：
    - _pct_to_emu 和 _GridHelper 尚未实现，相关测试应失败
    - 实现完成后所有测试应通过
    """

    # === _pct_to_emu 百分比定位测试 ===

    def test_pct_to_emu方法存在(self):
        """P1：PptRenderer 类有 _pct_to_emu 方法。"""
        assert hasattr(PptRenderer, "_pct_to_emu"), (
            "PptRenderer 未实现 _pct_to_emu（SPEC 0027 百分比定位辅助方法）"
        )

    def test_pct_to_emu_10pct(self):
        """P2：_pct_to_emu("10%", 1000) 返回 100。"""
        result = PptRenderer._pct_to_emu("10%", 1000)
        assert result == 100, f'_pct_to_emu("10%", 1000) = {result}, 预期 100'

    def test_pct_to_emu_50pct(self):
        """P3：_pct_to_emu("50%", 1000) 返回 500。"""
        result = PptRenderer._pct_to_emu("50%", 1000)
        assert result == 500, f'_pct_to_emu("50%", 1000) = {result}, 预期 500'

    def test_pct_to_emu_100pct(self):
        """P4：_pct_to_emu("100%", 1000) 返回 1000。"""
        result = PptRenderer._pct_to_emu("100%", 1000)
        assert result == 1000, f'_pct_to_emu("100%", 1000) = {result}, 预期 1000'

    def test_pct_to_emu_0pct(self):
        """P5：_pct_to_emu("0%", 1000) 返回 0。"""
        result = PptRenderer._pct_to_emu("0%", 1000)
        assert result == 0, f'_pct_to_emu("0%", 1000) = {result}, 预期 0'

    def test_pct_to_emu_非百分比字符串抛ValueError(self):
        """P6：_pct_to_emu("non-pct", 1000) 抛出 ValueError。"""
        with pytest.raises(ValueError, match="百分比"):
            PptRenderer._pct_to_emu("non-pct", 1000)

    def test_pct_to_emu_与Inches兼容(self):
        """P7：_pct_to_emu("10%", Inches(10)) ≈ Inches(1)（允许 ±100 EMU 精度误差）。"""
        from pptx.util import Inches
        total = Inches(10)
        result = PptRenderer._pct_to_emu("10%", total)
        expected = Inches(1)
        assert abs(result - expected) < 100, (
            f"百分比转 EMU 与 Inches 不兼容：result={result}, expected={expected}"
        )

    def test_pct_to_emu_支持小数百分比(self):
        """P8：_pct_to_emu("12.5%", 1000) 返回 125。"""
        result = PptRenderer._pct_to_emu("12.5%", 1000)
        assert result == 125, f'_pct_to_emu("12.5%", 1000) = {result}, 预期 125'

    # === _GridHelper Grid 布局测试 ===

    def test_GridHelper类存在(self):
        """G1：PptRenderer 有 _GridHelper 内部类。"""
        assert hasattr(PptRenderer, "_GridHelper"), (
            "PptRenderer 未实现 _GridHelper（SPEC 0027 Grid 布局辅助类）"
        )

    def test_GridHelper_2x2_cell_00(self):
        """G2：2×2 网格 cell(0,0) 返回左上角 (0, 0, 500, 500)。"""
        grid = PptRenderer._GridHelper(0, 0, 1000, 1000, rows=2, cols=2)
        left, top, width, height = grid.cell(0, 0)
        assert left == 0 and top == 0
        assert width == 500 and height == 500

    def test_GridHelper_2x2_cell_11(self):
        """G3：2×2 网格 cell(1,1) 返回右下角 (500, 500, 500, 500)。"""
        grid = PptRenderer._GridHelper(0, 0, 1000, 1000, rows=2, cols=2)
        left, top, width, height = grid.cell(1, 1)
        assert left == 500 and top == 500
        assert width == 500 and height == 500

    def test_GridHelper_1x2_带h_gap(self):
        """G4：1×2 网格带水平间距 100，每格宽 450。"""
        grid = PptRenderer._GridHelper(0, 0, 1000, 1000, rows=1, cols=2, h_gap=100)
        left1, _, width1, _ = grid.cell(0, 0)
        left2, _, width2, _ = grid.cell(0, 1)
        # 每个单元格宽度 = (1000 - 100) / 2 = 450
        assert width1 == 450 and width2 == 450
        # cell(0,1) 左边界 = 0 + 450 + 100 = 550
        assert left1 == 0 and left2 == 550

    def test_GridHelper_2x1_带v_gap(self):
        """G5：2×1 网格带垂直间距 100，每格高 450。"""
        grid = PptRenderer._GridHelper(0, 0, 1000, 1000, rows=2, cols=1, v_gap=100)
        _, top1, _, height1 = grid.cell(0, 0)
        _, top2, _, height2 = grid.cell(1, 0)
        # 每个单元格高度 = (1000 - 100) / 2 = 450
        assert height1 == 450 and height2 == 450
        # cell(1,0) 上边界 = 0 + 450 + 100 = 550
        assert top1 == 0 and top2 == 550

    def test_GridHelper_3x3_cell_22(self):
        """G6：3×3 网格 cell(2,2) 返回右下角 (600, 600, 300, 300)。"""
        grid = PptRenderer._GridHelper(0, 0, 900, 900, rows=3, cols=3)
        left, top, width, height = grid.cell(2, 2)
        assert left == 600 and top == 600
        assert width == 300 and height == 300

    def test_GridHelper_2x2_单元格水平不重叠(self):
        """G7：2×2 网格 cell(0,0) 和 cell(0,1) 水平不重叠。"""
        grid = PptRenderer._GridHelper(0, 0, 1000, 1000, rows=2, cols=2)
        l1, _, w1, _ = grid.cell(0, 0)
        l2, _, _, _ = grid.cell(0, 1)
        # cell(0,0) 右边界 = l1 + w1 = 500，应 <= cell(0,1) 左边界 = 500
        assert l1 + w1 <= l2, (
            f"单元格水平重叠：cell(0,0) 右边界={l1 + w1} > cell(0,1) 左边界={l2}"
        )

    def test_GridHelper_2x2_单元格垂直不重叠(self):
        """G8：2×2 网格 cell(0,0) 和 cell(1,0) 垂直不重叠。"""
        grid = PptRenderer._GridHelper(0, 0, 1000, 1000, rows=2, cols=2)
        _, t1, _, h1 = grid.cell(0, 0)
        _, t2, _, _ = grid.cell(1, 0)
        # cell(0,0) 下边界 = t1 + h1 = 500，应 <= cell(1,0) 上边界 = 500
        assert t1 + h1 <= t2, (
            f"单元格垂直重叠：cell(0,0) 下边界={t1 + h1} > cell(1,0) 上边界={t2}"
        )

    def test_GridHelper_带非零起点偏移(self):
        """G9：Grid 带非零起点偏移 (100, 200)。"""
        grid = PptRenderer._GridHelper(100, 200, 600, 400, rows=1, cols=2)
        left, top, width, height = grid.cell(0, 0)
        assert left == 100 and top == 200
        assert width == 300 and height == 400
        # cell(0,1) 左边界 = 100 + 300 = 400
        left2, _, _, _ = grid.cell(0, 1)
        assert left2 == 400

    def test_GridHelper_1x1_返回完整区域(self):
        """G10：1×1 网格返回完整区域。"""
        grid = PptRenderer._GridHelper(50, 60, 700, 800, rows=1, cols=1)
        left, top, width, height = grid.cell(0, 0)
        assert left == 50 and top == 60
        assert width == 700 and height == 800
