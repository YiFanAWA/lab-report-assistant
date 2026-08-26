import pytest

"""SPEC 0043：答辩 PPT 出版级版面专项测试。

这些测试只验证 renderer 投影层，不改变论文语义规划器或案例输入。
"""

from pathlib import Path
import pytest
from types import SimpleNamespace

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from app.infrastructure.renderers import ppt_renderer as renderer_module
from app.infrastructure.renderers.ppt_renderer import PptRenderer
from app.modules.outlines.layout_planner import LayoutKind


def _make_chart(path: Path) -> str:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1200, 700), (242, 246, 250)).save(path, format="PNG")
    return str(path)


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def _render_defense_fixture(tmp_path: Path, monkeypatch):
    chart_path = Path(_make_chart(tmp_path / "primary_result.png"))
    artifact = {
        "artifact_type": "CHART_PNG",
        "file_path": str(chart_path),
        "name": "primary_result",
        "figure_caption": "图1 主要结果的区间估计",
        "figure_lead": "主要结果在预设分组中保持同方向。",
        "figure_takeaway": "结果支持描述性差异，不能外推为因果效应。",
        "scientific_asset_ids": [],
    }
    result_slide = SimpleNamespace(
        role="result",
        title="主要结果：分组差异",
        content="结果显示不同分组之间存在可观察差异。",
        layout_kind=LayoutKind.RESULT_FOCUS,
        chart_artifacts=(artifact,),
        steps=(),
        metrics=(),
        figure_family="statistics",
        figure_lead=artifact["figure_lead"],
        figure_takeaway=artifact["figure_takeaway"],
    )
    limitation_slide = SimpleNamespace(
        role="limitation",
        title="解释范围",
        content="本地复核使用公开数据，结论受样本筛选与观察性设计限制。",
        layout_kind=LayoutKind.NARRATIVE,
        chart_artifacts=(),
        steps=(),
        metrics=(),
        figure_family="",
        figure_lead="",
        figure_takeaway="",
    )
    monkeypatch.setattr(
        renderer_module,
        "plan_defense_deck",
        lambda *_args, **_kwargs: SimpleNamespace(
            slides=(result_slide, limitation_slide),
        ),
    )

    output = tmp_path / "spec0043_defense.pptx"
    PptRenderer().render(
        project_name="论文解读",
        project_topic="公开数据的结果复核",
        outline_sections=[],
        execution_artifacts=[artifact],
        output_path=str(output),
        config={"ppt_workflow": "academic", "theme_preset": "PACIFIC_DEEP"},
    )
    return Presentation(str(output))


def test_defense_typography_contract_is_publication_scale():
    assert renderer_module.DEFENSE_TITLE_MIN_PT >= 35
    assert renderer_module.DEFENSE_BODY_MIN_PT >= 18
    assert renderer_module.DEFENSE_CAPTION_MIN_PT >= 12


def test_defense_result_slide_prioritizes_visual_and_removes_formulaic_copy(
    tmp_path, monkeypatch
):
    prs = _render_defense_fixture(tmp_path, monkeypatch)
    result_slide = prs.slides[1]

    pictures = [
        shape
        for shape in result_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert pictures, "答辩结果页必须包含主视觉"
    picture = pictures[0]
    assert picture.width == Inches(10.2)
    assert picture.height == Inches(5.95)
    assert picture.width / picture.height == pytest.approx(1.714, abs=0.002)
    slide_area = prs.slide_width * prs.slide_height
    visual_area = picture.width * picture.height
    assert visual_area / slide_area >= 0.60

    all_text = "\n".join(_slide_text(slide) for slide in prs.slides)
    assert "解释边界" not in all_text
    assert "补充图表用于从不同角度核对同一结果" not in all_text


def test_defense_slide_text_sizes_meet_spec(tmp_path, monkeypatch):
    prs = _render_defense_fixture(tmp_path, monkeypatch)
    for slide in list(prs.slides)[1:]:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    assert run.font.size is not None
                    if shape.top < Inches(1.2):
                        assert run.font.size >= Pt(renderer_module.DEFENSE_TITLE_MIN_PT)
                    elif "图" in run.text or "来源" in run.text:
                        assert run.font.size >= Pt(renderer_module.DEFENSE_CAPTION_MIN_PT)
                    else:
                        assert run.font.size >= Pt(renderer_module.DEFENSE_BODY_MIN_PT)
