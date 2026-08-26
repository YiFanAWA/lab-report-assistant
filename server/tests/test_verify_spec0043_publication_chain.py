import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from scripts import verify_spec0043_publication_chain as validator


def _write_manifest(path: Path, payload: dict) -> Path:
    path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    return path


def _add_text(slide, text: str, top: float, font_pt: float):
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(8), Inches(0.6))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    return shape


def _write_font_fixture(path: Path, body_pt: float = 18, title_pt: float = 35, caption_pt: float = 12) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    presentation.slides.add_slide(blank)
    slide = presentation.slides.add_slide(blank)
    _add_text(slide, "页面标题", 0.6, title_pt)
    _add_text(slide, "正文内容", 2.0, body_pt)
    _add_text(slide, "图注：执行结果", 3.0, caption_pt)
    presentation.save(path)
    return path


def test_current_deliverables_manifest_is_supported(tmp_path):
    manifest = _write_manifest(
        tmp_path / "publication_manifest.json",
        {
            "spec": "0043",
            "deliverables": {
                "docx": "publication.docx",
                "pdf": "publication.pdf",
                "pptx": "publication.pptx",
            },
        },
    )

    result = validator.verify_manifest(manifest)
    paths = validator.resolve_manifest_paths(manifest, json.loads(manifest.read_text(encoding="utf-8")))

    assert result.status == "PASS"
    assert result.details["container"] == "deliverables"
    assert paths == {
        "docx": tmp_path / "publication.docx",
        "pdf": tmp_path / "publication.pdf",
        "pptx": tmp_path / "publication.pptx",
    }


def test_legacy_artifacts_list_remains_supported(tmp_path):
    manifest = _write_manifest(
        tmp_path / "legacy_manifest.json",
        {
            "spec": "0043",
            "artifacts": [
                {"artifact_type": "DOCX", "file_path": "legacy.docx"},
                {"artifact_type": "PDF", "file_path": "legacy.pdf"},
                {"artifact_type": "PPTX", "file_path": "legacy.pptx"},
            ],
        },
    )

    result = validator.verify_manifest(manifest)
    paths = validator.resolve_manifest_paths(manifest, json.loads(manifest.read_text(encoding="utf-8")))

    assert result.status == "PASS"
    assert paths["docx"] == tmp_path / "legacy.docx"
    assert paths["pdf"] == tmp_path / "legacy.pdf"
    assert paths["pptx"] == tmp_path / "legacy.pptx"


def test_invalid_manifest_does_not_scan_root_or_tmp(tmp_path, monkeypatch, capsys):
    manifest = _write_manifest(
        tmp_path / "invalid_manifest.json",
        {"spec": "0043", "deliverables": {"docx": "only.docx"}},
    )
    (tmp_path / ".tmp").mkdir()
    (tmp_path / ".tmp" / "old.docx").write_bytes(b"old")

    def fail_if_scanned(*args, **kwargs):
        pytest.fail("manifest 已提供时不应扫描 root/.tmp")

    monkeypatch.setattr(validator, "find_artifact", fail_if_scanned)
    assert validator.main(["--manifest", str(manifest), "--root", str(tmp_path)]) == 1
    output = json.loads(capsys.readouterr().out)

    assert output["status"] == "FAIL"
    assert output["artifacts"]["docx"] is None


def test_pptx_font_contract_allows_12pt_caption_but_requires_18pt_body(tmp_path):
    pptx = _write_font_fixture(tmp_path / "font_contract.pptx")
    checks = {item.name: item for item in validator.verify_pptx(pptx, minimum_slides=1)}

    assert checks["pptx.font.title"].status == "PASS"
    assert checks["pptx.font.body"].status == "PASS"
    assert checks["pptx.font.caption"].status == "PASS"
    assert checks["pptx.font"].details["observed_caption_pt"] == 12.0

    low_body = _write_font_fixture(tmp_path / "low_body.pptx", body_pt=17)
    low_checks = {item.name: item for item in validator.verify_pptx(low_body, minimum_slides=1)}

    assert low_checks["pptx.font.body"].status == "FAIL"
    assert low_checks["pptx.font"].status == "FAIL"
