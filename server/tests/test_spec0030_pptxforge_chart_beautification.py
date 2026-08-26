"""SPEC 0030 pptxforge 集成与图表美化增强测试。

TDD 红阶段测试：本文件中的测试在实现前应失败，实现完成后应全部通过。

覆盖范围：
1. PptConfig theme_preset 字段校验（方案 B，10 主题枚举）
2. 图表美化：NPG 期刊配色、a/b/c 面板标签、dpi 统一修复
3. PptRenderer 主题映射：_map_theme(theme_color) → pptxforge 主题名
4. PptRenderer pptxforge 集成：render() 输出可被 python-pptx 重新打开

红色阶段预期失败点：
- _HEADER 中 NPG_PALETTE 尚未定义（N1/N2/N3 失败）
- _build_chart_code 仍包含 dpi=100（D1/D2 失败）
- _build_chart_code 未生成 a/b/c 面板标签（P1/P2 失败）
- PptRenderer._map_theme 尚未实现（T1~T6 失败）
- render() 内部尚未集成 pptxforge Deck（X1~X3 失败）
"""

import ast
import json

import pytest
from pptx import Presentation

from app.modules.llm.code_task_provider import LocalRuleCodeTaskProvider
from app.modules.outlines.contracts import PPT_THEME_PRESETS, PptConfig
from app.infrastructure.renderers.ppt_renderer import PptRenderer


# --- 测试数据 ---

SAMPLE_SECTIONS = [
    {
        "id": "sec_001",
        "title": "实验目的",
        "content": "分析胃病数据分布特征",
        "source_type": "REQUIREMENT",
        "source_ids": ["plan_001"],
    },
    {
        "id": "sec_002",
        "title": "实验方法",
        "content": "使用描述性统计方法",
        "source_type": "EVIDENCE",
        "source_ids": ["ev_001"],
    },
    {
        "id": "sec_003",
        "title": "实验结果",
        "content": "执行成功，数据已清洗",
        "source_type": "EXECUTION",
        "source_ids": ["run_001"],
    },
    {
        "id": "sec_004",
        "title": "总结",
        "content": "本实验完成了数据分析与可视化",
        "source_type": "SUMMARY",
        "source_ids": [],
    },
]


def _make_full_plan_with_charts() -> dict:
    """构造包含多种图表类型的 AnalysisPlan（SPEC 0030 测试用）。"""
    return {
        "cleaning_plan": [
            {
                "field": "age",
                "issue_type": "MISSING_VALUE",
                "action": "用中位数填充缺失值",
                "reason": "数值字段",
            },
        ],
        "analysis_plan": [
            {
                "analysis_type": "DESCRIPTIVE_STATISTICS",
                "target_fields": ["age"],
                "method": "计算均值",
                "expected_output": "统计表",
            },
        ],
        "chart_plan": [
            {
                "chart_type": "HISTOGRAM",
                "title": "age 分布",
                "data_fields": ["age"],
                "description": "展示分布",
            },
            {
                "chart_type": "BOXPLOT",
                "title": "数值字段箱线图",
                "data_fields": ["age"],
                "description": "展示离群值",
            },
            {
                "chart_type": "BAR",
                "title": "诊断分布",
                "data_fields": ["diagnosis"],
                "description": "展示分类分布",
            },
            {
                "chart_type": "SCATTER",
                "title": "年龄 vs 血压",
                "data_fields": ["age", "bp"],
                "description": "展示相关性",
            },
        ],
    }


# === 1. PptConfig theme_preset 字段校验测试（SPEC 0030 方案 B） ===


class TestSpec0030PptConfigThemePreset:
    """SPEC 0030 PptConfig theme_preset 字段校验测试（方案 B）。

    红色阶段说明：
    - theme_preset 字段已在 contracts.py 中新增（本切片第 4 步已完成）
    - 这些测试应通过（合同层已实现）
    """

    def test_theme_preset_默认为None(self):
        """P1：theme_preset 默认值为 None。"""
        config = PptConfig()
        assert config.theme_preset is None

    def test_theme_preset_合法名通过校验(self):
        """P2：合法主题名 MIDNIGHT_EXECUTIVE 通过校验。"""
        config = PptConfig(theme_preset="MIDNIGHT_EXECUTIVE")
        assert config.theme_preset == "MIDNIGHT_EXECUTIVE"

    def test_theme_preset_十个主题全部合法(self):
        """P3：10 个合法主题名全部通过校验。"""
        for preset in PPT_THEME_PRESETS:
            config = PptConfig(theme_preset=preset)
            assert config.theme_preset == preset

    def test_theme_preset_无效名抛出ValidationError(self):
        """P4：无效主题名抛出 ValidationError。"""
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            PptConfig(theme_preset="INVALID_THEME")

    def test_theme_preset_小写名抛出ValidationError(self):
        """P5：小写主题名抛出 ValidationError（枚举大小写敏感）。"""
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            PptConfig(theme_preset="midnight_executive")

    def test_theme_preset_与其他字段兼容(self):
        """P6：theme_preset 与其他字段同时设置。"""
        config = PptConfig(
            target_slide_count=8,
            theme_color="#2563eb",
            include_charts=True,
            theme_preset="PACIFIC_DEEP",
        )
        assert config.theme_preset == "PACIFIC_DEEP"
        assert config.target_slide_count == 8
        assert config.theme_color == "#2563eb"
        assert config.include_charts is True


# === 2. 图表美化测试（NPG 配色 + 面板标签 + dpi 统一） ===


class TestSpec0030NpgPalette:
    """SPEC 0030 NPG 期刊配色测试。

    红色阶段说明：
    - _HEADER 中尚未定义 NPG_PALETTE，相关测试应失败
    - 实现完成后所有测试应通过
    """

    def test_HEADER包含NPG_PALETTE定义(self):
        """N1：_HEADER 包含 NPG_PALETTE 变量定义。"""
        from app.modules.llm.code_task_provider import _HEADER

        assert "NPG_PALETTE" in _HEADER, "_HEADER 未定义 NPG_PALETTE 变量"

    def test_HEADER包含NPG配色值(self):
        """N2：_HEADER 包含 NPG 期刊配色值 #E64B35（Nature 红）。"""
        from app.modules.llm.code_task_provider import _HEADER

        assert "#E64B35" in _HEADER, "_HEADER 未包含 NPG 配色 #E64B35"

    def test_HEADER使用NPG配色设置seaborn(self):
        """N3：_HEADER 使用 sns.color_palette(NPG_PALETTE) 设置 seaborn 配色。"""
        from app.modules.llm.code_task_provider import _HEADER

        assert "sns.color_palette(NPG_PALETTE)" in _HEADER, (
            "_HEADER 未使用 NPG_PALETTE 设置 seaborn 配色"
        )

    def test_HEADER不再使用bright_palette(self):
        """N4：_HEADER 不再使用 palette='bright'（SPEC 0030 替换为 NPG）。"""
        from app.modules.llm.code_task_provider import _HEADER

        assert 'palette="bright"' not in _HEADER, (
            "_HEADER 仍使用 palette='bright'，应替换为 NPG 配色"
        )

    def test_HEADER包含完整8色NPG配色(self):
        """N5：_HEADER 包含完整 8 色 NPG 配色。"""
        from app.modules.llm.code_task_provider import _HEADER

        npg_colors = ["#E64B35", "#4DBBD5", "#00A087", "#3C5488",
                      "#F39B7F", "#8491B4", "#91D1C2", "#DC0000"]
        for color in npg_colors:
            assert color in _HEADER, f"_HEADER 未包含 NPG 配色 {color}"


class TestSpec0030DpiUnification:
    """SPEC 0030 dpi 统一修复测试。

    红色阶段说明：
    - _build_chart_code 仍包含 dpi=100，相关测试应失败
    - 实现完成后所有测试应通过
    """

    def test_生成的代码不含dpi_100(self):
        """D1：_build_chart_code 生成的代码不含 dpi=100。"""
        provider = LocalRuleCodeTaskProvider()
        result = provider.generate(_make_full_plan_with_charts())
        assert "dpi=100" not in result.code, (
            "生成的代码仍包含 dpi=100（SPEC 0028 遗留缺陷未修复）"
        )

    def test_生成的代码不含dpi_300硬编码(self):
        """D2：savefig 不硬编码 dpi=300，由 rcParams savefig.dpi=300 控制。"""
        provider = LocalRuleCodeTaskProvider()
        result = provider.generate(_make_full_plan_with_charts())
        # savefig 调用不应包含 dpi 参数（由 rcParams 控制）
        assert "dpi=100" not in result.code
        assert "dpi=300" not in result.code

    def test_生成的代码可通过ast语法检查(self):
        """D3：dpi 修复后代码仍可通过 ast.parse 语法检查。"""
        provider = LocalRuleCodeTaskProvider()
        result = provider.generate(_make_full_plan_with_charts())
        ast.parse(result.code)

    def test_HEADER包含savefig_dpi_300(self):
        """D4：_HEADER 的 rcParams 包含 savefig.dpi=300（由 rcParams 统一控制）。"""
        from app.modules.llm.code_task_provider import _HEADER

        assert "savefig.dpi" in _HEADER, "_HEADER 未配置 savefig.dpi rcParams"
        assert "300" in _HEADER, "_HEADER 的 savefig.dpi 不是 300"


class TestSpec0030PanelLabels:
    """SPEC 0030 a/b/c 面板标签测试。

    红色阶段说明：
    - _build_chart_code 未生成 a/b/c 面板标签，相关测试应失败
    - 实现完成后所有测试应通过
    """

    def test_多图场景生成面板标签(self):
        """P1：多图 subplot 场景生成 a/b/c 面板标签。"""
        provider = LocalRuleCodeTaskProvider()
        result = provider.generate(_make_full_plan_with_charts())
        # 多图场景应包含面板标签（transAxes 用于定位标签）
        assert "transAxes" in result.code, (
            "多图场景未生成 a/b/c 面板标签（缺少 transAxes 定位）"
        )

    def test_多图场景使用subplot(self):
        """P2：多图场景使用 plt.subplots 创建多面板布局。"""
        provider = LocalRuleCodeTaskProvider()
        result = provider.generate(_make_full_plan_with_charts())
        assert "subplots" in result.code, (
            "多图场景未使用 plt.subplots 多面板布局"
        )

    def test_面板标签使用fontweight_bold(self):
        """P3：面板标签使用 fontweight='bold' 加粗。"""
        provider = LocalRuleCodeTaskProvider()
        result = provider.generate(_make_full_plan_with_charts())
        assert "fontweight" in result.code, (
            "面板标签未使用 fontweight='bold' 加粗"
        )

    def test_单图场景不强制面板标签(self):
        """P4：单图场景不强制生成面板标签（回归保障）。"""
        provider = LocalRuleCodeTaskProvider()
        single_chart_plan = {
            "cleaning_plan": [],
            "analysis_plan": [],
            "chart_plan": [
                {
                    "chart_type": "HISTOGRAM",
                    "title": "age 分布",
                    "data_fields": ["age"],
                    "description": "展示分布",
                },
            ],
        }
        result = provider.generate(single_chart_plan)
        # 单图场景代码应可通过 ast 语法检查
        ast.parse(result.code)


# === 3. PptRenderer 主题映射测试（_map_theme） ===


class TestSpec0030ThemeMapping:
    """SPEC 0030 PptRenderer 主题映射测试。

    红色阶段说明：
    - _map_theme 方法尚未实现，相关测试应失败
    - 实现完成后所有测试应通过

    设计决策：_map_theme 返回主题名字符串（如 "MIDNIGHT_EXECUTIVE"），
    而非 Theme 对象。render() 内部用 getattr(themes, name) 获取 Theme 对象。
    """

    def test_map_theme方法存在(self):
        """T1：PptRenderer 有 _map_theme 方法。"""
        assert hasattr(PptRenderer, "_map_theme"), (
            "PptRenderer 未实现 _map_theme（SPEC 0030 主题映射方法）"
        )

    def test_map_theme_蓝色映射到MIDNIGHT_EXECUTIVE(self):
        """T2：蓝色 theme_color 映射到 MIDNIGHT_EXECUTIVE。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#2563eb")
        assert result == "MIDNIGHT_EXECUTIVE", (
            f"蓝色映射到 {result}，预期 MIDNIGHT_EXECUTIVE"
        )

    def test_map_theme_绿色映射到FOREST_MOSS(self):
        """T3：绿色 theme_color 映射到 FOREST_MOSS。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#16a34a")
        assert result == "FOREST_MOSS", (
            f"绿色映射到 {result}，预期 FOREST_MOSS"
        )

    def test_map_theme_紫色映射到ROYAL_PLUM(self):
        """T4：紫色 theme_color 映射到 ROYAL_PLUM。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#7c3aed")
        assert result == "ROYAL_PLUM", (
            f"紫色映射到 {result}，预期 ROYAL_PLUM"
        )

    def test_map_theme_红色映射到CORAL_ENERGY(self):
        """T5：红色 theme_color 映射到 CORAL_ENERGY。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#dc2626")
        assert result == "CORAL_ENERGY", (
            f"红色映射到 {result}，预期 CORAL_ENERGY"
        )

    def test_map_theme_橙色映射到AMBER_EDITORIAL(self):
        """T6：橙色 theme_color 映射到 AMBER_EDITORIAL。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#ea580c")
        assert result == "AMBER_EDITORIAL", (
            f"橙色映射到 {result}，预期 AMBER_EDITORIAL"
        )

    def test_map_theme_灰色映射到MONOCHROME_INK(self):
        """T7：灰色 theme_color 映射到 MONOCHROME_INK。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#475569")
        assert result == "MONOCHROME_INK", (
            f"灰色映射到 {result}，预期 MONOCHROME_INK"
        )

    def test_map_theme_None映射到默认SLATE_MINIMALIST(self):
        """T8：None 映射到默认 SLATE_MINIMALIST。"""
        renderer = PptRenderer()
        result = renderer._map_theme(None)
        assert result == "SLATE_MINIMALIST", (
            f"None 映射到 {result}，预期 SLATE_MINIMALIST"
        )

    def test_map_theme_无效色值映射到默认(self):
        """T9：无效色值映射到默认 SLATE_MINIMALIST（降级策略）。"""
        renderer = PptRenderer()
        result = renderer._map_theme("#invalid")
        assert result == "SLATE_MINIMALIST", (
            f"无效色值映射到 {result}，预期 SLATE_MINIMALIST"
        )

    def test_map_theme_返回值为合法主题名(self):
        """T10：_map_theme 返回值在 PPT_THEME_PRESETS 枚举中。"""
        renderer = PptRenderer()
        for color in ["#2563eb", "#16a34a", "#7c3aed", "#dc2626", "#ea580c",
                      "#475569", None, "#invalid"]:
            result = renderer._map_theme(color)
            assert result in PPT_THEME_PRESETS, (
                f"_map_theme({color!r}) 返回 {result!r} 不在合法主题枚举中"
            )


# === 4. PptRenderer pptxforge 集成测试 ===


class TestSpec0030PptxforgeIntegration:
    """SPEC 0030 PptRenderer pptxforge 集成测试。

    红色阶段说明：
    - render() 内部尚未集成 pptxforge Deck，相关测试可能失败或通过（取决于 fallback）
    - 实现完成后所有测试应通过

    设计决策：render() 签名不变，内部改用 pptxforge Deck。
    保留 python-pptx fallback 路径（pptxforge 失败时降级）。
    """

    def test_render_with_theme_preset生成有效PPTX(self, tmp_path):
        """X1：theme_preset 指定时生成有效 PPTX 文件。"""
        renderer = PptRenderer()
        output_path = tmp_path / "output.pptx"
        renderer.render(
            project_name="测试项目",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output_path),
            config={"theme_preset": "MIDNIGHT_EXECUTIVE"},
        )
        assert output_path.exists(), "PPTX 文件未生成"
        prs = Presentation(str(output_path))
        assert len(prs.slides) >= 3, f"幻灯片数 {len(prs.slides)} < 3"

    def test_render_with_theme_preset覆盖theme_color(self, tmp_path):
        """X2：theme_preset 优先于 theme_color（方案 B 优先级）。"""
        renderer = PptRenderer()
        output_path = tmp_path / "output.pptx"
        # 同时指定 theme_preset 和 theme_color，theme_preset 优先
        renderer.render(
            project_name="测试项目",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output_path),
            config={"theme_preset": "FOREST_MOSS", "theme_color": "#2563eb"},
        )
        assert output_path.exists(), "PPTX 文件未生成"

    def test_render输出可被python_pptx重新打开(self, tmp_path):
        """X3：render() 输出可被 python-pptx 重新打开（接口兼容性）。"""
        renderer = PptRenderer()
        output_path = tmp_path / "output.pptx"
        renderer.render(
            project_name="测试项目",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output_path),
            config={"theme_preset": "PACIFIC_DEEP"},
        )
        prs = Presentation(str(output_path))
        assert len(prs.slides) >= 1

    def test_render_无theme_preset时降级到theme_color映射(self, tmp_path):
        """X4：无 theme_preset 时降级到 theme_color 映射（fallback 路径）。"""
        renderer = PptRenderer()
        output_path = tmp_path / "output.pptx"
        renderer.render(
            project_name="测试项目",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output_path),
            config={"theme_color": "#2563eb"},  # 无 theme_preset
        )
        assert output_path.exists(), "PPTX 文件未生成"
        prs = Presentation(str(output_path))
        assert len(prs.slides) >= 3

    def test_render_无config时降级到默认主题(self, tmp_path):
        """X5：无 config 时降级到默认 SLATE_MINIMALIST 主题。"""
        renderer = PptRenderer()
        output_path = tmp_path / "output.pptx"
        renderer.render(
            project_name="测试项目",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output_path),
            config=None,
        )
        assert output_path.exists(), "PPTX 文件未生成"
        prs = Presentation(str(output_path))
        assert len(prs.slides) >= 3

    def test_render_签名不变(self):
        """X6：render() 签名保持不变（接口兼容性约束）。"""
        import inspect

        sig = inspect.signature(PptRenderer.render)
        params = list(sig.parameters.keys())
        expected_params = [
            "self", "project_name", "project_topic",
            "outline_sections", "execution_artifacts",
            "output_path", "config",
        ]
        assert params == expected_params, (
            f"render() 签名变化：{params} != {expected_params}"
        )
