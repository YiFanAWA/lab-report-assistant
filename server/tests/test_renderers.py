"""Word 和 PPT 渲染器测试。

验证：
- WordRenderer 从已确认大纲生成 .docx 文件
- PptRenderer 从同一份大纲生成 .pptx 文件
- 文件实际生成且可被对应库重新打开
- 执行产物（CSV/PNG）正确嵌入
- 渲染失败抛出 AppError 结构化错误
"""

import csv
import hashlib
from pathlib import Path
import zipfile

import pytest
from docx import Document
from pptx import Presentation

from app.core.errors import AppError
from app.infrastructure.renderers.word_renderer import WordRenderer
from app.infrastructure.renderers.ppt_renderer import PptRenderer


SAMPLE_SECTIONS = [
    {
        "id": "sec_001",
        "title": "实验目的",
        "content": "分析胃病数据分布特征",
        "source_type": "REQUIREMENT",
        "source_ids": ["plan_001"],
    },
    {
        "id": "sec_002",
        "title": "实验背景",
        "content": "胃病发病率近年上升",
        "source_type": "EVIDENCE",
        "source_ids": ["card_001"],
    },
    {
        "id": "sec_003",
        "title": "数据描述",
        "content": "数据集规模：100 行 × 3 列",
        "source_type": "DATASET",
        "source_ids": ["ver_001"],
    },
    {
        "id": "sec_004",
        "title": "分析方案",
        "content": "清洗方案：去除缺失值\n分析方案：描述性统计",
        "source_type": "ANALYSIS",
        "source_ids": ["plan_a"],
    },
    {
        "id": "sec_005",
        "title": "实验结果",
        "content": "执行成功，输出统计表",
        "source_type": "EXECUTION",
        "source_ids": ["run_001"],
    },
    {
        "id": "sec_006",
        "title": "结论与讨论",
        "content": "本实验完成既定分析目标。",
        "source_type": "SUMMARY",
        "source_ids": [],
    },
]


def _make_csv(path: Path, rows: list[list[str]]) -> str:
    """写入 CSV 文件，返回绝对路径。"""
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        for row in rows:
            writer.writerow(row)
    return str(path)


def _make_png(path: Path) -> str:
    """写入最小 PNG 文件，返回绝对路径。"""
    from PIL import Image
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGBA", (10, 10), (255, 0, 0, 255))
    img.save(str(path), format="PNG")
    return str(path)


# --- WordRenderer 测试 ---


class TestWordRenderer:
    """Word 渲染器测试。"""

    def test_generates_docx_file(self, tmp_path):
        """成功生成 .docx 文件且文件存在。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"

        result = renderer.render(
            project_name="胃病数据分析",
            project_topic="胃病数据分析实验报告",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        assert Path(result).exists()
        assert result.endswith(".docx")

    def test_docx_can_be_reopened(self, tmp_path):
        """生成的 docx 可被 python-docx 重新打开。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="测试项目",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        doc = Document(str(output))
        # 封面用 Title 样式（level=0），章节用 Heading 1 样式
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "测试课题" in all_text
        assert "实验目的" in all_text
        assert "结论与讨论" in all_text

    def test_cover_contains_project_name(self, tmp_path):
        """封面包含项目名称。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="封面项目",
            project_topic="封面课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        doc = Document(str(output))
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "封面项目" in all_text

    def test_section_content_rendered(self, tmp_path):
        """章节内容写入文档。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        doc = Document(str(output))
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "分析胃病数据分布特征" in all_text
        assert "胃病发病率近年上升" in all_text

    def test_csv_artifact_embedded_as_table(self, tmp_path):
        """EXECUTION 章节关联的 CSV 产物嵌入为表格。"""
        csv_path = _make_csv(tmp_path / "artifacts" / "result.csv",
                             [["col_a", "col_b"], ["1", "2"]])
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[
                {
                    "name": "result.csv",
                    "artifact_type": "TABLE_CSV",
                    "file_path": csv_path,
                    "execution_run_id": "run_001",
                },
            ],
            output_path=str(output),
        )

        doc = Document(str(output))
        # 应至少有一个表格
        assert len(doc.tables) >= 1
        assert doc.tables[0].rows[0].cells[0].text == "col_a"

    def test_png_artifact_embedded_as_image(self, tmp_path):
        """EXECUTION 章节关联的 PNG 产物嵌入为图片。"""
        png_path = _make_png(tmp_path / "artifacts" / "chart.png")
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[
                {
                    "name": "chart.png",
                    "artifact_type": "CHART_PNG",
                    "file_path": png_path,
                    "execution_run_id": "run_001",
                },
            ],
            output_path=str(output),
        )

        doc = Document(str(output))
        # 图片以 inline shape 形式存在
        assert len(doc.inline_shapes) >= 1

    def test_appendix_contains_artifact_index(self, tmp_path):
        """附录包含执行产物索引。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[
                {
                    "name": "chart.png",
                    "artifact_type": "CHART_PNG",
                    "file_path": "/tmp/chart.png",
                    "execution_run_id": "run_001",
                },
            ],
            output_path=str(output),
        )

        doc = Document(str(output))
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "附录" in all_text
        assert "chart.png" in all_text

    def test_creates_output_directory(self, tmp_path):
        """输出目录不存在时自动创建。"""
        renderer = WordRenderer()
        output = tmp_path / "deep" / "nested" / "dir" / "out.docx"

        result = renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        assert Path(result).exists()

    def test_empty_sections_renders_without_error(self, tmp_path):
        """空章节列表不报错。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        result = renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=[],
            execution_artifacts=[],
            output_path=str(output),
        )
        assert Path(result).exists()

    def test_missing_csv_file_uses_placeholder_text(self, tmp_path):
        """CSV 文件不存在时写入占位文本而非抛错。"""
        renderer = WordRenderer()
        output = tmp_path / "out.docx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[
                {
                    "name": "missing.csv",
                    "artifact_type": "TABLE_CSV",
                    "file_path": "/nonexistent/missing.csv",
                    "execution_run_id": "run_001",
                },
            ],
            output_path=str(output),
        )
        doc = Document(str(output))
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "表格文件不存在" in all_text

    def test_formal_academic_profile_renders_paper_structure(self, tmp_path):
        """正式论文 profile 输出章节、题名、引用和参考文献。"""
        renderer = WordRenderer()
        output = tmp_path / "formal.docx"
        png_path = _make_png(tmp_path / "artifacts" / "formal-chart.png")
        renderer.render(
            project_name="论文项目",
            project_topic="论文主题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[
                {
                    "name": "formal-chart.png",
                    "artifact_type": "CHART_PNG",
                    "file_path": png_path,
                    "execution_run_id": "run_001",
                },
            ],
            output_path=str(output),
            config={
                "document_profile": "formal_academic",
                "formal_title": "正式论文题名",
                "formal_subtitle": "论文副标题",
                "formal_metadata": {"研究类型": "教学性分析"},
                "abstract": "这是正式论文摘要。",
                "reference_catalog": {
                    "card_001": "作者甲等. 研究论文. 2024;1:1.",
                    "run_001": "项目执行记录. 2026.",
                },
            },
        )

        doc = Document(str(output))
        all_text = "\n".join(p.text for p in doc.paragraphs)
        assert "正式论文题名" in all_text
        assert "摘要" in all_text
        assert "第 1 章  绪论" in all_text
        assert "第 3 章  结果" in all_text
        assert "参考文献" in all_text
        assert "[1] 作者甲等. 研究论文. 2024;1:1." in all_text
        assert "[2] 项目执行记录. 2026." in all_text
        assert not any(p.text == "图表与统计表" and p.style.name == "Heading 2" for p in doc.paragraphs)


# --- PptRenderer 测试 ---


class TestPptRenderer:
    """PPT 渲染器测试。"""

    def test_generates_pptx_file(self, tmp_path):
        """成功生成 .pptx 文件。"""
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"

        result = renderer.render(
            project_name="胃病数据分析",
            project_topic="胃病数据分析实验报告",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        assert Path(result).exists()
        assert result.endswith(".pptx")

    def test_word_scientific_asset_source_contains_traceability(self, tmp_path):
        image_path = _make_png(tmp_path / "schematic.png")
        output = tmp_path / "scientific.docx"
        WordRenderer().render(
            project_name="p",
            project_topic="t",
            outline_sections=[{
                "title": "流程",
                "content": "科研示意图",
                "source_type": "EXECUTION",
                "source_ids": ["run-spec0042"],
            }],
            execution_artifacts=[{
                "name": "科研流程",
                "artifact_type": "CHART_PNG",
                "file_path": str(image_path),
                "execution_run_id": "run-spec0042",
                "figure_note": "来源：本地执行",
                "scientific_asset_ids": ["bioicons-cc0-algorithm"],
                "scientific_asset_attributions": ["Bioicons / CC0-1.0"],
                "scientific_asset_image_sha256": "b" * 64,
                "scientific_asset_render_metadata": "figure.json",
            }],
            output_path=str(output),
        )
        text = "\n".join(paragraph.text for paragraph in Document(output).paragraphs)
        assert "bioicons-cc0-algorithm" in text
        assert "Bioicons / CC0-1.0" in text
        assert "图像 SHA-256" in text
        assert "figure.json" in text
    def test_pptx_can_be_reopened(self, tmp_path):
        """生成的 pptx 可被 python-pptx 重新打开。"""
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"
        renderer.render(
            project_name="p",
            project_topic="测试课题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        prs = Presentation(str(output))
        # 标题页 + 内容页 + 总结页
        assert len(prs.slides) >= 2

    def test_apply_source_notes_persists_standard_notes_xml(self, tmp_path):
        path = tmp_path / "notes.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(path)
        PptRenderer._apply_source_notes(
            path,
            [(1, "[Sources]\n- Bioicons / CC0-1.0")],
        )
        reopened = Presentation(path)
        assert reopened.slides[1].notes_slide.notes_text_frame.text == (
            "[Sources]\n- Bioicons / CC0-1.0"
        )
    def test_scientific_asset_sources_are_written_to_speaker_notes(self):
        renderer = PptRenderer()
        notes = renderer._build_source_notes((
            {
                "scientific_asset_ids": ["bioicons-cc0-algorithm"],
                "scientific_asset_attributions": ["Bioicons / CC0-1.0"],
                "scientific_asset_image_sha256": "a" * 64,
                "scientific_asset_render_metadata": "figure.json",
                "execution_run_id": "run-spec0042",
                "figure_argument": {"evidence_refs": ["dataset:UCI-296"]},
            },
        ))
        assert notes == (
            "[Sources]\n"
            "- assets: bioicons-cc0-algorithm\n"
            f"- image-sha256: {'a' * 64}\n"
            "- render-metadata: figure.json\n"
            "- execution-run: run-spec0042\n"
            "- Bioicons / CC0-1.0\n"
            "- dataset:UCI-296"
        )
    def test_long_chinese_cover_title_is_split_for_pptxforge(self):
        renderer = PptRenderer()
        title = "开放许可科研图形资产库与科研示意图组件系统"
        fitted = renderer._fit_cover_title(title)
        assert "\n" in fitted
        assert fitted.replace("\n", "") == title
    def test_title_slide_contains_topic(self, tmp_path):
        """标题页包含课题（SPEC 0024：空白版式，遍历 shapes 找文本）。"""
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"
        renderer.render(
            project_name="p",
            project_topic="胃病数据分析",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        prs = Presentation(str(output))
        first_slide = prs.slides[0]
        # SPEC 0024 空白版式无 title placeholder，遍历 shapes 找文本
        all_text = ""
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text
        assert "胃病数据分析" in all_text

    def test_content_groups_match_artifacts_by_outline_contract(self):
        renderer = PptRenderer()
        artifact_a = {
            "name": "a.png",
            "artifact_type": "CHART_PNG",
            "execution_run_id": "run-a",
            "artifact_group": "group-a",
        }
        artifact_b = {
            "name": "b.png",
            "artifact_type": "CHART_PNG",
            "execution_run_id": "run-b",
            "artifact_group": "group-b",
        }
        groups = renderer._build_content_groups(
            [{
                "title": "结果 B",
                "content": "B",
                "source_type": "EXECUTION",
                "source_ids": ["run-b"],
                "artifact_group": "group-b",
            }],
            [artifact_a, artifact_b],
        )
        assert groups[0]["chart"] is artifact_b

    def test_python_pptx_fallback_persists_scientific_asset_sources(self, tmp_path):
        image_path = _make_png(tmp_path / "schematic.png")
        output = tmp_path / "fallback.pptx"
        renderer = PptRenderer()
        renderer._render_with_python_pptx(
            project_name="p",
            project_topic="t",
            outline_sections=[{
                "title": "流程",
                "content": "科研流程",
                "source_type": "EXECUTION",
                "source_ids": ["run-spec0042"],
                "artifact_group": "scientific",
            }],
            execution_artifacts=[{
                "name": "科研流程",
                "artifact_type": "CHART_PNG",
                "file_path": str(image_path),
                "execution_run_id": "run-spec0042",
                "artifact_group": "scientific",
                "scientific_asset_ids": ["bioicons-cc0-algorithm"],
                "scientific_asset_attributions": ["Bioicons / CC0-1.0"],
                "scientific_asset_image_sha256": "c" * 64,
                "scientific_asset_render_metadata": "figure.json",
            }],
            output_path=str(output),
            target_slide_count=None,
            theme_color=None,
            include_charts=True,
        )
        reopened = Presentation(output)
        notes = reopened.slides[1].notes_slide.notes_text_frame.text
        assert "bioicons-cc0-algorithm" in notes
        assert "run-spec0042" in notes
        assert "figure.json" in notes

    def test_word_and_ppt_embed_the_same_scientific_png_bytes(self, tmp_path):
        image_path = Path(_make_png(tmp_path / "schematic.png"))
        image_sha256 = hashlib.sha256(image_path.read_bytes()).hexdigest()
        section = {
            "title": "流程",
            "content": "科研流程",
            "source_type": "EXECUTION",
            "source_ids": ["run-spec0042"],
            "artifact_group": "scientific",
        }
        artifact = {
            "name": "科研流程",
            "artifact_type": "CHART_PNG",
            "file_path": str(image_path),
            "execution_run_id": "run-spec0042",
            "artifact_group": "scientific",
            "scientific_asset_ids": ["bioicons-cc0-algorithm"],
            "scientific_asset_attributions": ["Bioicons / CC0-1.0"],
            "scientific_asset_image_sha256": image_sha256,
            "scientific_asset_render_metadata": "figure.json",
        }
        word_path = tmp_path / "same-source.docx"
        ppt_path = tmp_path / "same-source.pptx"
        WordRenderer().render(
            project_name="p",
            project_topic="t",
            outline_sections=[section],
            execution_artifacts=[artifact],
            output_path=str(word_path),
        )
        PptRenderer()._render_with_python_pptx(
            project_name="p",
            project_topic="t",
            outline_sections=[section],
            execution_artifacts=[artifact],
            output_path=str(ppt_path),
            target_slide_count=None,
            theme_color=None,
            include_charts=True,
        )

        def embedded_png_hashes(path: Path, media_prefix: str) -> set[str]:
            with zipfile.ZipFile(path) as archive:
                return {
                    hashlib.sha256(archive.read(name)).hexdigest()
                    for name in archive.namelist()
                    if name.startswith(media_prefix) and name.lower().endswith(".png")
                }

        assert image_sha256 in embedded_png_hashes(word_path, "word/media/")
        assert image_sha256 in embedded_png_hashes(ppt_path, "ppt/media/")
        assert image_sha256 in "\n".join(
            paragraph.text for paragraph in Document(word_path).paragraphs
        )
        assert image_sha256 in Presentation(ppt_path).slides[1].notes_slide.notes_text_frame.text

    def test_content_slides_from_sections(self, tmp_path):
        """内容页按 source_type 分组渲染。"""
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        prs = Presentation(str(output))
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)

        # 课题与问题页应包含实验目的
        joined = "\n".join(all_text)
        assert "课题与问题" in joined or "实验目的" in joined
        # 方法与数据页应包含数据描述
        assert "方法与数据" in joined or "数据描述" in joined

    def test_summary_slide_present(self, tmp_path):
        """总结页存在（SPEC 0024：空白版式，遍历 shapes 找文本）。

        SPEC 0030：pptxforge 主路径使用 closing_slide 承载总结内容，
        检查总结内容文本而非 "总结" 标题词（pptxforge closing_slide 无标题栏）。
        """
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )

        prs = Presentation(str(output))
        last_slide = prs.slides[len(prs.slides) - 1]
        # SPEC 0024 空白版式无 title placeholder，遍历 shapes 找文本
        all_text = ""
        for shape in last_slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text
        # SPEC 0030：pptxforge closing_slide 承载总结内容（"本实验完成既定分析目标。"）
        # python-pptx 降级路径承载 "总结" 标题。两者均合规。
        assert "总结" in all_text or "本实验完成既定分析目标" in all_text, (
            f"总结页未包含总结文本：{all_text!r}"
        )

    def test_chart_slide_with_png_artifact(self, tmp_path):
        """有 PNG 产物时生成关键图表页。"""
        png_path = _make_png(tmp_path / "artifacts" / "chart.png")
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"
        renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[
                {
                    "name": "chart.png",
                    "artifact_type": "CHART_PNG",
                    "file_path": png_path,
                    "execution_run_id": "run_001",
                },
            ],
            output_path=str(output),
        )

        prs = Presentation(str(output))
        # 应找到包含图片的幻灯片
        has_picture = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    has_picture = True
                    break
            if has_picture:
                break
        assert has_picture

    def test_empty_sections_renders_without_error(self, tmp_path):
        """空章节列表不报错。"""
        renderer = PptRenderer()
        output = tmp_path / "out.pptx"
        result = renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=[],
            execution_artifacts=[],
            output_path=str(output),
        )
        assert Path(result).exists()

    def test_creates_output_directory(self, tmp_path):
        """输出目录不存在时自动创建。"""
        renderer = PptRenderer()
        output = tmp_path / "deep" / "nested" / "out.pptx"
        result = renderer.render(
            project_name="p",
            project_topic="t",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(output),
        )
        assert Path(result).exists()

    def test_academic_diagnosis_stratified_keeps_chart_visual(self, tmp_path):
        """显式诊断分层答辩页必须保留其结果图表。"""
        png_path = _make_png(tmp_path / "artifacts" / "diagnosis.png")
        sections = [
            {
                "id": "diagnosis_stratified",
                "title": "主要诊断分层与 HbA1c 交互",
                "content": "诊断分层结果用于展示描述性异质性。",
                "presentation_role": "diagnosis_stratified",
                "presentation_content": "9 类主要诊断分层；交互检验 P 0.01。",
                "figure_lead": "先阅读分层率及区间。",
                "figure_takeaway": "结果不替代完整模型或因果解释。",
                "source_type": "EXECUTION",
                "source_ids": ["run_001"],
                "artifact_group": "diagnosis_stratified",
            },
        ]
        artifact = {
            "name": "diagnosis.png",
            "artifact_type": "CHART_PNG",
            "file_path": png_path,
            "execution_run_id": "run_001",
            "artifact_group": "diagnosis_stratified",
            "chart_kind": "bar",
        }
        output = tmp_path / "diagnosis.pptx"

        PptRenderer().render(
            project_name="论文解读",
            project_topic="诊断分层复核",
            outline_sections=sections,
            execution_artifacts=[artifact],
            output_path=str(output),
            config={"ppt_workflow": "academic"},
        )

        prs = Presentation(str(output))
        pictures = [
            shape
            for shape in prs.slides[1].shapes
            if shape.shape_type == 13
        ]
        assert pictures, "诊断分层答辩页必须保留图表主视觉"
