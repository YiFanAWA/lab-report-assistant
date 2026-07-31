"""SPEC 0025 验收验证脚本：生成 6 种预设色 PPT 并提取色彩/布局信息。

生成内容：
- 6 种预设色（蓝/紫/绿/红/橙/灰）的 PPT 文件
- 每种色的三角色彩派生结果
- 三明治结构 shape 信息（标题栏/页脚栏/辅助色背景位置和颜色）
- 对比说明文本
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

# 添加 server 目录到 path
server_root = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(server_root))

from app.infrastructure.renderers.ppt_renderer import PptRenderer

# 6 种预设色
PRESET_COLORS = {
    "blue": "#2563eb",
    "purple": "#7c3aed",
    "green": "#16a34a",
    "red": "#dc2626",
    "orange": "#ea580c",
    "gray": "#475569",
}

# 示例大纲数据
SAMPLE_SECTIONS = [
    {
        "id": "sec_001",
        "title": "实验目的",
        "content": "分析胃病数据分布特征，探索不同诊断类型的检验指标差异",
        "source_type": "REQUIREMENT",
        "source_ids": ["plan_001"],
    },
    {
        "id": "sec_002",
        "title": "实验背景",
        "content": "胃病发病率近年上升，早期诊断和数据分析对医学研究具有重要意义",
        "source_type": "EVIDENCE",
        "source_ids": ["card_001"],
    },
    {
        "id": "sec_003",
        "title": "数据描述",
        "content": "数据集规模：200 行 × 15 列，包含血常规、生化指标、肿瘤标志物",
        "source_type": "DATASET",
        "source_ids": ["ver_001"],
    },
    {
        "id": "sec_004",
        "title": "分析方案",
        "content": "描述性统计 + 分组对比 + 相关性分析",
        "source_type": "ANALYSIS",
        "source_ids": ["plan_a"],
    },
    {
        "id": "sec_005",
        "title": "实验结果",
        "content": "执行成功，生成诊断分布图、年龄分布直方图等 4 张图表",
        "source_type": "EXECUTION",
        "source_ids": ["run_001"],
    },
    {
        "id": "sec_006",
        "title": "结论与讨论",
        "content": "本实验完成既定分析目标，各诊断类型指标差异显著。",
        "source_type": "SUMMARY",
        "source_ids": [],
    },
]


def emu_to_inches(emu):
    """将 EMU 转为英寸。"""
    if emu is None:
        return 0.0
    return round(Emu(emu).inches, 2)


def rgb_to_hex(rgb):
    """将 RGBColor 转为 hex 字符串。"""
    return str(rgb).upper()


def extract_slide_info(pptx_path):
    """提取 PPT 中每页的 shape 信息。"""
    prs = Presentation(str(pptx_path))
    slides_info = []

    for idx, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            info = {
                "shape_type": str(shape.shape_type),
                "name": shape.name,
                "left": emu_to_inches(shape.left),
                "top": emu_to_inches(shape.top),
                "width": emu_to_inches(shape.width),
                "height": emu_to_inches(shape.height),
            }
            # 提取填充色
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    fill_color = shape.fill.fore_color
                    if fill_color.type is not None:
                        info["fill_color"] = str(fill_color.rgb).upper()
                except Exception:
                    pass
            # 提取文字
            if shape.has_text_frame:
                text = shape.text_frame.text[:80]
                if text:
                    info["text"] = text
                # 提取文字颜色
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            font_color = run.font.color
                            if font_color.type is not None:
                                info["font_color"] = str(font_color.rgb).upper()
                                break
                        except Exception:
                            pass
                    if "font_color" in info:
                        break
            shapes_info.append(info)
        slides_info.append({
            "slide_index": idx,
            "shapes_count": len(slide.shapes),
            "shapes": shapes_info,
        })
    return slides_info


def main():
    output_dir = server_root.parent / "dev-docs" / "e2e-screenshots" / "spec0025"
    output_dir.mkdir(parents=True, exist_ok=True)

    renderer = PptRenderer()
    results = {}

    for color_name, theme_color in PRESET_COLORS.items():
        print(f"\n=== 生成 {color_name} 主题色 PPT ({theme_color}) ===")

        # 派生色彩
        theme_rgb = renderer._resolve_theme_color(theme_color)
        primary, auxiliary, accent, title_text_color = (
            renderer._derive_color_palette(theme_rgb)
        )

        color_info = {
            "theme_color": theme_color,
            "primary": rgb_to_hex(primary),
            "auxiliary": rgb_to_hex(auxiliary),
            "accent": rgb_to_hex(accent),
            "title_text_color": rgb_to_hex(title_text_color),
        }
        print(f"  派生色彩: {color_info}")

        # 生成 PPT
        pptx_path = output_dir / f"spec0025-{color_name}.pptx"
        renderer.render(
            project_name=f"胃病数据分析-{color_name}",
            project_topic="胃病数据分析实验报告",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=[],
            output_path=str(pptx_path),
            config={"theme_color": theme_color},
        )
        print(f"  PPT 文件: {pptx_path}")

        # 提取 shape 信息
        slides_info = extract_slide_info(pptx_path)
        print(f"  幻灯片数: {len(slides_info)}")

        results[color_name] = {
            "color_info": color_info,
            "pptx_file": str(pptx_path.name),
            "slides": slides_info,
        }

    # 保存 JSON 结果
    json_path = output_dir / "spec0025-verification.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    print(f"\n=== 验证结果已保存到 {json_path} ===")

    # 打印汇总表
    print("\n=== 6 种预设色派生结果汇总 ===")
    print(f"{'色彩':<8} {'theme_color':<12} {'主色':<8} {'辅助色':<8} {'强调色':<8} {'标题文字':<8}")
    print("-" * 60)
    for color_name, data in results.items():
        ci = data["color_info"]
        print(f"{color_name:<8} {ci['theme_color']:<12} {ci['primary']:<8} {ci['auxiliary']:<8} {ci['accent']:<8} {ci['title_text_color']:<8}")

    # 打印三明治结构验证
    print("\n=== 三明治结构验证（以 blue 为例） ===")
    blue_slides = results["blue"]["slides"]
    for slide in blue_slides[:3]:  # 前 3 页
        print(f"\n  幻灯片 {slide['slide_index']} ({slide['shapes_count']} shapes):")
        for s in slide["shapes"]:
            fill = s.get("fill_color", "无")
            text = s.get("text", "")[:30]
            print(f"    {s['name']:<20} fill={fill:<8} pos=({s['left']},{s['top']}) size=({s['width']}x{s['height']}) text={text}")


if __name__ == "__main__":
    main()
