"""SPEC 0027 八图表截断逻辑验证脚本。

验证内容：
1. 生成 8 张图表，渲染 PPT
2. 提取截断注释精确属性（文本、位置、颜色、字号）
3. 验证截断逻辑正确性

运行方式（从 server 目录）：
    .venv/Scripts/python.exe -m scripts.verify_8charts_truncation
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import scienceplots  # noqa: F401
plt.style.use(["science", "no-latex", "cjk-sc-font", "bright"])
import seaborn as sns
sns.set_theme(style="whitegrid", palette="bright", font="Microsoft YaHei")
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Emu

from app.infrastructure.renderers.ppt_renderer import PptRenderer


def generate_8_charts(output_dir: Path) -> list[dict]:
    """生成 8 张示例图表 PNG。"""
    output_dir.mkdir(parents=True, exist_ok=True)

    np.random.seed(42)
    n = 200
    data = {
        "age": np.random.normal(50, 15, n),
        "bmi": np.random.normal(24, 4, n),
        "symptom_score": np.random.randint(0, 10, n),
        "diagnosis": np.random.choice(["胃炎", "溃疡", "正常", "其他"], n),
        "treatment_days": np.random.randint(1, 30, n),
        "blood_pressure": np.random.normal(120, 20, n),
        "heart_rate": np.random.normal(75, 12, n),
        "cholesterol": np.random.normal(5.0, 1.2, n),
    }
    df = pd.DataFrame(data)

    chart_configs = [
        ("age_distribution", "年龄分布直方图",
         lambda: sns.histplot(data=df, x="age", kde=True, bins=20)),
        ("diagnosis_count", "诊断类别分布",
         lambda: sns.countplot(data=df, x="diagnosis")),
        ("symptom_boxplot", "症状评分箱线图",
         lambda: sns.boxplot(data=df[["symptom_score", "treatment_days"]])),
        ("age_bmi_scatter", "年龄与BMI散点图",
         lambda: sns.scatterplot(data=df, x="age", y="bmi")),
        ("bp_distribution", "血压分布直方图",
         lambda: sns.histplot(data=df, x="blood_pressure", kde=True, bins=20)),
        ("corr_heatmap", "相关性热图",
         lambda: sns.heatmap(df.corr(numeric_only=True), annot=True, cmap="coolwarm", fmt=".2f")),
        ("heart_rate_hist", "心率分布直方图",
         lambda: sns.histplot(data=df, x="heart_rate", kde=True, bins=20)),
        ("age_diagnosis_violin", "年龄与诊断小提琴图",
         lambda: sns.violinplot(data=df, x="diagnosis", y="age")),
    ]

    charts = []
    for name, title, plot_fn in chart_configs:
        plt.figure(figsize=(8, 5))
        plot_fn()
        plt.title(title)
        plt.tight_layout()
        file_path = output_dir / f"{name}.png"
        plt.savefig(file_path, dpi=100, bbox_inches="tight")
        plt.close()
        charts.append({
            "name": title,
            "file_path": str(file_path),
            "artifact_type": "CHART_PNG",
        })

    return charts


def extract_truncation_note(ppt_path: str) -> dict | None:
    """从 PPT 中提取截断注释的精确属性。"""
    prs = Presentation(ppt_path)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if "共" in text and "图表" in text and "已展示前" in text:
                info = {
                    "slide_idx": slide_idx + 1,
                    "text": text,
                    "left_inches": round(Emu(shape.left).inches, 3),
                    "top_inches": round(Emu(shape.top).inches, 3),
                    "width_inches": round(Emu(shape.width).inches, 3),
                    "height_inches": round(Emu(shape.height).inches, 3),
                    "font_size_pt": None,
                    "font_color_hex": None,
                    "alignment": None,
                }
                for para in shape.text_frame.paragraphs:
                    if para.alignment is not None:
                        info["alignment"] = str(para.alignment)
                    for run in para.runs:
                        if run.font.size is not None:
                            info["font_size_pt"] = round(run.font.size.pt, 1)
                        if run.font.color is not None and run.font.color.rgb is not None:
                            info["font_color_hex"] = str(run.font.color.rgb)
                        break
                    break
                return info
    return None


def extract_all_slides_summary(ppt_path: str) -> list[dict]:
    """提取所有幻灯片的图表数量和注释信息。"""
    prs = Presentation(ppt_path)
    slides_info = []
    for slide_idx, slide in enumerate(prs.slides):
        chart_count = sum(1 for s in slide.shapes if s.shape_type == 13)
        note = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "共" in text and "图表" in text and "已展示前" in text:
                    note = text
        shapes_count = len(slide.shapes)
        slides_info.append({
            "slide_idx": slide_idx + 1,
            "shapes_count": shapes_count,
            "chart_count": chart_count,
            "note": note,
        })
    return slides_info


def main():
    project_root = Path(__file__).parent.parent.parent
    output_dir = project_root / "dev-docs" / "e2e-screenshots" / "spec0027"
    output_dir.mkdir(parents=True, exist_ok=True)

    print("=" * 60)
    print("SPEC 0027 八图表截断逻辑验证")
    print("=" * 60)

    # 1. 生成 8 张图表
    print("\n1. 生成 8 张示例图表（Seaborn + SciencePlots）...")
    charts_dir = output_dir / "charts_8_demo"
    charts = generate_8_charts(charts_dir)
    print(f"   生成 {len(charts)} 张图表:")
    for i, c in enumerate(charts):
        print(f"   {i+1}. {c['name']}")

    # 2. 渲染 8 图表 PPT
    print(f"\n2. 渲染 8 图表 PPT...")
    ppt_path = str(output_dir / "grid_layout_8charts.pptx")
    renderer = PptRenderer()
    renderer.render(
        project_name="八图表截断验证",
        project_topic="胃病数据分析报告",
        outline_sections=[
            {
                "title": "数据分析概述",
                "source_type": "REQUIREMENT",
                "content": "本实验对胃病数据进行分析，包括描述性统计和可视化。",
            },
            {
                "title": "总结",
                "source_type": "SUMMARY",
                "content": "通过数据分析发现年龄与BMI存在相关性，胃炎是最常见诊断类别。",
            },
        ],
        execution_artifacts=charts,
        output_path=ppt_path,
        config={"theme_color": "2563EB", "include_charts": True},
    )
    print(f"   已保存: {ppt_path}")

    # 3. 提取所有幻灯片摘要
    print(f"\n3. 提取 PPT 所有幻灯片信息...")
    slides = extract_all_slides_summary(ppt_path)
    for s in slides:
        note_str = f", 注释: \"{s['note']}\"" if s["note"] else ""
        print(f"   第{s['slide_idx']}页: {s['shapes_count']}个形状, {s['chart_count']}张图表{note_str}")

    # 4. 提取截断注释属性
    print(f"\n4. 提取截断注释精确属性...")
    note = extract_truncation_note(ppt_path)
    if note:
        print(f"   ✅ 找到截断注释:")
        print(f"      所在页:     第 {note['slide_idx']} 页")
        print(f"      文本:       \"{note['text']}\"")
        print(f"      位置:       left={note['left_inches']}\", top={note['top_inches']}\"")
        print(f"      尺寸:       {note['width_inches']}\" × {note['height_inches']}\"")
        print(f"      字号:       {note['font_size_pt']}pt")
        print(f"      颜色:       #{note['font_color_hex']}")
        print(f"      对齐:       {note['alignment']}")
    else:
        print(f"   ❌ 未找到截断注释")

    # 5. 验证结论
    print(f"\n5. 验证结论")
    print(f"   {'=' * 50}")
    expected_note = "共 7 张图表，已展示前 4 张"
    if note:
        checks = [
            ("截断注释存在", note is not None),
            ("文本正确", note["text"] == expected_note),
            ("注释总数=7（8-1=7，1张被内容页吸收）", "共 7 张" in note["text"]),
            ("展示前4张", "已展示前 4 张" in note["text"]),
            ("位置在底部 top=6.5\"", note["top_inches"] == 6.5),
            ("字号 12pt", note["font_size_pt"] == 12.0),
            ("颜色灰色 #888888", note["font_color_hex"] == "888888"),
            ("对齐居中", "CENTER" in note["alignment"]),
        ]
        all_pass = True
        for name, passed in checks:
            status = "✅" if passed else "❌"
            if not passed:
                all_pass = False
            print(f"   {status} {name}")

        print(f"\n   {'✅ 全部验证通过！' if all_pass else '❌ 存在验证失败项！'}")

    print(f"\n{'=' * 60}")
    print(f"验证完成！")
    print(f"   PPT 文件: {ppt_path}")


if __name__ == "__main__":
    main()
