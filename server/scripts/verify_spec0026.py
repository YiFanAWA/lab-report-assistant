"""SPEC 0026 真实文件验证脚本。

生成 6 种预设色的 PPT 文件，程序化验证：
- 渐变填充（封面/标题栏/页脚栏）
- 圆角矩形（左栏背景）
- 外阴影（右栏图表）
- 细边框（右栏图表）

验证完成后输出 JSON 报告。
"""

import json
import sys
from pathlib import Path

# 添加 server 目录到 path
sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from app.infrastructure.renderers.ppt_renderer import PptRenderer


# 6 种预设色
PRESET_COLORS = [
    ("蓝", "#2563eb"),
    ("紫", "#7c3aed"),
    ("绿", "#16a34a"),
    ("红", "#dc2626"),
    ("橙", "#ea580c"),
    ("灰", "#475569"),
]

# 示例大纲
SAMPLE_SECTIONS = [
    {
        "id": "sec_001",
        "title": "实验目的",
        "content": "分析胃病数据分布特征，掌握描述性统计方法在医学数据分析中的应用。",
        "source_type": "REQUIREMENT",
        "source_ids": ["plan_001"],
    },
    {
        "id": "sec_002",
        "title": "实验方法",
        "content": "使用 Python pandas 进行数据清洗，matplotlib 绘制分布图。",
        "source_type": "EVIDENCE",
        "source_ids": ["ev_001"],
    },
    {
        "id": "sec_003",
        "title": "实验结果",
        "content": "成功生成诊断分布柱状图和年龄分布直方图，数据清洗完成。",
        "source_type": "EXECUTION",
        "source_ids": ["run_001"],
    },
    {
        "id": "sec_004",
        "title": "总结",
        "content": "本实验完成了胃病数据分析与可视化，验证了描述性统计方法的有效性。",
        "source_type": "SUMMARY",
        "source_ids": [],
    },
]


def make_chart_png(path):
    """生成测试用 PNG 图表。"""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (600, 400), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    # 绘制简单柱状图
    draw.rectangle([50, 300, 150, 350], fill=(37, 99, 235))
    draw.rectangle([200, 200, 300, 350], fill=(37, 99, 235))
    draw.rectangle([350, 250, 450, 350], fill=(37, 99, 235))
    draw.line([50, 350, 550, 350], fill=(0, 0, 0), width=2)
    img.save(str(path))


def verify_ppt(pptx_path, theme_color):
    """验证单个 PPT 文件的视觉效果。"""
    prs = Presentation(str(pptx_path))
    result = {
        "slides_count": len(prs.slides),
        "gradient_found": False,
        "gradient_slides": [],
        "rounded_rect_found": False,
        "rounded_rect_slides": [],
        "shadow_found": False,
        "shadow_slides": [],
        "border_found": False,
        "border_slides": [],
        "file_valid": True,
    }

    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 检查渐变填充
            try:
                if shape.fill.type == MSO_FILL_TYPE.GRADIENT:
                    result["gradient_found"] = True
                    if idx not in result["gradient_slides"]:
                        result["gradient_slides"].append(idx)
            except Exception:
                pass

            # 检查圆角矩形
            try:
                if getattr(shape, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE:
                    result["rounded_rect_found"] = True
                    if idx not in result["rounded_rect_slides"]:
                        result["rounded_rect_slides"].append(idx)
            except Exception:
                pass

            # 检查图片阴影和边框
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 阴影
                try:
                    spPr = shape._element.spPr
                    effectLst = spPr.find(qn('a:effectLst'))
                    if effectLst is not None:
                        outerShdw = effectLst.find(qn('a:outerShdw'))
                        if outerShdw is not None:
                            result["shadow_found"] = True
                            if idx not in result["shadow_slides"]:
                                result["shadow_slides"].append(idx)
                except Exception:
                    pass

                # 边框
                try:
                    if shape.line.color.rgb is not None:
                        result["border_found"] = True
                        if idx not in result["border_slides"]:
                            result["border_slides"].append(idx)
                except Exception:
                    pass

    return result


def main():
    """主函数：生成 6 种预设色 PPT 并验证。"""
    output_dir = Path(__file__).parent.parent / "tests" / "tmp_spec0026_verify"
    output_dir.mkdir(parents=True, exist_ok=True)

    # 生成测试图表
    chart_path = output_dir / "chart.png"
    make_chart_png(chart_path)

    artifacts = [
        {"name": "chart1.png", "artifact_type": "CHART_PNG",
         "file_path": str(chart_path)},
    ]

    all_results = {}

    for color_name, theme_color in PRESET_COLORS:
        pptx_path = output_dir / f"spec0026_{color_name}.pptx"
        renderer = PptRenderer()
        renderer.render(
            project_name=f"SPEC0026验证_{color_name}",
            project_topic=f"胃病数据分析 - {color_name}主题",
            outline_sections=SAMPLE_SECTIONS,
            execution_artifacts=artifacts,
            output_path=str(pptx_path),
            config={"theme_color": theme_color},
        )

        # 验证文件可重新打开（XML 完整性）
        try:
            result = verify_ppt(pptx_path, theme_color)
        except Exception as e:
            result = {"file_valid": False, "error": str(e)}

        all_results[color_name] = result
        print(f"[{color_name}] {theme_color}: "
              f"渐变={'✓' if result.get('gradient_found') else '✗'} "
              f"圆角={'✓' if result.get('rounded_rect_found') else '✗'} "
              f"阴影={'✓' if result.get('shadow_found') else '✗'} "
              f"边框={'✓' if result.get('border_found') else '✗'} "
              f"页数={result.get('slides_count', 0)}")

    # 输出 JSON 报告
    report_path = output_dir / "spec0026_verify_report.json"
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2)

    # 汇总
    all_pass = all(
        r.get("file_valid") and r.get("gradient_found")
        and r.get("rounded_rect_found") and r.get("shadow_found")
        and r.get("border_found")
        for r in all_results.values()
    )
    print(f"\n报告已保存：{report_path}")
    print(f"总体结果：{'全部通过 ✓' if all_pass else '存在失败 ✗'}")

    return 0 if all_pass else 1


if __name__ == "__main__":
    sys.exit(main())
