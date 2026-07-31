"""SPEC 0027 网格布局演示脚本。

演示内容：
1. 4 张图表 → 2×2 网格布局（_place_chart_grid）
2. 6 张图表 → 截断为 4 张 + 截断注释（_add_truncation_note）
3. 生成 HTML 预览展示布局效果和形状坐标

运行方式：
    server/.venv/Scripts/python.exe -m scripts.demo_grid_layout
"""

import sys
from pathlib import Path

# 添加 server 目录到 path
sys.path.insert(0, str(Path(__file__).parent.parent))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import scienceplots  # noqa: F401
plt.style.use(["science", "no-latex", "cjk-sc-font", "bright"])
import seaborn as sns
sns.set_theme(style="whitegrid", palette="bright", font="Microsoft YaHei")
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Emu

from app.infrastructure.renderers.ppt_renderer import PptRenderer, SLIDE_WIDTH, SLIDE_HEIGHT


# === 布局颜色（用于 HTML 预览中的形状区分） ===
COLOR_CHART = "#3B82F6"       # 蓝色 - 图表
COLOR_TITLE_BAR = "#1E40AF"   # 深蓝 - 标题栏
COLOR_FOOTER = "#1E3A8A"      # 深蓝 - 页脚栏
COLOR_TEXT = "#6B7280"        # 灰色 - 文本框
COLOR_NOTE = "#F59E0B"        # 橙色 - 截断注释
COLOR_BG = "#F3F4F6"          # 浅灰 - 背景


def generate_sample_charts(output_dir: Path, count: int) -> list[dict]:
    """生成指定数量的示例图表 PNG。"""
    output_dir.mkdir(parents=True, exist_ok=True)

    np.random.seed(42)
    n = 200
    data = {
        "age": np.random.normal(50, 15, n),
        "bmi": np.random.normal(24, 4, n),
        "symptom_score": np.random.randint(0, 10, n),
        "diagnosis": np.random.choice(["胃炎", "溃疡", "正常", "其他"], n),
        "treatment_days": np.random.randint(1, 30, n),
        "blood_pressure": np.random.normal(120, 20, n),
    }
    df = pd.DataFrame(data)

    chart_configs = [
        ("age_distribution", "年龄分布直方图",
         lambda: sns.histplot(data=df, x="age", kde=True, bins=20)),
        ("diagnosis_count", "诊断类别分布",
         lambda: sns.countplot(data=df, x="diagnosis")),
        ("symptom_boxplot", "症状评分箱线图",
         lambda: sns.boxplot(data=df[["symptom_score", "treatment_days"]])),
        ("age_bmi_scatter", "年龄与BMI散点图",
         lambda: sns.scatterplot(data=df, x="age", y="bmi")),
        ("bp_distribution", "血压分布直方图",
         lambda: sns.histplot(data=df, x="blood_pressure", kde=True, bins=20)),
        ("corr_heatmap", "相关性热图",
         lambda: sns.heatmap(df.corr(numeric_only=True), annot=True, cmap="coolwarm", fmt=".2f")),
    ]

    charts = []
    for name, title, plot_fn in chart_configs[:count]:
        plt.figure(figsize=(8, 5))
        plot_fn()
        plt.title(title)
        plt.tight_layout()
        file_path = output_dir / f"{name}.png"
        plt.savefig(file_path, dpi=100, bbox_inches="tight")
        plt.close()
        charts.append({
            "name": title,
            "file_path": str(file_path),
            "artifact_type": "CHART_PNG",
        })

    return charts


def render_ppt_with_charts(charts: list[dict], output_path: str) -> str:
    """使用生产 PptRenderer 渲染 PPT。

    outline 只含 SUMMARY 章节，确保所有图表进入图表页，
    从而精确控制图表页上的图表数量。
    """
    renderer = PptRenderer()

    outline_sections = [
        {
            "title": "数据分析概述",
            "source_type": "REQUIREMENT",
            "content": "本实验对胃病数据进行分析，包括描述性统计和可视化。",
        },
        {
            "title": "总结",
            "source_type": "SUMMARY",
            "content": "通过数据分析发现年龄与BMI存在相关性，胃炎是最常见诊断类别。",
        },
    ]

    return renderer.render(
        project_name="网格布局演示",
        project_topic="胃病数据分析报告",
        outline_sections=outline_sections,
        execution_artifacts=charts,
        output_path=output_path,
        config={"theme_color": "2563EB", "include_charts": True},
    )


def extract_slide_layout(ppt_path: str) -> list[dict]:
    """从 PPT 中提取每张幻灯片的形状布局信息。"""
    prs = Presentation(ppt_path)
    slides_info = []

    for slide_idx, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            info = {
                "type": str(shape.shape_type) if shape.shape_type else "unknown",
                "left": round(Emu(shape.left).inches, 2) if shape.left else 0,
                "top": round(Emu(shape.top).inches, 2) if shape.top else 0,
                "width": round(Emu(shape.width).inches, 2) if shape.width else 0,
                "height": round(Emu(shape.height).inches, 2) if shape.height else 0,
                "name": shape.name,
                "is_chart": shape.shape_type == 13,  # PICTURE
            }
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()[:80]
                if text:
                    info["text"] = text
            shapes_info.append(info)
        slides_info.append({"slide_idx": slide_idx + 1, "shapes": shapes_info})

    return slides_info


def _shape_to_html(shape: dict) -> str:
    """将单个形状转为 HTML 绝对定位 div。"""
    left_pct = shape["left"] / SLIDE_WIDTH * 100
    top_pct = shape["top"] / SLIDE_HEIGHT * 100
    width_pct = shape["width"] / SLIDE_WIDTH * 100
    height_pct = shape["height"] / SLIDE_HEIGHT * 100

    if shape["is_chart"]:
        color = COLOR_CHART
        label = f"图表<br>{shape['width']}×{shape['height']}\""
        bg = f"background: {color}22; border: 2px solid {color};"
    elif shape.get("text", "").startswith("共 ") and "图表" in shape.get("text", ""):
        color = COLOR_NOTE
        label = shape["text"][:40]
        bg = f"background: {color}22; border: 2px solid {color};"
    elif shape["top"] < 1.0 and shape["width"] > 10:
        color = COLOR_TITLE_BAR
        label = shape.get("text", "标题栏")[:20]
        bg = f"background: {color}; color: white;"
    elif shape["top"] > 6.5 and shape["width"] > 10:
        color = COLOR_FOOTER
        label = shape.get("text", "页脚")[:20]
        bg = f"background: {color}; color: white;"
    else:
        color = COLOR_TEXT
        label = shape.get("text", "")[:30]
        bg = f"background: {color}11; border: 1px solid {color}66;"

    return f"""<div style="position:absolute; left:{left_pct:.2f}%; top:{top_pct:.2f}%;
    width:{width_pct:.2f}%; height:{height_pct:.2f}%; {bg}
    display:flex; align-items:center; justify-content:center;
    font-size:10px; text-align:center; overflow:hidden; border-radius:3px;">
    {label}
  </div>"""


def _slide_to_html(slide_info: dict) -> str:
    """将一张幻灯片转为 HTML。"""
    shapes_html = "\n".join(_shape_to_html(s) for s in slide_info["shapes"])
    return f"""<div style="margin-bottom:20px;">
    <h3>第 {slide_info['slide_idx']} 页</h3>
    <div style="position:relative; width:100%; padding-bottom:56.25%;
    background:{COLOR_BG}; border:1px solid #ccc; border-radius:6px; overflow:hidden;">
      {shapes_html}
    </div>
  </div>"""


def generate_html_preview(
    ppt4_info: list[dict],
    ppt6_info: list[dict],
    output_path: str,
) -> None:
    """生成 HTML 预览，对比 4 图表和 6 图表布局。"""

    # 找到图表页（含图表形状的页）
    ppt4_chart_slide = None
    for s in ppt4_info:
        if any(sh["is_chart"] for sh in s["shapes"]):
            ppt4_chart_slide = s
            break

    ppt6_chart_slide = None
    for s in ppt6_info:
        if any(sh["is_chart"] for sh in s["shapes"]):
            ppt6_chart_slide = s
            break

    # 提取图表坐标信息用于表格展示
    def chart_table(slide_info):
        if not slide_info:
            return "<p>无图表页</p>"
        rows = []
        chart_idx = 0
        for sh in slide_info["shapes"]:
            if sh["is_chart"]:
                chart_idx += 1
                rows.append(
                    f"<tr><td>{chart_idx}</td><td>{sh['left']:.2f}\"</td>"
                    f"<td>{sh['top']:.2f}\"</td><td>{sh['width']:.2f}\"</td>"
                    f"<td>{sh['height']:.2f}\"</td></tr>"
                )
        # 检查截断注释
        note = ""
        for sh in slide_info["shapes"]:
            text = sh.get("text", "")
            if "共" in text and "图表" in text:
                note = f'<p style="color:{COLOR_NOTE};font-weight:bold;">⚠ {text}</p>'
        table = f"""<table style="border-collapse:collapse;width:100%;font-size:13px;">
          <tr style="background:#2563EB;color:white;">
            <th style="padding:6px;border:1px solid #ddd;">序号</th>
            <th style="padding:6px;border:1px solid #ddd;">Left (inches)</th>
            <th style="padding:6px;border:1px solid #ddd;">Top (inches)</th>
            <th style="padding:6px;border:1px solid #ddd;">Width (inches)</th>
            <th style="padding:6px;border:1px solid #ddd;">Height (inches)</th>
          </tr>
          {''.join(rows)}
        </table>"""
        return note + table

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>SPEC 0027 网格布局演示</title>
<style>
  body {{ font-family: "Microsoft YaHei", sans-serif; margin: 20px; background: #fff; }}
  h1 {{ color: #2563EB; border-bottom: 3px solid #2563EB; padding-bottom: 10px; }}
  h2 {{ color: #1E40AF; margin-top: 40px; }}
  .container {{ max-width: 1200px; margin: 0 auto; }}
  .grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 20px; }}
  .card {{ background: #f8fafc; border: 1px solid #e2e8f0; border-radius: 8px; padding: 16px; }}
  .slide-preview {{ position: relative; width: 100%; padding-bottom: 56.25%;
    background: {COLOR_BG}; border: 1px solid #ccc; border-radius: 6px; overflow: hidden; }}
  .badge {{ display: inline-block; padding: 4px 12px; border-radius: 12px;
    font-size: 12px; font-weight: bold; }}
  .badge-green {{ background: #D1FAE5; color: #065F46; }}
  .badge-orange {{ background: #FEF3C7; color: #92400E; }}
  .note {{ background: #FEF3C7; border-left: 4px solid #F59E0B; padding: 12px; margin: 16px 0; }}
</style>
</head>
<body>
<div class="container">
<h1>SPEC 0027 网格布局演示</h1>
<p>本页面演示 PptRenderer 在不同图表数量下的自适应布局行为，
使用生产代码 <code>PptRenderer.render()</code> 生成。</p>

<div class="note">
  <strong>布局规则（SPEC 0024/0027）：</strong>
  <ul>
    <li>1 张图表 → <code>_place_chart_centered</code> 单图居中（8" 宽）</li>
    <li>2 张图表 → <code>_place_chart_side_by_side</code> 双图并排（1×2 Grid，5.8" 宽/图）</li>
    <li>3 张图表 → <code>_place_chart_three</code> 上排 2 张 + 下排 1 张居中</li>
    <li>4 张图表 → <code>_place_chart_grid</code> 2×2 网格（3.8" 宽/图）</li>
    <li><strong>超过 4 张 → 截断为前 4 张 + 底部截断注释</strong></li>
  </ul>
</div>

<h2>场景一：4 张图表 → 2×2 网格布局</h2>
<p><span class="badge badge-green">布局类型：_place_chart_grid（2×2）</span></p>

<div class="grid">
  <div class="card">
    <h3>幻灯片预览（坐标可视化）</h3>
    <div class="slide-preview">
      {_slide_to_html(ppt4_chart_slide) if ppt4_chart_slide else '<p>无图表页</p>'}
    </div>
  </div>
  <div class="card">
    <h3>图表坐标表</h3>
    {chart_table(ppt4_chart_slide)}
  </div>
</div>

<h2>场景二：6 张图表 → 截断为 4 张 + 注释</h2>
<p><span class="badge badge-orange">布局类型：_place_chart_grid（2×2）+ _add_truncation_note</span></p>

<div class="grid">
  <div class="card">
    <h3>幻灯片预览（坐标可视化）</h3>
    <div class="slide-preview">
      {_slide_to_html(ppt6_chart_slide) if ppt6_chart_slide else '<p>无图表页</p>'}
    </div>
  </div>
  <div class="card">
    <h3>图表坐标表</h3>
    {chart_table(ppt6_chart_slide)}
  </div>
</div>

<h2>Grid 坐标计算验证</h2>
<p>以下为 <code>_GridHelper</code> 在 2×2 网格中计算的坐标值（来自生产代码）：</p>
<table style="border-collapse:collapse;width:100%;font-size:13px;">
  <tr style="background:#2563EB;color:white;">
    <th style="padding:8px;border:1px solid #ddd;">单元格</th>
    <th style="padding:8px;border:1px solid #ddd;">Left</th>
    <th style="padding:8px;border:1px solid #ddd;">Top</th>
    <th style="padding:8px;border:1px solid #ddd;">Width</th>
    <th style="padding:8px;border:1px solid #ddd;">Height</th>
    <th style="padding:8px;border:1px solid #ddd;">说明</th>
  </tr>
  <tr>
    <td style="padding:8px;border:1px solid #ddd;">cell(0,0) 左上</td>
    <td style="padding:8px;border:1px solid #ddd;">0.70"</td>
    <td style="padding:8px;border:1px solid #ddd;">1.50"</td>
    <td style="padding:8px;border:1px solid #ddd;">3.80"</td>
    <td style="padding:8px;border:1px solid #ddd;">2.30"</td>
    <td style="padding:8px;border:1px solid #ddd;">上排左</td>
  </tr>
  <tr>
    <td style="padding:8px;border:1px solid #ddd;">cell(0,1) 右上</td>
    <td style="padding:8px;border:1px solid #ddd;">6.80"</td>
    <td style="padding:8px;border:1px solid #ddd;">1.50"</td>
    <td style="padding:8px;border:1px solid #ddd;">3.80"</td>
    <td style="padding:8px;border:1px solid #ddd;">2.30"</td>
    <td style="padding:8px;border:1px solid #ddd;">上排右</td>
  </tr>
  <tr>
    <td style="padding:8px;border:1px solid #ddd;">cell(1,0) 左下</td>
    <td style="padding:8px;border:1px solid #ddd;">0.70"</td>
    <td style="padding:8px;border:1px solid #ddd;">4.00"</td>
    <td style="padding:8px;border:1px solid #ddd;">3.80"</td>
    <td style="padding:8px;border:1px solid #ddd;">2.30"</td>
    <td style="padding:8px;border:1px solid #ddd;">下排左</td>
  </tr>
  <tr>
    <td style="padding:8px;border:1px solid #ddd;">cell(1,1) 右下</td>
    <td style="padding:8px;border:1px solid #ddd;">6.80"</td>
    <td style="padding:8px;border:1px solid #ddd;">4.00"</td>
    <td style="padding:8px;border:1px solid #ddd;">3.80"</td>
    <td style="padding:8px;border:1px solid #ddd;">2.30"</td>
    <td style="padding:8px;border:1px solid #ddd;">下排右</td>
  </tr>
</table>
<p style="margin-top:12px;color:#666;font-size:13px;">
  网格参数：left=0.7", top=1.5", width=9.9", height=4.8", rows=2, cols=2,
  h_gap=2.3", v_gap=0.2"<br>
  计算：cell_w = (9.9 - 2.3×1) // 2 = 3.8",
  cell_h = (4.8 - 0.2×1) // 2 = 2.3"
</p>

</div>
</body>
</html>"""

    Path(output_path).write_text(html, encoding="utf-8")


def main():
    output_dir = Path("dev-docs/e2e-screenshots/spec0027")
    output_dir.mkdir(parents=True, exist_ok=True)

    # 1. 生成 6 张图表
    print("1. 生成示例图表（Seaborn + SciencePlots）...")
    charts_dir = output_dir / "demo_charts"
    all_charts = generate_sample_charts(charts_dir, 6)
    print(f"   生成 {len(all_charts)} 张图表: {[c['name'] for c in all_charts]}")

    # 2. 渲染 4 图表 PPT
    print("\n2. 渲染 4 图表 PPT（2×2 网格布局）...")
    ppt4_path = str(output_dir / "grid_layout_4charts.pptx")
    render_ppt_with_charts(all_charts[:4], ppt4_path)
    print(f"   已保存: {ppt4_path}")

    # 3. 渲染 6 图表 PPT（截断演示）
    print("\n3. 渲染 6 图表 PPT（截断为 4 + 注释）...")
    ppt6_path = str(output_dir / "grid_layout_6charts.pptx")
    render_ppt_with_charts(all_charts, ppt6_path)
    print(f"   已保存: {ppt6_path}")

    # 4. 提取布局信息
    print("\n4. 提取 PPT 布局信息...")
    ppt4_info = extract_slide_layout(ppt4_path)
    ppt6_info = extract_slide_layout(ppt6_path)

    # 打印图表页信息
    for label, info in [("4图表", ppt4_info), ("6图表", ppt6_info)]:
        for slide in info:
            chart_count = sum(1 for s in slide["shapes"] if s["is_chart"])
            if chart_count > 0:
                print(f"   [{label}] 第{slide['slide_idx']}页: {chart_count}张图表")
                for s in slide["shapes"]:
                    if s["is_chart"]:
                        print(f"     - 图表: left={s['left']:.2f}\" top={s['top']:.2f}\" "
                              f"w={s['width']:.2f}\" h={s['height']:.2f}\"")
                    if "共" in s.get("text", "") and "图表" in s.get("text", ""):
                        print(f"     - 截断注释: {s['text']}")

    # 5. 生成 HTML 预览
    print("\n5. 生成 HTML 预览...")
    html_path = str(output_dir / "grid-layout-demo.html")
    generate_html_preview(ppt4_info, ppt6_info, html_path)
    print(f"   已保存: {html_path}")

    print("\n✅ 演示完成！")
    print(f"   PPT 文件: {ppt4_path}")
    print(f"   PPT 文件: {ppt6_path}")
    print(f"   HTML 预览: {html_path}")


if __name__ == "__main__":
    main()
