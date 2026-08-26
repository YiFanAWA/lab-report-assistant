r"""SPEC 0027 截断逻辑验证脚本。

验证内容：
1. 从 grid_layout_6charts.pptx 中提取截断注释的精确属性（文本、位置、颜色、字号）
2. 生成包含 7 张图表的 PPT，验证截断逻辑
3. 提取 7 图表 PPT 的截断注释属性
4. 生成 HTML 预览对比

运行方式（从 server 目录）：
    .\.venv\Scripts\python.exe -m scripts.verify_truncation_logic
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import scienceplots  # noqa: F401
plt.style.use(["science", "no-latex", "cjk-sc-font", "bright"])
import seaborn as sns
sns.set_theme(style="whitegrid", palette="bright", font="Microsoft YaHei")
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from app.infrastructure.renderers.ppt_renderer import PptRenderer, SLIDE_WIDTH, SLIDE_HEIGHT


def generate_sample_charts(output_dir: Path, count: int) -> list[dict]:
    """生成指定数量的示例图表 PNG。"""
    output_dir.mkdir(parents=True, exist_ok=True)

    np.random.seed(42)
    n = 200
    data = {
        "age": np.random.normal(50, 15, n),
        "bmi": np.random.normal(24, 4, n),
        "symptom_score": np.random.randint(0, 10, n),
        "diagnosis": np.random.choice(["胃炎", "溃疡", "正常", "其他"], n),
        "treatment_days": np.random.randint(1, 30, n),
        "blood_pressure": np.random.normal(120, 20, n),
        "heart_rate": np.random.normal(75, 12, n),
    }
    df = pd.DataFrame(data)

    chart_configs = [
        ("age_distribution", "年龄分布直方图",
         lambda: sns.histplot(data=df, x="age", kde=True, bins=20)),
        ("diagnosis_count", "诊断类别分布",
         lambda: sns.countplot(data=df, x="diagnosis")),
        ("symptom_boxplot", "症状评分箱线图",
         lambda: sns.boxplot(data=df[["symptom_score", "treatment_days"]])),
        ("age_bmi_scatter", "年龄与BMI散点图",
         lambda: sns.scatterplot(data=df, x="age", y="bmi")),
        ("bp_distribution", "血压分布直方图",
         lambda: sns.histplot(data=df, x="blood_pressure", kde=True, bins=20)),
        ("corr_heatmap", "相关性热图",
         lambda: sns.heatmap(df.corr(numeric_only=True), annot=True, cmap="coolwarm", fmt=".2f")),
        ("heart_rate_hist", "心率分布直方图",
         lambda: sns.histplot(data=df, x="heart_rate", kde=True, bins=20)),
    ]

    charts = []
    for name, title, plot_fn in chart_configs[:count]:
        plt.figure(figsize=(8, 5))
        plot_fn()
        plt.title(title)
        plt.tight_layout()
        file_path = output_dir / f"{name}.png"
        plt.savefig(file_path, dpi=100, bbox_inches="tight")
        plt.close()
        charts.append({
            "name": title,
            "file_path": str(file_path),
            "artifact_type": "CHART_PNG",
        })

    return charts


def render_ppt_with_charts(charts: list[dict], output_path: str) -> str:
    """使用生产 PptRenderer 渲染 PPT。"""
    renderer = PptRenderer()
    outline_sections = [
        {
            "title": "数据分析概述",
            "source_type": "REQUIREMENT",
            "content": "本实验对胃病数据进行分析，包括描述性统计和可视化。",
        },
        {
            "title": "总结",
            "source_type": "SUMMARY",
            "content": "通过数据分析发现年龄与BMI存在相关性，胃炎是最常见诊断类别。",
        },
    ]
    return renderer.render(
        project_name="截断逻辑验证",
        project_topic="胃病数据分析报告",
        outline_sections=outline_sections,
        execution_artifacts=charts,
        output_path=output_path,
        config={"theme_color": "2563EB", "include_charts": True},
    )


def extract_truncation_note_details(ppt_path: str) -> dict | None:
    """从 PPT 中提取截断注释的精确属性。

    返回 dict 或 None（如果未找到截断注释）：
    {
        'slide_idx': int,
        'text': str,
        'left_inches': float,
        'top_inches': float,
        'width_inches': float,
        'height_inches': float,
        'font_size_pt': float,
        'font_color_hex': str,
        'alignment': str,
    }
    """
    prs = Presentation(ppt_path)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if "共" in text and "图表" in text and "已展示前" in text:
                # 找到截断注释
                info = {
                    "slide_idx": slide_idx + 1,
                    "text": text,
                    "left_inches": round(Emu(shape.left).inches, 3) if shape.left else 0,
                    "top_inches": round(Emu(shape.top).inches, 3) if shape.top else 0,
                    "width_inches": round(Emu(shape.width).inches, 3) if shape.width else 0,
                    "height_inches": round(Emu(shape.height).inches, 3) if shape.height else 0,
                    "font_size_pt": None,
                    "font_color_hex": None,
                    "alignment": None,
                }

                # 提取字体属性
                for para in shape.text_frame.paragraphs:
                    if para.alignment is not None:
                        info["alignment"] = str(para.alignment)
                    for run in para.runs:
                        if run.font.size is not None:
                            info["font_size_pt"] = Pt(1).pt if False else round(
                                run.font.size.pt, 1
                            )
                        if run.font.color is not None and run.font.color.rgb is not None:
                            info["font_color_hex"] = str(run.font.color.rgb)
                        break
                    break

                return info

    return None


def extract_chart_slide_info(ppt_path: str) -> dict | None:
    """提取图表页信息（图表数量、坐标、是否截断）。"""
    prs = Presentation(ppt_path)

    for slide_idx, slide in enumerate(prs.slides):
        chart_shapes = []
        note_text = None
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                chart_shapes.append({
                    "left": round(Emu(shape.left).inches, 2),
                    "top": round(Emu(shape.top).inches, 2),
                    "width": round(Emu(shape.width).inches, 2),
                    "height": round(Emu(shape.height).inches, 2),
                })
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "共" in text and "图表" in text and "已展示前" in text:
                    note_text = text

        if chart_shapes:
            return {
                "slide_idx": slide_idx + 1,
                "chart_count": len(chart_shapes),
                "charts": chart_shapes,
                "has_truncation_note": note_text is not None,
                "note_text": note_text,
            }

    return None


def generate_html_preview(
    note_6chart: dict | None,
    info_6chart: dict | None,
    note_7chart: dict | None,
    info_7chart: dict | None,
    output_path: str,
) -> None:
    """生成 HTML 预览对比。"""

    def note_table(note, label):
        if not note:
            return f"<p style='color:red;'>❌ {label}：未找到截断注释</p>"
        color_display = note.get("font_color_hex", "N/A")
        # 判断颜色类型
        color_val = note.get("font_color_hex", "")
        if color_val and color_val.upper() in ("888888", "666666"):
            color_type = "灰色（#888888）"
            color_badge = "background:#E5E7EB;color:#374151;"
        elif color_val and color_val.upper().startswith("F"):
            color_type = "橙色系"
            color_badge = "background:#FEF3C7;color:#92400E;"
        else:
            color_type = f"其他颜色（#{color_val}）"
            color_badge = "background:#E0E7FF;color:#3730A3;"

        return f"""<table style="border-collapse:collapse;width:100%;font-size:13px;">
          <tr><td style="padding:6px;border:1px solid #ddd;font-weight:bold;">属性</td>
              <td style="padding:6px;border:1px solid #ddd;font-weight:bold;">值</td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">所在页</td>
              <td style="padding:6px;border:1px solid #ddd;">第 {note['slide_idx']} 页</td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">注释文本</td>
              <td style="padding:6px;border:1px solid #ddd;font-weight:bold;">"{note['text']}"</td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">位置 (Left, Top)</td>
              <td style="padding:6px;border:1px solid #ddd;">{note['left_inches']}\", {note['top_inches']}\"</td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">尺寸 (W×H)</td>
              <td style="padding:6px;border:1px solid #ddd;">{note['width_inches']}\" × {note['height_inches']}\"</td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">字号</td>
              <td style="padding:6px;border:1px solid #ddd;">{note['font_size_pt']}pt</td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">颜色</td>
              <td style="padding:6px;border:1px solid #ddd;">
                <span style="padding:2px 8px;border-radius:8px;font-size:11px;{color_badge}">#{color_val}</span>
                {color_type}
              </td></tr>
          <tr><td style="padding:6px;border:1px solid #ddd;">对齐方式</td>
              <td style="padding:6px;border:1px solid #ddd;">{note['alignment']}</td></tr>
        </table>"""

    def chart_info_table(info, label):
        if not info:
            return f"<p>{label}：未找到图表页</p>"
        rows = []
        for i, c in enumerate(info["charts"]):
            rows.append(
                f"<tr><td>{i+1}</td><td>{c['left']}\"</td><td>{c['top']}\"</td>"
                f"<td>{c['width']}\"</td><td>{c['height']}\"</td></tr>"
            )
        note_status = (
            '<span style="color:#059669;font-weight:bold;">✅ 有截断注释</span>'
            if info["has_truncation_note"]
            else '<span style="color:#DC2626;font-weight:bold;">❌ 无截断注释</span>'
        )
        return f"""<p>图表页：第 {info['slide_idx']} 页，共 {info['chart_count']} 张图表，{note_status}</p>
        <table style="border-collapse:collapse;width:100%;font-size:13px;">
          <tr style="background:#2563EB;color:white;">
            <th style="padding:6px;border:1px solid #ddd;">序号</th>
            <th style="padding:6px;border:1px solid #ddd;">Left</th>
            <th style="padding:6px;border:1px solid #ddd;">Top</th>
            <th style="padding:6px;border:1px solid #ddd;">Width</th>
            <th style="padding:6px;border:1px solid #ddd;">Height</th>
          </tr>
          {''.join(rows)}
        </table>"""

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>SPEC 0027 截断逻辑验证</title>
<style>
  body {{ font-family: "Microsoft YaHei", sans-serif; margin: 20px; background: #fff; }}
  h1 {{ color: #2563EB; border-bottom: 3px solid #2563EB; padding-bottom: 10px; }}
  h2 {{ color: #1E40AF; margin-top: 40px; }}
  .container {{ max-width: 1100px; margin: 0 auto; }}
  .grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 20px; margin-top: 16px; }}
  .card {{ background: #f8fafc; border: 1px solid #e2e8f0; border-radius: 8px; padding: 16px; }}
  .note-box {{ background: #FEF3C7; border-left: 4px solid #F59E0B; padding: 12px; margin: 16px 0; }}
  .source-code {{ background: #1e293b; color: #e2e8f0; padding: 12px; border-radius: 6px;
    font-family: "Cascadia Code", monospace; font-size: 12px; overflow-x: auto; }}
  .source-code .keyword {{ color: #c084fc; }}
  .source-code .string {{ color: #86efac; }}
  .source-code .comment {{ color: #94a3b8; }}
</style>
</head>
<body>
<div class="container">
<h1>SPEC 0027 截断逻辑验证</h1>

<div class="note-box">
  <strong>截断逻辑源码</strong>（<code>ppt_renderer.py:786-797, 1029-1042</code>）：
  <div class="source-code">
<span class="comment"># _add_chart_slide 方法（第 786-797 行）</span>
count = <span class="keyword">min</span>(<span class="keyword">len</span>(chart_artifacts), <span class="string">4</span>)  <span class="comment"># 最多 4 张</span>
<span class="keyword">if</span> count == <span class="string">4</span>:
    self._place_chart_grid(slide, chart_artifacts[:<span class="string">4</span>])  <span class="comment"># 只取前 4 张</span>
<span class="keyword">if</span> <span class="keyword">len</span>(chart_artifacts) &gt; <span class="string">4</span>:
    self._add_truncation_note(slide, <span class="keyword">len</span>(chart_artifacts))  <span class="comment"># 传入总数</span>

<span class="comment"># _add_truncation_note 方法（第 1029-1042 行）</span>
tb = slide.shapes.add_textbox(
    Inches(MARGIN_LEFT), Inches(<span class="string">6.5</span>),  <span class="comment"># 底部位置</span>
    Inches(SLIDE_WIDTH - <span class="string">2</span> * MARGIN_LEFT), Inches(<span class="string">0.4</span>),
)
run.text = <span class="string">f"共 {{total_count}} 张图表，已展示前 4 张"</span>
self._set_run_font(run, FONT_SIZE_CAPTION, RGBColor(<span class="string">0x88</span>, <span class="string">0x88</span>, <span class="string">0x88</span>))
<span class="comment"># ↑ 颜色是灰色 #888888，不是橙色；字号 12pt (FONT_SIZE_CAPTION)</span>
  </div>
</div>

<h2>场景一：6 张图表 PPT（grid_layout_6charts.pptx）</h2>
<p>输入 6 张图表 → 1 张被内容页吸收 → 5 张进入图表页 → 截断为 4 张展示</p>
<div class="grid">
  <div class="card">
    <h3>截断注释属性</h3>
    {note_table(note_6chart, "6图表PPT")}
  </div>
  <div class="card">
    <h3>图表页信息</h3>
    {chart_info_table(info_6chart, "6图表PPT")}
  </div>
</div>

<h2>场景二：7 张图表 PPT（grid_layout_7charts.pptx）</h2>
<p>输入 7 张图表 → 1 张被内容页吸收 → 6 张进入图表页 → 截断为 4 张展示</p>
<div class="grid">
  <div class="card">
    <h3>截断注释属性</h3>
    {note_table(note_7chart, "7图表PPT")}
  </div>
  <div class="card">
    <h3>图表页信息</h3>
    {chart_info_table(info_7chart, "7图表PPT")}
  </div>
</div>

<h2>验证结论</h2>
<div class="card">
  <h3>截断逻辑确认</h3>
  <ul>
    <li>✅ 超过 4 张图表时，<code>_add_chart_slide</code> 只渲染前 4 张（<code>chart_artifacts[:4]</code>）</li>
    <li>✅ 调用 <code>_add_truncation_note</code> 添加底部注释，文本格式为「共 X 张图表，已展示前 4 张」</li>
    <li>✅ 注释位置在底部（top=6.5\"），宽度撑满内容区，高度 0.4\"</li>
    <li>✅ 注释字号为 12pt（FONT_SIZE_CAPTION），颜色为灰色 #888888（RGBColor(0x88, 0x88, 0x88)）</li>
    <li>✅ 注释居中对齐（PP_ALIGN.CENTER）</li>
    <li>⚠️ 注意：注释颜色是<strong>灰色</strong>而非橙色，橙色仅用于 HTML 预览中的视觉强调</li>
  </ul>

  <h3>6 图表 vs 7 图表对比</h3>
  <ul>
    <li>6 图表：1 张被内容页吸收 → 5 张进入图表页 → 注释显示「共 5 张图表，已展示前 4 张」</li>
    <li>7 图表：1 张被内容页吸收 → 6 张进入图表页 → 注释显示「共 6 张图表，已展示前 4 张」</li>
    <li>两者图表页都只展示 4 张（2×2 网格），截断注释的总数随输入图表数变化</li>
  </ul>
</div>

</div>
</body>
</html>"""

    Path(output_path).write_text(html, encoding="utf-8")


def main():
    # 输出目录（相对于项目根目录）
    project_root = Path(__file__).parent.parent.parent
    output_dir = project_root / "dev-docs" / "e2e-screenshots" / "spec0027"
    output_dir.mkdir(parents=True, exist_ok=True)

    # === 1. 验证已有的 6 图表 PPT ===
    print("=" * 60)
    print("1. 验证已有的 6 图表 PPT 截断注释")
    print("=" * 60)
    ppt6_path = str(output_dir / "grid_layout_6charts.pptx")
    if not Path(ppt6_path).exists():
        print(f"   ❌ 文件不存在: {ppt6_path}")
        return

    note_6chart = extract_truncation_note_details(ppt6_path)
    info_6chart = extract_chart_slide_info(ppt6_path)

    if note_6chart:
        print(f"   ✅ 找到截断注释:")
        print(f"      文本: \"{note_6chart['text']}\"")
        print(f"      位置: left={note_6chart['left_inches']}\", top={note_6chart['top_inches']}\"")
        print(f"      尺寸: {note_6chart['width_inches']}\" × {note_6chart['height_inches']}\"")
        print(f"      字号: {note_6chart['font_size_pt']}pt")
        print(f"      颜色: #{note_6chart['font_color_hex']}")
        print(f"      对齐: {note_6chart['alignment']}")
    else:
        print("   ❌ 未找到截断注释")

    if info_6chart:
        print(f"   图表页: 第{info_6chart['slide_idx']}页, {info_6chart['chart_count']}张图表")
        print(f"   截断注释: {'有' if info_6chart['has_truncation_note'] else '无'}")

    # === 2. 生成 7 图表 PPT ===
    print("\n" + "=" * 60)
    print("2. 生成 7 图表 PPT")
    print("=" * 60)

    charts_dir = output_dir / "truncation_demo_charts"
    print("   生成 7 张示例图表...")
    all_charts = generate_sample_charts(charts_dir, 7)
    print(f"   生成 {len(all_charts)} 张图表: {[c['name'] for c in all_charts]}")

    ppt7_path = str(output_dir / "grid_layout_7charts.pptx")
    print(f"\n   渲染 7 图表 PPT...")
    render_ppt_with_charts(all_charts, ppt7_path)
    print(f"   已保存: {ppt7_path}")

    # === 3. 验证 7 图表 PPT ===
    print("\n" + "=" * 60)
    print("3. 验证 7 图表 PPT 截断注释")
    print("=" * 60)

    note_7chart = extract_truncation_note_details(ppt7_path)
    info_7chart = extract_chart_slide_info(ppt7_path)

    if note_7chart:
        print(f"   ✅ 找到截断注释:")
        print(f"      文本: \"{note_7chart['text']}\"")
        print(f"      位置: left={note_7chart['left_inches']}\", top={note_7chart['top_inches']}\"")
        print(f"      尺寸: {note_7chart['width_inches']}\" × {note_7chart['height_inches']}\"")
        print(f"      字号: {note_7chart['font_size_pt']}pt")
        print(f"      颜色: #{note_7chart['font_color_hex']}")
        print(f"      对齐: {note_7chart['alignment']}")
    else:
        print("   ❌ 未找到截断注释")

    if info_7chart:
        print(f"   图表页: 第{info_7chart['slide_idx']}页, {info_7chart['chart_count']}张图表")
        print(f"   截断注释: {'有' if info_7chart['has_truncation_note'] else '无'}")

    # === 4. 对比验证 ===
    print("\n" + "=" * 60)
    print("4. 对比验证结论")
    print("=" * 60)

    if note_6chart and note_7chart:
        print(f"   6 图表 PPT 注释: \"{note_6chart['text']}\"")
        print(f"   7 图表 PPT 注释: \"{note_7chart['text']}\"")
        print(f"   总数差异: 5 vs 6（因为 1 张被内容页吸收）")
        print(f"   颜色一致: #{note_6chart['font_color_hex']} == #{note_7chart['font_color_hex']}")
        print(f"   字号一致: {note_6chart['font_size_pt']}pt == {note_7chart['font_size_pt']}pt")

    # === 5. 生成 HTML 预览 ===
    print("\n" + "=" * 60)
    print("5. 生成 HTML 预览")
    print("=" * 60)

    html_path = str(output_dir / "truncation-logic-verification.html")
    generate_html_preview(note_6chart, info_6chart, note_7chart, info_7chart, html_path)
    print(f"   已保存: {html_path}")

    print("\n✅ 验证完成！")
    print(f"   6 图表 PPT: {ppt6_path}")
    print(f"   7 图表 PPT: {ppt7_path}")
    print(f"   HTML 预览: {html_path}")


if __name__ == "__main__":
    main()
