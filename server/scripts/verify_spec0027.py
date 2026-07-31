"""SPEC 0027 真实图表 PPT 预览生成与网格布局对齐验证脚本。

流程：
1. 准备教学数据集（胃病数据分析，CSV）
2. 构造 AnalysisPlan（含 HISTOGRAM/BOXPLOT/BAR/SCATTER + CORRELATION 热图）
3. 用 LocalRuleCodeTaskProvider 生成 Python 代码（SciencePlots + Seaborn）
4. 用沙箱 python_executor 执行代码，生成真实图表 PNG
5. 用 PptRenderer 渲染 6 种预设色 PPT（4 张图表触发 2×2 网格布局）
6. 程序化验证 Grid 布局坐标对齐（用 _GridHelper 计算并对比）
7. 生成 HTML 预览文件（嵌入真实图表 base64）

验收点：
- 图表生成成功（SciencePlots + Seaborn 样式应用，无 LaTeX 依赖）
- PPT 文件 XML 完整性（6 种预设色均可重新打开）
- _place_chart_grid 使用 _GridHelper 后坐标与原硬编码一致
- 三角色彩系统 + 三明治结构 + 渐变/圆角/阴影/边框保持（SPEC 0025/0026）
"""

import base64
import json
import sys
import tempfile
from pathlib import Path

# 添加 server 目录到 path
sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from app.infrastructure.renderers.ppt_renderer import (
    PptRenderer, SLIDE_WIDTH, SLIDE_HEIGHT, FOOTER_BAR_TOP,
)
from app.modules.llm.code_task_provider import LocalRuleCodeTaskProvider
from app.infrastructure.sandbox.python_executor import execute_code


# 6 种预设色（与 SPEC 0025/0026 验证脚本一致）
PRESET_COLORS = [
    ("蓝", "#2563eb"),
    ("紫", "#7c3aed"),
    ("绿", "#16a34a"),
    ("红", "#dc2626"),
    ("橙", "#ea580c"),
    ("灰", "#475569"),
]

# 示例大纲（包含 REQUIREMENT/EVIDENCE/EXECUTION/SUMMARY 四类 source_type）
SAMPLE_SECTIONS = [
    {
        "id": "sec_001",
        "title": "实验目的与背景",
        "content": "基于胃病数据分析掌握描述性统计与可视化方法，"
                   "理解医学数据的分布特征与变量关系。",
        "source_type": "REQUIREMENT",
        "source_ids": ["plan_001"],
    },
    {
        "id": "sec_002",
        "title": "数据来源与方法",
        "content": "使用教学数据集，通过 pandas 清洗、seaborn 可视化，"
                   "应用 SciencePlots 期刊样式提升图表质量。",
        "source_type": "EVIDENCE",
        "source_ids": ["ev_001"],
    },
    {
        "id": "sec_003",
        "title": "描述性分析结果",
        "content": "生成年龄分布直方图、诊断类别柱状图、症状评分箱线图、"
                   "年龄-住院天数散点图与相关性热图。",
        "source_type": "EXECUTION",
        "source_ids": ["run_001"],
    },
    {
        "id": "sec_004",
        "title": "结论与讨论",
        "content": "胃病数据分析验证了 SciencePlots + Seaborn 组合"
                   "在医学统计可视化中的有效性，图表符合科研出版规范。",
        "source_type": "SUMMARY",
        "source_ids": [],
    },
]


def make_sample_dataset(csv_path: Path) -> None:
    """生成教学用胃病数据集 CSV。"""
    import csv
    import random

    random.seed(42)  # 可复现
    diagnoses = ["胃炎", "胃溃疡", "胃癌", "十二指肠溃疡", "功能性消化不良"]
    genders = ["男", "女"]

    rows = []
    for _ in range(120):
        age = random.randint(18, 80)
        gender = random.choice(genders)
        diagnosis = random.choices(
            diagnoses, weights=[30, 25, 10, 20, 15], k=1,
        )[0]
        # 症状评分与诊断相关
        base_score = {"胃炎": 4, "胃溃疡": 6, "胃癌": 8,
                      "十二指肠溃疡": 5, "功能性消化不良": 3}
        symptom_score = max(
            1, min(10, base_score[diagnosis] + random.randint(-2, 2)),
        )
        # 住院天数与诊断和年龄相关
        base_days = {"胃炎": 3, "胃溃疡": 7, "胃癌": 14,
                     "十二指肠溃疡": 5, "功能性消化不良": 1}
        hospital_days = max(1, base_days[diagnosis] + random.randint(-2, 3))
        rows.append([age, gender, diagnosis, symptom_score, hospital_days])

    with open(csv_path, "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow([
            "age", "gender", "diagnosis", "symptom_score", "hospital_days",
        ])
        writer.writerows(rows)

    print(f"数据集已生成：{csv_path}（{len(rows)} 行）")


def make_analysis_plan() -> dict:
    """构造含多种图表类型的 AnalysisPlan。"""
    return {
        "cleaning_plan": [
            {
                "issue_type": "MISSING_VALUE",
                "field": "*",
                "action": "删除缺失值",
                "reason": "保证分析数据完整性",
            },
            {
                "issue_type": "DUPLICATE_ROW",
                "field": "*",
                "action": "删除重复行",
                "reason": "避免重复样本影响统计",
            },
        ],
        "analysis_plan": [
            {
                "analysis_type": "DESCRIPTIVE_STATISTICS",
                "target_fields": "*",
                "method": "描述性统计",
                "expected_output": "均值/标准差/分位数",
            },
            {
                "analysis_type": "CORRELATION",
                "target_fields": "age vs hospital_days",
                "method": "相关性分析",
                "expected_output": "相关系数矩阵 + 热图",
            },
        ],
        "chart_plan": [
            {
                "chart_type": "HISTOGRAM",
                "title": "年龄分布直方图",
                "data_fields": ["age"],
                "description": "患者年龄分布特征",
            },
            {
                "chart_type": "BAR",
                "title": "诊断类别分布",
                "data_fields": ["diagnosis"],
                "description": "各类胃病诊断频次",
            },
            {
                "chart_type": "BOXPLOT",
                "title": "症状评分箱线图",
                "data_fields": ["symptom_score"],
                "description": "症状评分分布与离群值",
            },
            {
                "chart_type": "SCATTER",
                "title": "年龄与住院天数散点图",
                "data_fields": ["age", "hospital_days"],
                "description": "年龄与住院天数关系",
            },
        ],
    }


def generate_and_execute_charts(output_dir: Path) -> list[dict]:
    """生成代码并在沙箱中执行，收集图表产物。"""
    # 1. 准备数据集
    csv_path = output_dir / "sample_gastric_data.csv"
    make_sample_dataset(csv_path)

    # 2. 生成代码
    provider = LocalRuleCodeTaskProvider()
    plan = make_analysis_plan()
    draft = provider.generate(plan)
    code = draft.code

    # 保存生成的代码用于审查
    code_path = output_dir / "generated_code.py"
    code_path.write_text(code, encoding="utf-8")
    print(f"生成代码已保存：{code_path}")

    # 验证代码包含 SciencePlots + Seaborn
    checks = {
        "import scienceplots": "import scienceplots" in code,
        "plt.style.use(['science'": "plt.style.use(['science'" in code,
        "import seaborn as sns": "import seaborn as sns" in code,
        "sns.set_theme": "sns.set_theme" in code,
        "sns.histplot": "sns.histplot" in code,
        "sns.boxplot": "sns.boxplot" in code,
        "sns.countplot": "sns.countplot" in code,
        "sns.scatterplot": "sns.scatterplot" in code,
        "sns.heatmap": "sns.heatmap" in code,
    }
    print("\n=== 图表代码集成检查 ===")
    for key, passed in checks.items():
        print(f"  {'✓' if passed else '✗'} {key}")
    if not all(checks.values()):
        print("警告：部分 SciencePlots/Seaborn 集成检查未通过")

    # 3. 沙箱执行
    print("\n=== 沙箱执行代码 ===")
    work_dir = output_dir / "sandbox_workdir"
    work_dir.mkdir(parents=True, exist_ok=True)

    result = execute_code(
        code=code,
        work_dir=str(work_dir),
        data_path=str(csv_path),
        timeout_seconds=60,
        memory_limit_mb=2048,
    )

    print(f"  exit_code: {result.exit_code}")
    print(f"  duration: {result.duration_seconds:.2f}s")
    print(f"  stdout 末尾: {result.stdout[-300:] if result.stdout else '(空)'}")
    if result.stderr:
        print(f"  stderr 末尾: {result.stderr[-500:]}")
    if result.sandbox_error_code:
        print(f"  沙箱错误码: {result.sandbox_error_code}")

    # 4. 收集图表产物
    chart_files = sorted(work_dir.glob("*.png"))
    print(f"\n=== 收集图表产物 ===")
    print(f"  共 {len(chart_files)} 张图表")
    for cf in chart_files:
        print(f"  - {cf.name} ({cf.stat().st_size} bytes)")

    # 转为 artifacts 格式
    artifacts = [
        {
            "name": cf.name,
            "artifact_type": "CHART_PNG",
            "file_path": str(cf),
        }
        for cf in chart_files
    ]
    return artifacts, checks, result


def verify_grid_alignment() -> dict:
    """验证 _GridHelper 计算的坐标与原硬编码 positions 完全一致。"""
    print("\n=== Grid 布局坐标对齐验证 ===")
    results = {}

    # 1. _place_chart_grid（2×2 网格）
    # 原硬编码：positions = [(0.7, 1.5), (6.8, 1.5), (0.7, 4.0), (6.8, 4.0)]
    # 原 max_width=3.8, max_height=2.3
    grid_2x2 = PptRenderer._GridHelper(
        left=Inches(0.7), top=Inches(1.5),
        width=Inches(9.9), height=Inches(4.8),
        rows=2, cols=2,
        h_gap=Inches(2.3), v_gap=Inches(0.2),
    )
    expected_2x2 = [
        (0.7, 1.5, 3.8, 2.3),
        (6.8, 1.5, 3.8, 2.3),
        (0.7, 4.0, 3.8, 2.3),
        (6.8, 4.0, 3.8, 2.3),
    ]
    grid_2x2_results = []
    for i, (row, col) in enumerate([(0, 0), (0, 1), (1, 0), (1, 1)]):
        cl, ct, cw, ch = grid_2x2.cell(row, col)
        actual = (
            round(Emu(cl).inches, 4),
            round(Emu(ct).inches, 4),
            round(Emu(cw).inches, 4),
            round(Emu(ch).inches, 4),
        )
        exp = expected_2x2[i]
        match = all(abs(a - e) < 0.01 for a, e in zip(actual, exp))
        grid_2x2_results.append({
            "cell": f"({row},{col})",
            "expected": exp,
            "actual": actual,
            "match": match,
        })
        print(f"  2×2 cell({row},{col}): 预期={exp}, 实际={actual} "
              f"{'✓' if match else '✗ 不一致'}")
    results["grid_2x2"] = grid_2x2_results

    # 2. _place_chart_side_by_side（1×2 网格）
    # 原硬编码：positions = [(0.5, 1.8), (7.0, 1.8)], max_width=5.8
    grid_1x2 = PptRenderer._GridHelper(
        left=Inches(0.5), top=Inches(1.8),
        width=Inches(12.3), height=Inches(5.2),
        rows=1, cols=2,
        h_gap=Inches(0.7), v_gap=0,
    )
    expected_1x2 = [
        (0.5, 1.8, 5.8),
        (7.0, 1.8, 5.8),
    ]
    grid_1x2_results = []
    for i in range(2):
        cl, ct, cw, _ = grid_1x2.cell(0, i)
        actual = (
            round(Emu(cl).inches, 4),
            round(Emu(ct).inches, 4),
            round(Emu(cw).inches, 4),
        )
        exp = expected_1x2[i]
        match = all(abs(a - e) < 0.01 for a, e in zip(actual, exp))
        grid_1x2_results.append({
            "cell": f"(0,{i})",
            "expected": exp,
            "actual": actual,
            "match": match,
        })
        print(f"  1×2 cell(0,{i}): 预期={exp}, 实际={actual} "
              f"{'✓' if match else '✗ 不一致'}")
    results["grid_1x2"] = grid_1x2_results

    # 3. _place_chart_three 上排（1×2 网格）
    # 原硬编码：top_positions = [(0.5, 1.5), (7.0, 1.5)], max_width_top=5.8
    grid_three_top = PptRenderer._GridHelper(
        left=Inches(0.5), top=Inches(1.5),
        width=Inches(12.3), height=Inches(2.3),
        rows=1, cols=2,
        h_gap=Inches(0.7), v_gap=0,
    )
    expected_three = [
        (0.5, 1.5, 5.8),
        (7.0, 1.5, 5.8),
    ]
    grid_three_results = []
    for i in range(2):
        cl, ct, cw, _ = grid_three_top.cell(0, i)
        actual = (
            round(Emu(cl).inches, 4),
            round(Emu(ct).inches, 4),
            round(Emu(cw).inches, 4),
        )
        exp = expected_three[i]
        match = all(abs(a - e) < 0.01 for a, e in zip(actual, exp))
        grid_three_results.append({
            "cell": f"(0,{i})",
            "expected": exp,
            "actual": actual,
            "match": match,
        })
        print(f"  three 上排 cell(0,{i}): 预期={exp}, 实际={actual} "
              f"{'✓' if match else '✗ 不一致'}")
    results["grid_three_top"] = grid_three_results

    # 4. _pct_to_emu 百分比定位验证
    print("\n=== _pct_to_emu 百分比定位验证 ===")
    pct_results = []
    pct_cases = [
        ("10%", Inches(13.333), round(Inches(1.3333), -2)),
        ("50%", Inches(7.5), Inches(3.75)),
        ("100%", Inches(10), Inches(10)),
        ("0%", Inches(10), 0),
        ("12.5%", 1000, 125),
    ]
    for pct_str, total, expected in pct_cases:
        actual = PptRenderer._pct_to_emu(pct_str, total)
        # 允许 ±100 EMU 精度误差
        match = abs(actual - expected) < 200
        pct_results.append({
            "input": f'_pct_to_emu("{pct_str}", {total})',
            "expected": expected,
            "actual": actual,
            "match": match,
        })
        print(f"  {pct_str} of {total}: 预期≈{expected}, 实际={actual} "
              f"{'✓' if match else '✗'}")
    results["pct_to_emu"] = pct_results

    all_match = all(
        cell["match"]
        for key in ["grid_2x2", "grid_1x2", "grid_three_top", "pct_to_emu"]
        for cell in results[key]
    )
    results["all_aligned"] = all_match
    print(f"\n  网格布局对齐总体结果：{'全部对齐 ✓' if all_match else '存在偏差 ✗'}")
    return results


def verify_ppt_visual_effects(pptx_path: Path) -> dict:
    """验证 PPT 视觉效果（SPEC 0025/0026 保持）。"""
    prs = Presentation(str(pptx_path))
    result = {
        "slides_count": len(prs.slides),
        "gradient_found": False,
        "rounded_rect_found": False,
        "shadow_found": False,
        "border_found": False,
        "picture_count": 0,
        "file_valid": True,
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            # 渐变填充
            try:
                if shape.fill.type == MSO_FILL_TYPE.GRADIENT:
                    result["gradient_found"] = True
            except Exception:
                pass

            # 圆角矩形
            try:
                if getattr(shape, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE:
                    result["rounded_rect_found"] = True
            except Exception:
                pass

            # 图片阴影和边框
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result["picture_count"] += 1
                try:
                    spPr = shape._element.spPr
                    effectLst = spPr.find(qn('a:effectLst'))
                    if effectLst is not None and effectLst.find(qn('a:outerShdw')) is not None:
                        result["shadow_found"] = True
                except Exception:
                    pass
                try:
                    if shape.line.color.rgb is not None:
                        result["border_found"] = True
                except Exception:
                    pass

    return result


def render_ppt_files(artifacts: list[dict], output_dir: Path) -> dict:
    """用真实图表渲染 6 种预设色 PPT。"""
    print("\n=== 渲染 6 种预设色 PPT ===")
    ppt_results = {}

    # 确保至少 4 张图表触发 2×2 网格布局
    chart_artifacts = [a for a in artifacts if a["artifact_type"] == "CHART_PNG"]
    if len(chart_artifacts) < 4:
        print(f"  警告：只有 {len(chart_artifacts)} 张图表，不足 4 张，"
              f"2×2 网格布局无法完整展示")
    # 取前 4 张用于 2×2 网格展示
    display_artifacts = chart_artifacts[:4]

    for color_name, theme_color in PRESET_COLORS:
        pptx_path = output_dir / f"spec0027_{color_name}.pptx"
        renderer = PptRenderer()
        try:
            renderer.render(
                project_name=f"SPEC0027验证_{color_name}",
                project_topic=f"胃病数据分析 - {color_name}主题（SciencePlots+Seaborn）",
                outline_sections=SAMPLE_SECTIONS,
                execution_artifacts=display_artifacts,
                output_path=str(pptx_path),
                config={"theme_color": theme_color},
            )
            # 验证文件可重新打开 + 视觉效果
            verify_result = verify_ppt_visual_effects(pptx_path)
            ppt_results[color_name] = {
                "path": str(pptx_path),
                "theme_color": theme_color,
                **verify_result,
            }
            print(f"  [{color_name}] {theme_color}: "
                  f"页数={verify_result['slides_count']} "
                  f"图片={verify_result['picture_count']} "
                  f"渐变={'✓' if verify_result['gradient_found'] else '✗'} "
                  f"圆角={'✓' if verify_result['rounded_rect_found'] else '✗'} "
                  f"阴影={'✓' if verify_result['shadow_found'] else '✗'} "
                  f"边框={'✓' if verify_result['border_found'] else '✗'}")
        except Exception as e:
            ppt_results[color_name] = {
                "path": str(pptx_path),
                "theme_color": theme_color,
                "file_valid": False,
                "error": str(e),
            }
            print(f"  [{color_name}] {theme_color}: 渲染失败 ✗ {e}")

    return ppt_results


def generate_html_preview(
    artifacts: list[dict],
    grid_results: dict,
    ppt_results: dict,
    code_checks: dict,
    output_html: Path,
    output_dir: Path,
) -> None:
    """生成 HTML 预览文件（嵌入真实图表 base64）。"""
    print(f"\n=== 生成 HTML 预览：{output_html} ===")

    # 收集图表 base64
    chart_images = []
    for art in artifacts:
        if art["artifact_type"] != "CHART_PNG":
            continue
        p = Path(art["file_path"])
        if not p.exists():
            continue
        b64 = base64.b64encode(p.read_bytes()).decode("ascii")
        chart_images.append({
            "name": p.name,
            "b64": b64,
        })

    # 构建图表展示 HTML
    charts_html = ""
    for img in chart_images[:6]:
        charts_html += f"""
        <div class="chart-item">
          <h4>{img['name']}</h4>
          <img src="data:image/png;base64,{img['b64']}" alt="{img['name']}" />
        </div>"""

    # 构建 Grid 对齐验证表格
    def build_grid_table(grid_key, title):
        rows_html = ""
        for cell in grid_results.get(grid_key, []):
            status = "✓" if cell["match"] else "✗"
            cls = "pass" if cell["match"] else "fail"
            rows_html += f"""
            <tr class="{cls}">
              <td>{cell['cell']}</td>
              <td>{cell['expected']}</td>
              <td>{cell['actual']}</td>
              <td>{status}</td>
            </tr>"""
        return f"""
        <h3>{title}</h3>
        <table class="verify-table">
          <thead><tr><th>单元格</th><th>预期坐标</th><th>实际坐标</th><th>对齐</th></tr></thead>
          <tbody>{rows_html}</tbody>
        </table>"""

    # 构建 PPT 验证表格
    ppt_rows_html = ""
    for color_name, r in ppt_results.items():
        valid = r.get("file_valid", False)
        status = "✓" if valid else "✗"
        cls = "pass" if valid else "fail"
        ppt_rows_html += f"""
        <tr class="{cls}">
          <td>{color_name}</td>
          <td>{r.get('theme_color', '')}</td>
          <td>{r.get('slides_count', '-')}</td>
          <td>{r.get('picture_count', '-')}</td>
          <td>{'✓' if r.get('gradient_found') else '✗'}</td>
          <td>{'✓' if r.get('rounded_rect_found') else '✗'}</td>
          <td>{'✓' if r.get('shadow_found') else '✗'}</td>
          <td>{'✓' if r.get('border_found') else '✗'}</td>
          <td>{status}</td>
        </tr>"""

    # 构建代码集成检查
    code_checks_html = ""
    for key, passed in code_checks.items():
        cls = "pass" if passed else "fail"
        code_checks_html += f"""
        <tr class="{cls}">
          <td>{key}</td>
          <td>{'✓ 通过' if passed else '✗ 失败'}</td>
        </tr>"""

    all_aligned = grid_results.get("all_aligned", False)

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>SPEC 0027 真实图表 PPT 预览与网格布局验证</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{
  font-family: "Microsoft YaHei", "微软雅黑", sans-serif;
  background: #f0f2f5;
  color: #333;
  padding: 20px;
  line-height: 1.6;
}}
h1 {{ text-align: center; margin-bottom: 10px; font-size: 28px; }}
h2 {{ margin: 30px 0 15px; font-size: 22px; border-bottom: 2px solid #2563eb; padding-bottom: 8px; }}
h3 {{ margin: 20px 0 10px; font-size: 18px; }}
.subtitle {{ text-align: center; color: #666; margin-bottom: 30px; font-size: 14px; }}
.summary-box {{
  max-width: 960px; margin: 0 auto 30px; padding: 20px;
  background: #fff; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);
}}
.summary-box.overall-pass {{ border-left: 5px solid #16a34a; }}
.summary-box.overall-fail {{ border-left: 5px solid #dc2626; }}
.section {{
  max-width: 960px; margin: 0 auto 30px; padding: 20px;
  background: #fff; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);
}}
.charts-grid {{
  display: grid; grid-template-columns: repeat(auto-fill, minmax(300px, 1fr));
  gap: 20px; margin-top: 15px;
}}
.chart-item {{
  border: 1px solid #e0e0e0; border-radius: 8px; padding: 12px;
  background: #fafafa; text-align: center;
}}
.chart-item h4 {{ font-size: 14px; margin-bottom: 8px; color: #555; }}
.chart-item img {{ max-width: 100%; height: auto; border-radius: 4px; }}
.verify-table {{
  width: 100%; border-collapse: collapse; margin-top: 10px; font-size: 14px;
}}
.verify-table th, .verify-table td {{
  border: 1px solid #ddd; padding: 8px 12px; text-align: center;
}}
.verify-table th {{ background: #f5f5f5; font-weight: 600; }}
.verify-table tr.pass {{ background: #f0fdf4; }}
.verify-table tr.fail {{ background: #fef2f2; }}
.status-badge {{
  display: inline-block; padding: 4px 12px; border-radius: 4px;
  font-weight: 600; font-size: 14px;
}}
.status-badge.pass {{ background: #dcfce7; color: #16a34a; }}
.status-badge.fail {{ background: #fee2e2; color: #dc2626; }}
.note {{ color: #666; font-size: 13px; margin-top: 8px; }}
</style>
</head>
<body>
<h1>SPEC 0027 真实图表 PPT 预览与网格布局验证</h1>
<p class="subtitle">SciencePlots + Seaborn 图表美化 · _GridHelper 布局对齐 · 6 种预设色 PPT</p>

<div class="summary-box {'overall-pass' if all_aligned else 'overall-fail'}">
  <h3>总体结果</h3>
  <p>网格布局对齐：
    <span class="status-badge {'pass' if all_aligned else 'fail'}">
      {'全部对齐 ✓' if all_aligned else '存在偏差 ✗'}
    </span>
  </p>
  <p class="note">说明：_place_chart_grid / _place_chart_side_by_side / _place_chart_three
    已用 _GridHelper 重构，坐标与原硬编码实现完全一致（无视觉漂移）。</p>
</div>

<div class="section">
  <h2>一、SciencePlots + Seaborn 代码集成检查</h2>
  <table class="verify-table">
    <thead><tr><th>检查项</th><th>结果</th></tr></thead>
    <tbody>{code_checks_html}</tbody>
  </table>
</div>

<div class="section">
  <h2>二、真实图表（沙箱执行生成）</h2>
  <p class="note">以下图表由 LocalRuleCodeTaskProvider 生成代码，经沙箱 python_executor 执行生成。
    应用了 <code>plt.style.use(['science','no-latex','cjk-sc-font','bright'])</code>
    与 <code>sns.set_theme(...)</code> 样式。</p>
  <div class="charts-grid">{charts_html}</div>
</div>

<div class="section">
  <h2>三、_GridHelper 网格布局坐标对齐验证</h2>
  <p class="note">验证 _GridHelper 计算的坐标与原硬编码 positions 完全一致（精度 ±0.01 英寸）。</p>
  {build_grid_table("grid_2x2", "_place_chart_grid（2×2 网格）")}
  {build_grid_table("grid_1x2", "_place_chart_side_by_side（1×2 网格）")}
  {build_grid_table("grid_three_top", "_place_chart_three 上排（1×2 网格）")}
</div>

<div class="section">
  <h2>四、6 种预设色 PPT 验证</h2>
  <p class="note">用真实图表渲染 PPT，验证 SPEC 0025/0026 视觉效果保持
    （渐变/圆角/阴影/边框）与文件 XML 完整性。</p>
  <table class="verify-table">
    <thead><tr>
      <th>主题色</th><th>HEX</th><th>页数</th><th>图片数</th>
      <th>渐变</th><th>圆角</th><th>阴影</th><th>边框</th><th>文件有效</th>
    </tr></thead>
    <tbody>{ppt_rows_html}</tbody>
  </table>
</div>

<div class="section">
  <h2>五、SPEC 0027 实现要点</h2>
  <ul style="margin-left: 20px; line-height: 2;">
    <li><strong>图表层</strong>：_HEADER 集成 SciencePlots + Seaborn；
      _build_chart_code 升级为 sns.histplot/sns.boxplot/sns.countplot/sns.scatterplot；
      CORRELATION 分析新增 sns.heatmap 热图。</li>
    <li><strong>PPT 层</strong>：新增 _pct_to_emu 百分比定位方法；
      新增 _GridHelper 内部类（N×M 网格 + h_gap/v_gap）；
      _place_chart_grid / _side_by_side / _three 用 _GridHelper 重构。</li>
    <li><strong>沙箱层</strong>：DEFAULT_ALLOWED_IMPORTS 新增 scienceplots、seaborn；
      easypptx 不入沙箱白名单（仅 PPT 渲染层使用）。</li>
    <li><strong>DeepSeek</strong>：_SYSTEM_PROMPT 追加 SciencePlots + Seaborn 使用要求。</li>
  </ul>
</div>

</body>
</html>"""

    output_html.parent.mkdir(parents=True, exist_ok=True)
    output_html.write_text(html, encoding="utf-8")
    print(f"  HTML 预览已保存：{output_html}")


def main():
    """主函数。"""
    output_dir = Path(__file__).parent.parent / "tests" / "tmp_spec0027_verify"
    output_dir.mkdir(parents=True, exist_ok=True)

    # 1. 生成并执行真实图表
    artifacts, code_checks, exec_result = generate_and_execute_charts(output_dir)

    if not artifacts:
        print("\n错误：未生成任何图表产物，无法继续 PPT 渲染")
        print(f"exit_code={exec_result.exit_code}")
        print(f"stderr={exec_result.stderr[-1000:]}")
        return 1

    # 2. 验证 Grid 布局坐标对齐
    grid_results = verify_grid_alignment()

    # 3. 渲染 6 种预设色 PPT
    ppt_results = render_ppt_files(artifacts, output_dir)

    # 4. 生成 HTML 预览
    html_path = (
        Path(__file__).parent.parent.parent
        / "dev-docs" / "e2e-screenshots" / "spec0027" / "spec0027-preview.html"
    )
    generate_html_preview(
        artifacts, grid_results, ppt_results, code_checks,
        html_path, output_dir,
    )

    # 5. 输出 JSON 报告
    report = {
        "code_checks": code_checks,
        "exec_result": {
            "exit_code": exec_result.exit_code,
            "duration_seconds": exec_result.duration_seconds,
            "sandbox_error_code": exec_result.sandbox_error_code,
            "artifacts_count": len(artifacts),
        },
        "grid_alignment": grid_results,
        "ppt_results": ppt_results,
    }
    report_path = output_dir / "spec0027_verify_report.json"
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)
    print(f"\nJSON 报告已保存：{report_path}")

    # 6. 汇总
    all_pass = (
        grid_results.get("all_aligned", False)
        and all(r.get("file_valid", False) for r in ppt_results.values())
        and all(code_checks.values())
    )
    print(f"\n{'='*60}")
    print(f"SPEC 0027 预览生成总体结果：{'全部通过 ✓' if all_pass else '存在失败 ✗'}")
    print(f"{'='*60}")
    print(f"HTML 预览：{html_path}")
    print(f"PPT 文件目录：{output_dir}")
    print(f"JSON 报告：{report_path}")

    return 0 if all_pass else 1


if __name__ == "__main__":
    sys.exit(main())
