"""PPT 文档渲染器（SPEC 0030 pptxforge 集成）。

从同一份已确认大纲提炼生成 .pptx 文件。
主路径使用 pptxforge Deck（10 主题 + 视觉原语 + Morph 转场），
降级路径保留 python-pptx 命令式渲染（SPEC 0024-0026 布局与视觉）。

设计要点（SPEC 0024，降级路径保留）：
- 16:9 宽屏画布（13.333×7.5 英寸）
- 空白版式（slide_layouts[6]）+ add_textbox/add_picture/add_shape 精确定位
- 双栏内容页：左栏 40% 文本要点 + 右栏 60% 图表
- 图表自适应：单图居中放大、双图并排、3-4 图 2×2 网格
- 五级字号体系：36/28/20/16/12 pt

设计要点（SPEC 0025，降级路径保留）：
- 三角色彩系统：从单一 theme_color 用 colorsys 派生主色/辅助色/强调色
- 深浅对比三明治结构：深色标题栏 → 浅色内容区 → 深色页脚栏

设计要点（SPEC 0026，降级路径保留）：
- 渐变填充：封面顶部色块、标题栏、页脚栏改为线性渐变（主色 → 主色暗化）
- 圆角矩形：左栏背景衬托改为圆角矩形（半径 0.05），柔化硬边缘
- 外阴影效果：右栏图表添加柔和外阴影（oxml 操作 a:effectLst）

设计要点（SPEC 0030，主路径）：
- pptxforge Deck 接管视觉渲染：10 主题系统、视觉原语（TwoColumn/Grid/Centered/Text/Image）
- 主题优先级：theme_preset > theme_color 映射 > 默认 SLATE_MINIMALIST
- theme_color hex → pptxforge 主题映射（_map_theme 基于 HLS 色相和饱和度）
- Morph 转场（default_transition='morph'）
- 降级策略：pptxforge 失败（LayoutOverflowError 等）时降级到 python-pptx 路径
- render() 签名不变，PptConfig 合同扩展为四字段（方案 B）
"""

import colorsys
import logging
import shutil
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from app.core.errors import AppError
from app.modules.outlines.layout_planner import (
    LayoutKind,
    SectionLayoutPlan,
    plan_section_layouts,
)
from app.modules.outlines.document_planner import (
    DefenseSlidePlan,
    plan_defense_deck,
)
from app.modules.outlines.ppt_workflows import resolve_ppt_workflow

logger = logging.getLogger(__name__)


# === SPEC 0024 布局常量 ===

# 16:9 画布尺寸（英寸）
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# 页面边距
MARGIN_LEFT = 0.5
MARGIN_TOP = 0.5
MARGIN_BOTTOM = 0.5

# 双栏布局参数
CONTENT_LEFT_WIDTH = 5.3      # 左栏文本宽度（40%）
CONTENT_RIGHT_LEFT = 6.1      # 右栏起始位置
CONTENT_RIGHT_WIDTH = 6.7     # 右栏图表宽度（60%）

# 字号体系（Pt）
FONT_SIZE_MAIN_TITLE = 36
FONT_SIZE_PAGE_TITLE = 28
FONT_SIZE_SUBTITLE = 20
FONT_SIZE_BODY = 16
FONT_SIZE_CAPTION = 12
# SPEC 0043：答辩页独立于旧 native_editable 路径的最低字号合同。
DEFENSE_TITLE_MIN_PT = 35
DEFENSE_BODY_MIN_PT = 18
DEFENSE_CAPTION_MIN_PT = 12

# 字体
FONT_NAME_CN = "微软雅黑"
FONT_NAME_EN = "Calibri"

# 默认主题色（theme_color=None 时使用，深灰色）
DEFAULT_THEME_COLOR = "333333"

# 色块与分隔线
TITLE_BANNER_HEIGHT = 2.5     # 封面顶部色块高度
DIVIDER_HEIGHT = 0.04         # 分隔线粗细


# === SPEC 0025 三明治结构常量 ===

# 标题栏（深色背景栏）
TITLE_BAR_HEIGHT = 1.0           # 内容页/图表页/总结页标题栏高度

# 页脚栏（深色背景栏）
FOOTER_BAR_HEIGHT = 0.5          # 页脚栏高度
FOOTER_BAR_TOP = SLIDE_HEIGHT - FOOTER_BAR_HEIGHT  # 7.0

# 内容区（三明治夹心：标题栏和页脚栏之间）
CONTENT_TOP = TITLE_BAR_HEIGHT + 0.2    # 1.2（内容起始纵向位置）
CONTENT_BOTTOM = FOOTER_BAR_TOP - 0.2   # 6.8
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP  # 5.6

# 图表可用区域上边界（图表页标题栏下方）
CHART_AREA_TOP = TITLE_BAR_HEIGHT + 0.3  # 1.3

# 文字颜色
TEXT_COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_COLOR_MUTED = RGBColor(0x55, 0x55, 0x55)


class PptRenderer:
    """PPT 文档渲染器（SPEC 0026 视觉效果增强 + SPEC 0027 布局增强）。

    从同一份已确认大纲提炼生成 .pptx 文件。
    采用 16:9 画布 + 空白版式精确定位 + 双栏内容页 + 图表自适应布局
    + 三角色彩系统（主色/辅助色/强调色）+ 深浅对比三明治结构
    + 渐变填充 + 圆角矩形 + 外阴影 + 细边框（SPEC 0026）
    + 百分比定位 + Grid 布局辅助（SPEC 0027，借鉴 EasyPPTX 设计思路）。
    """

    # === SPEC 0027 布局增强辅助方法 ===

    @staticmethod
    def _pct_to_emu(pct_str: str, total_emu: int) -> int:
        """百分比字符串转 EMU（EasyPPTX 风格，SPEC 0027）。

        参数：
        - pct_str: 百分比字符串，如 "10%", "50.5%"
        - total_emu: 总长度（EMU），如 slide_width

        返回：EMU 整数

        异常：非百分比字符串抛出 ValueError。

        示例：
        - _pct_to_emu("10%", Inches(13.333)) → Inches(1.3333)
        - _pct_to_emu("50%", Inches(7.5)) → Inches(3.75)
        """
        if not isinstance(pct_str, str) or not pct_str.endswith("%"):
            raise ValueError(f"百分比字符串必须以 % 结尾：{pct_str}")
        pct = float(pct_str[:-1]) / 100.0
        return int(total_emu * pct)

    class _GridHelper:
        """EasyPPTX 风格的 Grid 布局辅助类（SPEC 0027）。

        给定区域 (left, top, width, height) 和 N×M 网格，
        计算每个单元格的 (left, top, width, height)。
        支持水平和垂直间距。
        """

        def __init__(
            self,
            left: int,
            top: int,
            width: int,
            height: int,
            rows: int,
            cols: int,
            h_gap: int = 0,
            v_gap: int = 0,
        ):
            self.left = left
            self.top = top
            self.width = width
            self.height = height
            self.rows = rows
            self.cols = cols
            self.h_gap = h_gap
            self.v_gap = v_gap

        def cell(self, row: int, col: int) -> tuple[int, int, int, int]:
            """返回 (row, col) 单元格的 (left, top, width, height)。"""
            cell_w = (self.width - self.h_gap * (self.cols - 1)) // self.cols
            cell_h = (self.height - self.v_gap * (self.rows - 1)) // self.rows
            cell_left = self.left + col * (cell_w + self.h_gap)
            cell_top = self.top + row * (cell_h + self.v_gap)
            return (cell_left, cell_top, cell_w, cell_h)

    def render(
        self,
        project_name: str,
        project_topic: str,
        outline_sections: list[dict],
        execution_artifacts: list[dict],
        output_path: str,
        config: dict | None = None,
    ) -> str:
        """渲染 PPT 文档。

        参数：
        - project_name: 项目名称（用于标题页和页脚）
        - project_topic: 项目课题（用于标题页）
        - outline_sections: 已确认大纲的 sections 列表（dict 形式）
        - execution_artifacts: 执行产物列表（含 file_path/name/artifact_type）
        - output_path: 输出文件绝对路径
        - config: PPT 配置（SPEC 0011 + SPEC 0030），可选。支持字段：
          - target_slide_count: 目标页数（5-20），None 表示默认
          - theme_color: 主题色 hex 值，None 表示使用默认深灰色
          - include_charts: 是否包含图表页，默认 True
          - theme_preset: pptxforge 主题名（SPEC 0030），None 时由 theme_color 映射

        返回：生成的文件路径

        异常：渲染失败抛出 AppError(code="PPT_RENDER_FAILED")。

        SPEC 0030：主路径使用 pptxforge Deck，失败时降级到 python-pptx 路径。
        主题优先级：theme_preset > theme_color 映射 > 默认 SLATE_MINIMALIST。
        """
        # 解析 config（SPEC 0011 + SPEC 0030）
        cfg = config or {}
        target_slide_count = cfg.get("target_slide_count")
        theme_color = cfg.get("theme_color")
        include_charts = cfg.get("include_charts", True)
        theme_preset = cfg.get("theme_preset")
        ppt_workflow = cfg.get("ppt_workflow")
        # 直接调用渲染器时也保持与 PptConfig 相同的模式解析边界。
        resolve_ppt_workflow(ppt_workflow)

        # 确保输出目录存在
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)

        # SPEC 0030：主路径 - pptxforge Deck
        try:
            return self._render_with_pptxforge(
                project_name=project_name,
                project_topic=project_topic,
                outline_sections=outline_sections,
                execution_artifacts=execution_artifacts,
                output_path=output_path,
                target_slide_count=target_slide_count,
                theme_color=theme_color,
                include_charts=include_charts,
                theme_preset=theme_preset,
                ppt_workflow=ppt_workflow,
            )
        except Exception as exc:
            # pptxforge 失败时降级到 python-pptx 路径（SPEC 0030 降级策略）
            logger.warning(
                "pptxforge 渲染失败，降级到 python-pptx 路径：%s", exc,
            )

        # 降级路径 - python-pptx 命令式渲染（SPEC 0024-0026）
        try:
            return self._render_with_python_pptx(
                project_name=project_name,
                project_topic=project_topic,
                outline_sections=outline_sections,
                execution_artifacts=execution_artifacts,
                output_path=output_path,
                target_slide_count=target_slide_count,
                theme_color=theme_color,
                include_charts=include_charts,
                ppt_workflow=ppt_workflow,
            )
        except AppError:
            raise
        except Exception as exc:
            raise AppError(
                code="PPT_RENDER_FAILED",
                message=f"PPT 文档生成失败：{exc}",
            ) from exc

    # === SPEC 0030 pptxforge 主路径 ===

    @staticmethod
    def _apply_source_notes(
        output_path: Path, slide_sources: list[tuple[int, str]],
    ) -> None:
        """在 pptxforge 保存后，用 python-pptx 持久化标准 speaker notes。"""

        if not slide_sources:
            return
        presentation = Presentation(str(output_path))
        for slide_index, notes in slide_sources:
            presentation.slides[slide_index].notes_slide.notes_text_frame.text = notes
        presentation.save(str(output_path))

    @staticmethod
    def _build_source_notes(artifacts: tuple[dict, ...] | list[dict]) -> str | None:
        """将外部科研资产和证据追溯写入 speaker notes 的 [Sources] 块。"""

        lines: list[str] = []
        for artifact in artifacts:
            asset_ids = [
                str(value).strip()
                for value in artifact.get("scientific_asset_ids", []) or []
                if str(value).strip()
            ]
            if asset_ids:
                lines.append("assets: " + ", ".join(asset_ids))
            image_sha256 = str(
                artifact.get("scientific_asset_image_sha256", "")
            ).strip()
            if image_sha256:
                lines.append(f"image-sha256: {image_sha256}")
            metadata_path = str(
                artifact.get("scientific_asset_render_metadata", "")
            ).strip()
            if metadata_path:
                lines.append(f"render-metadata: {metadata_path}")
            execution_run_id = str(artifact.get("execution_run_id", "")).strip()
            if execution_run_id:
                lines.append(f"execution-run: {execution_run_id}")
            for attribution in artifact.get("scientific_asset_attributions", []) or []:
                value = str(attribution).strip()
                if value and value not in lines:
                    lines.append(value)
            for source_id in (
                artifact.get("figure_argument", {}).get("evidence_refs", []) or []
            ):
                value = str(source_id).strip()
                if value and value not in lines:
                    lines.append(value)
        if not lines:
            return None
        return "[Sources]\n" + "\n".join(f"- {line}" for line in lines)

    @staticmethod
    def _fit_cover_title(title: str, *, max_chars_per_line: int = 16) -> str:
        """为 pptxforge hero 标题插入稳定换行，避免长中文封面触发溢出降级。"""

        normalized = " ".join(title.split())
        if "\n" in title or len(normalized) <= max_chars_per_line:
            return title.strip()
        punctuation = "：:—-，,"
        midpoint = len(normalized) // 2
        candidates = [
            index + 1
            for index, character in enumerate(normalized)
            if character in punctuation
            and 8 <= index + 1 <= len(normalized) - 6
        ]
        split_at = (
            min(candidates, key=lambda index: abs(index - midpoint))
            if candidates
            else midpoint
        )
        return f"{normalized[:split_at].rstrip()}\n{normalized[split_at:].lstrip()}"
    def _resolve_theme_preset(
        self,
        theme_preset: str | None,
        theme_color: str | None,
        ppt_workflow: str | None = None,
    ) -> str:
        """解析最终 pptxforge 主题名（SPEC 0030 方案 B 优先级）。

        优先级：theme_preset > theme_color 映射 > 默认 SLATE_MINIMALIST。
        """
        if theme_preset:
            return theme_preset
        workflow = resolve_ppt_workflow(ppt_workflow)
        if workflow.default_theme_preset:
            return workflow.default_theme_preset
        return self._map_theme(theme_color)

    def _get_pptxforge_theme(self, theme_name: str):
        """根据主题名获取 pptxforge Theme 对象。

        使用 getattr(themes, name) 获取预定义主题对象。
        无效主题名降级到 SLATE_MINIMALIST。
        """
        from pptxforge import themes

        theme_obj = getattr(themes, theme_name, None)
        if theme_obj is None:
            logger.warning(
                "pptxforge 主题 %s 不存在，降级到 SLATE_MINIMALIST",
                theme_name,
            )
            theme_obj = themes.SLATE_MINIMALIST
        return theme_obj

    def _render_with_pptxforge(
        self,
        project_name: str,
        project_topic: str,
        outline_sections: list[dict],
        execution_artifacts: list[dict],
        output_path: str,
        target_slide_count: int | None,
        theme_color: str | None,
        include_charts: bool,
        theme_preset: str | None,
        ppt_workflow: str | None,
    ) -> str:
        """pptxforge Deck 主路径渲染（SPEC 0030）。

        使用 pptxforge Deck + 主题系统 + 视觉原语生成 PPT。
        Deck.save() 输出到临时目录，再将 deck.pptx 移动到 output_path。
        """
        from pptxforge import Deck, layouts, transitions

        # 解析主题
        theme_name = self._resolve_theme_preset(
            theme_preset, theme_color, ppt_workflow,
        )
        theme_obj = self._get_pptxforge_theme(theme_name)

        # 收集图表产物
        chart_artifacts = (
            [a for a in execution_artifacts
             if a.get("artifact_type") == "CHART_PNG"]
            if include_charts else []
        )

        # 创建 pptxforge Deck（Morph 转场，SPEC 0030）
        deck = Deck(
            theme=theme_obj,
            title=project_name,
            author="实验报告助手",
            transition="morph",
        )

        # 1. 封面页
        deck.add_title_slide(
            title=self._fit_cover_title(project_topic or "实验报告"),
            subtitle=f"项目：{project_name}",
            eyebrow="实验报告助手",
        )

        workflow_id = ppt_workflow or "native_editable"
        if workflow_id in {"academic", "sjtu_academic"}:
            # SPEC 0034：答辩工作流消费共享页序规划，避免把论文段落机械分页。
            deck_plan = plan_defense_deck(
                outline_sections, chart_artifacts,
            )
            if target_slide_count is not None:
                available_slots = max(1, target_slide_count - 1)
                defense_slides = deck_plan.slides[:available_slots]
            else:
                defense_slides = deck_plan.slides

            used_chart_ids = {
                id(artifact)
                for slide in defense_slides
                for artifact in slide.chart_artifacts
            }
            rendered_slide_sources: list[tuple[int, str]] = []
            for slide_index, slide in enumerate(defense_slides, start=1):
                deck.add_content_slide(
                    title=slide.title,
                    layout=self._build_defense_layout(slide),
                )
                source_notes = self._build_source_notes(slide.chart_artifacts)
                if source_notes:
                    rendered_slide_sources.append((slide_index, source_notes))

            remaining_charts = [
                artifact for artifact in chart_artifacts
                if id(artifact) not in used_chart_ids
            ]
            if remaining_charts:
                deck.add_content_slide(
                    title="补充图表",
                    layout=self._build_chart_grid_layout(remaining_charts),
                )
            summary_text = None
        else:
            # native_editable 保留原有兼容布局与页数控制行为。
            content_groups = self._build_content_groups(
                outline_sections, chart_artifacts,
            )
            if target_slide_count is not None:
                content_groups = self._control_slide_count(
                    content_groups, target_slide_count,
                )
            used_chart_ids = {
                id(g["chart"])
                for g in content_groups
                if g["chart"] is not None
            }
            remaining_charts = [
                a for a in chart_artifacts if id(a) not in used_chart_ids
            ]
            rendered_slide_sources = []
            for slide_index, group in enumerate(content_groups, start=1):
                layout = self._build_content_layout(
                    group["sections"], group["chart"],
                )
                deck.add_content_slide(title=group["title"], layout=layout)
                source_notes = self._build_source_notes(
                    [group["chart"]] if group["chart"] else []
                )
                if source_notes:
                    rendered_slide_sources.append((slide_index, source_notes))
            if remaining_charts:
                chart_layout = self._build_chart_grid_layout(remaining_charts)
                deck.add_content_slide(title="关键图表", layout=chart_layout)
                source_notes = self._build_source_notes(remaining_charts)
                if source_notes:
                    rendered_slide_sources.append(
                        (len(content_groups) + 1, source_notes)
                    )
            summary_text = self._extract_summary_text(outline_sections)

        # SPEC 0030：原生工作流继续使用 closing slide；答辩工作流已经由
        # DefenseDeckPlan 生成结论页，不再追加重复的“总结”页。
        if summary_text is not None:
            if not summary_text:
                summary_text = "本实验已按既定方案完成数据分析与可视化。"

            if len(summary_text) <= 15:
                deck.add_closing_slide(message=summary_text)
            else:
                from pptxforge import layouts as _layouts
                deck.add_content_slide(
                    title="总结",
                    layout=_layouts.Text(summary_text, role="lead"),
                )

        # Deck.save() 输出到临时目录，再移动到 output_path
        # save() 会创建 deck.pptx + README.md，只需保留 deck.pptx
        output = Path(output_path)
        with tempfile.TemporaryDirectory() as tmp_dir:
            saved_path = deck.save(tmp_dir)
            # saved_path 是 deck.pptx 的完整路径
            shutil.move(str(saved_path), str(output))
        self._apply_source_notes(output, rendered_slide_sources)
        if workflow_id in {"academic", "sjtu_academic"}:
            result_slide_indices = tuple(
                index + 1
                for index, slide in enumerate(defense_slides)
                if slide.role in {
                    "result",
                    "primary",
                    "stratified",
                    "diagnosis_stratified",
                    "comparison",
                    "comparison_matrix",
                }
                and slide.chart_artifacts
            )
            self._postprocess_defense_presentation(
                output,
                result_slide_indices=result_slide_indices,
            )

        return str(output)

    @classmethod
    def _postprocess_defense_presentation(
        cls,
        output_path: Path,
        *,
        result_slide_indices: tuple[int, ...],
    ) -> None:
        """在 pptxforge 保存后收口答辩页的版式合同（SPEC 0043）。

        pptxforge 负责主题与语义布局，但其图片和字体会随布局内容发生
        缩放。这里以最终 .pptx 为真源做一次确定性收口：结果页主图统一为
        10.2×5.95 英寸（1.714 比例），并逐 run 修正答辩页最低字号。
        图片宽高分别设置而不是拉伸填满容器；固定比例与常见论文图表资源
        一致，避免非等比变形。
        """

        presentation = Presentation(str(output_path))
        result_indices = set(result_slide_indices)
        target_width = Inches(10.2)
        target_height = Inches(5.95)
        centered_left = int((presentation.slide_width - target_width) / 2)
        image_top = Inches(1.18)

        for slide_index, slide in enumerate(presentation.slides):
            if slide_index == 0:
                continue

            if slide_index in result_indices:
                pictures = [
                    shape
                    for shape in slide.shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                ]
                if pictures:
                    picture = max(
                        pictures,
                        key=lambda shape: shape.width * shape.height,
                    )
                    picture.left = centered_left
                    picture.top = image_top
                    picture.width = target_width
                    picture.height = target_height

            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if shape.text_frame.text.strip() in {
                    "解释边界",
                    "补充图表用于从不同角度核对同一结果",
                }:
                    shape.text_frame.clear()
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text or not run.text.strip():
                            continue
                        minimum = cls._defense_run_minimum_pt(shape, run.text)
                        current = run.font.size.pt if run.font.size else 0
                        if current < minimum:
                            run.font.size = Pt(minimum)

        presentation.save(str(output_path))

    @staticmethod
    def _defense_run_minimum_pt(shape, text: str) -> int:
        """根据最终形状位置和文案角色返回答辩页最低字号。"""

        if shape.top < Inches(1.2):
            return DEFENSE_TITLE_MIN_PT
        if "图" in text or "来源" in text:
            return DEFENSE_CAPTION_MIN_PT
        return DEFENSE_BODY_MIN_PT

    def _build_defense_layout(self, slide: DefenseSlidePlan):
        """构建高级答辩页：每页一个表达任务（SPEC 0034）。"""

        from pptxforge import layouts

        content = " ".join(slide.content.split()) or "（暂无内容）"
        if slide.role == "question":
            return layouts.TwoColumn(
                left=layouts.Stack(
                    children=[
                        layouts.Text("研究对象", role="h2"),
                        layouts.Text(content[:260], role="body", align="left"),
                    ],
                    gap=0.18,
                    weights=[0.45, 1.25],
                ),
                right=layouts.Stack(
                    children=[
                        layouts.Text("本次解读回答什么", role="h2"),
                        layouts.Text("问题  →  方法  →  证据  →  结论", role="body"),
                    ],
                    gap=0.22,
                    weights=[0.75, 1.0],
                ),
                split=0.62,
                gap=0.45,
            )

        if slide.role in {"source", "limitation"}:
            if slide.role == "source":
                return layouts.TwoColumn(
                    left=layouts.Callout(
                        "原论文：Cerner Health Facts 数据库；最终分析样本 69984 条。",
                        eyebrow="论文原文",
                        variant="info",
                    ),
                    right=layouts.Callout(
                        "公开复核：UCI 原始 CSV 101766 条；保留原始口径，不伪装成论文筛选样本。",
                        eyebrow="本地数据",
                        variant="key",
                    ),
                    split=0.5,
                    gap=0.35,
                )
            return layouts.Stack(
                children=[
                    layouts.Text(
                        "论文原文" if slide.role == "source" else "解释边界",
                        role="h2",
                    ),
                    layouts.Callout(
                        content[:180],
                        eyebrow="解释边界",
                        variant="warning",
                    ) if slide.role == "limitation" else layouts.Text(content[:420], role="body"),
                ],
                gap=0.35,
                weights=[0.55, 2.2],
            )

        if slide.role == "data" and slide.metrics:
            stats = [
                layouts.StatRowStat(value=self._compact_stat_value(value), label=label[:18])
                for label, value in slide.metrics
            ]
            return layouts.Stack(
                children=[
                    layouts.StatRow(stats),
                    layouts.Text(content[:260], role="body"),
                ],
                gap=0.34,
                weights=[1.0, 1.45],
            )

        if slide.role == "method" and slide.steps:
            items = [
                layouts.IconRowItem(
                    icon=layouts.Icon("circle"),
                    header=f"{index + 1:02d}",
                    description=step[:78],
                )
                for index, step in enumerate(slide.steps[:5])
            ]
            return layouts.Stack(
                children=[
                    layouts.Text("从数据到证据", role="h2"),
                    layouts.IconRow(items, gap=0.25),
                ],
                gap=0.36,
                weights=[0.55, 2.0],
            )

        if slide.role in {
            "sample_structure", "quality", "model", "primary",
            "stratified", "diagnosis_stratified", "comparison",
            "comparison_matrix", "evidence_chain",
            "data_pipeline", "relationship_graph",
        } and slide.chart_artifacts:
            plan = SectionLayoutPlan(
                section={
                    "title": slide.title,
                    "content": slide.content,
                    "figure_lead": slide.figure_lead,
                    "figure_takeaway": slide.figure_takeaway,
                },
                layout_kind=slide.layout_kind,
                chart_artifacts=slide.chart_artifacts,
                text_density="medium",
            )
            return self._build_adaptive_layout(plan)

        if slide.role == "sample":
            if slide.metrics:
                stats = [
                    layouts.StatRowStat(value=self._compact_stat_value(value), label=label[:16])
                    for label, value in slide.metrics[:4]
                ]
                return layouts.Stack(
                    children=[
                        layouts.StatRow(stats),
                        layouts.Callout(
                            "样本口径必须先于结果解释：公开 CSV 与论文最终队列不是同一筛选结果。",
                            eyebrow="口径先行",
                            variant="key",
                        ),
                    ],
                    gap=0.28,
                    weights=[1.0, 1.0],
                )
            return layouts.Stack(
                children=[
                    layouts.Text("口径先行", role="h2"),
                    layouts.Text(content[:360], role="body"),
                ],
                gap=0.35,
                weights=[0.55, 2.0],
            )

        if slide.role == "result" and slide.chart_artifacts:
            plan = SectionLayoutPlan(
                section={
                    "title": slide.title,
                    "content": slide.content,
                    "figure_lead": slide.figure_lead,
                    "figure_takeaway": slide.figure_takeaway,
                },
                layout_kind=slide.layout_kind,
                chart_artifacts=slide.chart_artifacts,
                text_density="medium",
            )
            return self._build_adaptive_layout(plan)

        if slide.role == "conclusion":
            return layouts.Stack(
                children=[
                    layouts.Text("结论先行", role="h2"),
                    layouts.Text(content[:300], role="lead"),
                    layouts.Text(
                        "证据范围：已确认大纲与真实执行产物；不延伸为因果或临床结论。",
                        role="caption",
                    ),
                ],
                gap=0.25,
                weights=[0.45, 2.1, 0.42],
            )

        return layouts.Text(content[:360], role="body")

    def _build_adaptive_layout(self, plan: SectionLayoutPlan):
        """根据语义计划选择 pptxforge 构图（SPEC 0033）。"""

        from pptxforge import layouts

        content = " ".join(plan.content.split()) or "（暂无内容）"

        # SPEC 0039：逻辑图与数据图共享 FigurePlan，但逻辑图使用单一
        # 主结论布局，避免被统计图的比较网格压缩成普通图片。
        if plan.chart_artifacts and any(
            str(artifact.get("figure_kind", "data_chart")) != "data_chart"
            for artifact in plan.chart_artifacts
        ):
            if str(plan.chart_artifacts[0].get("figure_visual_family", "")) == "matrix":
                return self._build_matrix_figure_layout(
                    plan.chart_artifacts[0], content,
                )
            return self._build_semantic_figure_layout(
                plan.chart_artifacts[0], content,
            )
        if plan.chart_artifacts and plan.chart_artifacts[0].get("figure_argument"):
            return self._build_argument_chart_layout(
                plan.chart_artifacts[0], content,
            )
        if plan.layout_kind == LayoutKind.NARRATIVE:
            role = "lead" if plan.text_density == "low" else "body"
            return layouts.Centered(
                child=layouts.Text(content[:360], role=role, align="left")
            )

        if plan.layout_kind == LayoutKind.DATA_OVERVIEW and plan.metrics:
            stats = [
                layouts.StatRowStat(value=self._compact_stat_value(value), label=label[:18])
                for label, value in plan.metrics
            ]
            return layouts.Stack(
                children=[
                    layouts.StatRow(stats),
                    layouts.Text(content[:240], role="body"),
                ],
                gap=0.35,
                weights=[1.0, 1.7],
            )

        if plan.layout_kind == LayoutKind.METHOD_FLOW and plan.steps:
            items = [
                layouts.IconRowItem(
                    icon=layouts.Icon("circle"),
                    header=f"步骤 {index + 1}",
                    description=step[:72],
                )
                for index, step in enumerate(plan.steps[:4])
            ]
            return layouts.Stack(
                children=[
                    layouts.Text("分析路径", role="h2"),
                    layouts.IconRow(items, gap=0.3),
                ],
                gap=0.4,
                weights=[0.55, 2.0],
            )

        if plan.layout_kind == LayoutKind.RESULT_FOCUS and plan.chart_artifacts:
            artifact = plan.chart_artifacts[0]
            chart_kind = str(artifact.get("chart_kind", ""))
            chart = layouts.Image(str(artifact.get("file_path", "")), fit="contain")
            caption = self._chart_caption(artifact)
            figure_lead = str(plan.section.get("figure_lead", "")).strip()
            figure_takeaway = str(plan.section.get("figure_takeaway", "")).strip()
            if chart_kind == "dumbbell":
                return layouts.Stack(
                    children=[
                        layouts.StatRow([
                            layouts.StatRowStat(value="18.4%", label="论文原文检测率"),
                            layouts.StatRowStat(value="16.7%", label="本地 CSV 检测率"),
                        ]),
                        chart,
                        layouts.Text(
                            figure_takeaway or "两者分析口径不同，图中只呈现配对差异，不宣称精确复现。",
                            role="caption",
                        ),
                    ],
                    gap=0.18,
                    weights=[1.3, 2.8, 0.5],
                )
            if chart_kind in {"flow", "stacked_composition"}:
                return layouts.Stack(
                    children=[
                    layouts.Text(figure_lead or "先确认分析口径，再阅读图中的阶段与比例。", role="body"),
                    layouts.Callout(
                        content[:170],
                        eyebrow="先读结论",
                        variant="key",
                        ),
                        chart,
                        layouts.Text(caption, role="caption"),
                        layouts.Text(figure_takeaway, role="caption") if figure_takeaway else layouts.Text("", role="caption"),
                    ],
                    gap=0.18,
                    weights=[0.56, 1.08, 2.62, 0.40, 0.28],
                )
            return layouts.Stack(
                children=[
                    layouts.Text(figure_lead or "先读图中估计值，再结合方法和边界解释。", role="body"),
                    layouts.TwoColumn(
                        left=layouts.Stack(
                            children=[
                                layouts.Callout(
                                    content[:160],
                                    eyebrow="结果解释",
                                    variant="key",
                                ),
                                layouts.Text(
                                    str(artifact.get("chart_rationale", ""))[:180],
                                    role="caption",
                                ),
                            ],
                            gap=0.25,
                            weights=[1.3, 0.8],
                        ),
                        right=chart,
                        split=0.38,
                        gap=0.35,
                    ),
                    layouts.Text(
                        "；".join(
                            value for value in (caption, figure_takeaway) if value
                        ),
                        role="caption",
                    ),
                ],
                gap=0.14,
                weights=[0.56, 3.95, 0.74],
            )

        if plan.layout_kind == LayoutKind.RESULT_COMPARE and plan.chart_artifacts:
            callout_text = content.split("。", 1)[0].strip()
            if callout_text:
                callout_text += "。"
            cards = [
                layouts.Stack(
                    children=[
                        layouts.Image(
                            str(artifact.get("file_path", "")),
                            fit="contain",
                        ),
                        layouts.Text(
                            self._chart_caption(artifact),
                            role="caption",
                        ),
                    ],
                    gap=0.08,
                    weights=[4.0, 1.15],
                )
                for artifact in plan.chart_artifacts[:4]
            ]
            chart_grid = layouts.Grid(
                children=cards[:2],
                cols=2,
                rows=1,
                gap=0.35,
            )
            return layouts.Stack(
                children=[
                    layouts.Text(
                        str(plan.section.get("figure_lead", "")).strip() or "先说明比较口径，再并列阅读图形中的差异。",
                        role="body",
                    ),
                    layouts.Callout(
                        callout_text[:72],
                        eyebrow="结果解释",
                        variant="key",
                    ),
                    chart_grid,
                    layouts.Text(
                        str(plan.section.get("figure_takeaway", "")).strip(),
                        role="caption",
                    ) if str(plan.section.get("figure_takeaway", "")).strip() else layouts.Text("", role="caption"),
                ],
                gap=0.18,
                weights=[0.56, 1.10, 2.90, 0.38],
            )

        if plan.layout_kind == LayoutKind.SUMMARY:
            return layouts.Centered(
                child=layouts.Text(
                    content[:280],
                    role="lead" if len(content) <= 150 else "body",
                    align="center",
                )
            )

        return layouts.Text(content[:360], role="body")

    @staticmethod
    def _chart_caption(artifact: dict) -> str:
        if artifact.get("figure_caption"):
            return str(artifact["figure_caption"])
        name = artifact.get("name") or "未命名图表"
        kind = str(artifact.get("chart_kind", "")).replace("_", " ").strip()
        suffix = f" · {kind}" if kind else ""
        return f"图表：{Path(str(name)).stem.replace('_', ' ')[:36]}{suffix}"

    def _build_semantic_figure_layout(self, artifact: dict, content: str):
        """用现有 pptxforge 原语呈现论证型流程/关系/证据链逻辑图。"""

        from pptxforge import layouts

        kind = str(artifact.get("figure_kind", "semantic_figure"))
        labels = {
            "evidence_chain": "证据链",
            "data_pipeline": "数据处理路径",
            "relationship_graph": "变量关系",
            "research_framework": "研究框架",
            "process_flow": "研究流程",
            "timeline": "研究时间线",
        }
        eyebrow = labels.get(kind, "论文逻辑图")
        argument = artifact.get("figure_argument") or {}
        callout_text = argument.get("claim") or {
            "evidence_chain": "四步证据链：来源 → 数据 → 复核 → 边界。",
            "data_pipeline": "五段数据管线：原始数据 → 质量 → 分组 → 结局 → 模型。",
            "relationship_graph": "观察性关联：不绘制确定性因果箭头。",
            "research_framework": "研究框架：问题、变量、证据与结论分层组织。",
            "process_flow": "按实际研究步骤读取流程，不把阶段误读成数值比较。",
            "timeline": "按时间顺序呈现研究阶段与事件节点。",
        }.get(kind, "逻辑图主线与证据边界见图下说明。")
        if kind == "evidence_chain":
            callout_text = "复核不等于复现。"
        elif kind == "data_pipeline":
            callout_text = "路径可复核。"
        elif kind == "relationship_graph":
            callout_text = "观察性关联。"
        note = str(artifact.get("figure_note", ""))
        caption = self._chart_caption(artifact)
        result = str(argument.get("result", ""))
        boundary = str(argument.get("boundary", ""))
        evidence_labels = {
            "paper:PMC3996476": "开放论文",
            "dataset:UCI-296": "UCI 数据集",
            "spec0039_semantic_figures": "本地执行批次",
        }
        evidence = "、".join(
            evidence_labels.get(str(value), Path(str(value)).stem)
            for value in argument.get("evidence_refs", [])[:3]
        )
        proof_line = "；".join(
            value for value in (
                f"证据：{evidence}" if evidence else "",
                f"结果：{result[:42]}" if result else "",
                f"边界：{boundary[:58]}" if boundary else note[:58],
            ) if value
        )
        # 多面板证据论证图本身已经承担了主要叙事，不能继续沿用
        # “大主张条 + 缩略图 + 说明”的普通逻辑图模板。否则 A-D
        # 四个面板会在 PPT 中被压缩成不可读的缩略图。
        if kind == "evidence_chain":
            return layouts.Stack(
                children=[
                    layouts.Image(str(artifact.get("file_path", "")), fit="contain"),
                    layouts.Text(
                        f"{caption}。{boundary[:58]}"[:120],
                        role="caption",
                    ),
                ],
                gap=0.12,
                weights=[5.2, 0.42],
            )
        return layouts.Stack(
            children=[
                layouts.Callout(
                    str(callout_text)[:180],
                    eyebrow=eyebrow,
                    variant="warning" if kind == "relationship_graph" else "info",
                ),
                layouts.Image(str(artifact.get("file_path", "")), fit="contain"),
                layouts.Text(
                    f"{caption}。边界：{boundary[:36]}"[:105],
                    role="caption",
                ),
            ],
            gap=0.18,
            weights=[1.34, 3.04, 0.55],
        )

    def _build_matrix_figure_layout(self, artifact: dict, content: str):
        """比较矩阵使用表格主视觉，不套用流程/关系图的节点布局。"""

        from pptxforge import layouts

        argument = artifact.get("figure_argument") or {}
        boundary = str(argument.get("boundary") or artifact.get("figure_note", ""))
        return layouts.Stack(
            children=[
                layouts.Text(
                    "论文与本地复核：可对照，不等同。",
                    role="h2",
                ),
                layouts.Image(str(artifact.get("file_path", "")), fit="contain"),
                layouts.Text(
                    f"{self._chart_caption(artifact)}。边界：{boundary[:72]}"[:150],
                    role="caption",
                ),
            ],
            gap=0.16,
            weights=[0.95, 4.45, 0.42],
        )

    def _build_argument_chart_layout(self, artifact: dict, content: str):
        """统计图也采用主张—证据—结果—边界的答辩页结构。"""

        from pptxforge import layouts

        argument = artifact.get("figure_argument") or {}
        chart_kind = str(artifact.get("chart_kind", ""))
        semantic_claims = {
            "flow": "口径先于比较。",
            "stacked_composition": "先读总体构成。",
            "horizontal_bar": "缺失结构并不均匀。",
            "dumbbell": "口径不同，不作复现误差。",
            "point_ci": "差异需结合置信区间。",
            "ordered_line": "分层差异不等于因果。",
            "forest": "简化模型不等同原模型。",
        }
        claim = semantic_claims.get(
            chart_kind,
            str(argument.get("claim") or content).strip()[:84],
        )
        evidence = "本地执行图表"
        result_lines = {
            "flow": "公开 CSV 与论文样本口径不同。",
            "stacked_composition": "三类结局构成清晰。",
            "horizontal_bar": "缺失率集中于少数字段。",
            "dumbbell": "论文 18.4%，本地 16.7%。",
            "point_ci": "已检测组率低于未检测组。",
            "ordered_line": "分层后存在结构差异。",
            "forest": "主要变量方向与区间可见。",
        }
        boundary_lines = {
            "flow": "不做口径等价。",
            "stacked_composition": "仅描述总体构成。",
            "horizontal_bar": "不做完整案例推断。",
            "dumbbell": "不解释为复现误差。",
            "point_ci": "不作因果解释。",
            "ordered_line": "不作临床风险分层。",
            "forest": "不等同原论文模型。",
        }
        result = result_lines.get(chart_kind, str(argument.get("result", "")).strip()[:26])
        boundary = boundary_lines.get(
            chart_kind,
            str(argument.get("boundary", artifact.get("figure_note", ""))).strip()[:26],
        )
        caption = self._chart_caption(artifact)
        summary = "\n".join(
            value for value in (
                f"证据：{evidence}" if evidence else "",
                f"结果：{result}" if result else "",
                f"边界：{boundary}" if boundary else "",
            ) if value
        )
        return layouts.Stack(
            children=[
                layouts.Callout(
                    claim[:180],
                    eyebrow="本页主张",
                    variant="key",
                ),
                layouts.TwoColumn(
                    left=layouts.Stack(
                        children=[
                            layouts.Text("证据与结果", role="h2"),
                            layouts.Text(summary[:260], role="body"),
                        ],
                        gap=0.2,
                        weights=[0.45, 1.0],
                    ),
                    right=layouts.Image(str(artifact.get("file_path", "")), fit="contain"),
                    split=0.35,
                    gap=0.35,
                ),
                layouts.Text(caption, role="caption"),
            ],
            gap=0.18,
            weights=[1.34, 3.16, 0.28],
        )

    @staticmethod
    def _compact_stat_value(value: str) -> str:
        """把长样本量压缩为答辩页可读值，精确值仍保留在正文。"""

        normalized = str(value or "")
        if normalized.endswith(" 条"):
            raw = normalized[:-2].strip()
            try:
                number = int(raw)
            except ValueError:
                return normalized[:14]
            if number >= 10000:
                return f"{number / 1000:.1f}k"
        return normalized[:14]

    def _build_content_layout(
        self,
        sections: list[dict],
        chart_artifact: dict | None,
    ):
        """构建内容页布局（pptxforge 视觉原语映射，SPEC 0030）。

        - 有图表：TwoColumn(left=Text, right=Stack(Image + caption), split=0.4)
        - 无图表：Text（全幅文本）

        内容截断策略（避免 LayoutOverflowError）：
        - TwoColumn 布局：文本列窄（40%），每段截断 120 字，最多 3 段
        - 全幅 Text 布局：文本列宽（100%），每段截断 200 字，最多 3 段
        """
        from pptxforge import layouts

        # 根据布局类型确定截断参数
        has_chart = (
            chart_artifact is not None
            and chart_artifact.get("file_path")
            and Path(chart_artifact["file_path"]).exists()
        )
        max_sections = 3
        max_chars_per_section = 120 if has_chart else 200

        def wrap_for_column(text: str, width: int = 16) -> str:
            """为中文窄栏显式换行，避免 pptxforge 将整段视为不可断开的长行。"""
            normalized = " ".join(str(text or "").split())
            if not normalized:
                return ""
            return "\n".join(
                normalized[index:index + width]
                for index in range(0, len(normalized), width)
            )

        # 合并章节正文；页面标题已经由 Deck 统一渲染，避免标题在正文中重复出现。
        text_parts: list[str] = []
        for section in sections[:max_sections]:
            title = section.get("title", "")
            content = section.get("content", "")
            # 截断过长内容（pptxforge LayoutOverflowError 保护）
            short_content = (
                content[:max_chars_per_section]
                + ("…" if len(content) > max_chars_per_section else "")
            )
            text_parts.append(
                wrap_for_column(short_content or title or "（暂无内容）")
            )
        combined_text = "\n\n".join(text_parts) or "（暂无内容）"

        # 有图表：双栏布局（左 40% 文本 + 右 60% 图表）
        if chart_artifact:
            if has_chart:
                left = layouts.Text(combined_text, role="body")
                chart_name = chart_artifact.get("name") or "未命名图表"
                right = layouts.Stack(
                    children=[
                        layouts.Image(str(chart_artifact["file_path"]), fit="contain"),
                        layouts.Text(f"图表：{chart_name[:28]}", role="caption"),
                    ],
                        gap=0.08,
                        weights=[5.0, 0.8],
                )
                return layouts.TwoColumn(
                    left=left, right=right, split=0.4, gap=0.3,
                )
            # 图表文件不存在时，用占位文本替代
            name = chart_artifact.get("name", "")
            left = layouts.Text(combined_text, role="body")
            right = layouts.Text(f"[图片文件不存在：{name}]", role="caption")
            return layouts.TwoColumn(
                left=left, right=right, split=0.4, gap=0.3,
            )

        # 无图表：全幅文本
        return layouts.Text(combined_text, role="body")

    def _build_chart_grid_layout(self, chart_artifacts: list[dict]):
        """构建图表页 Grid 布局（pptxforge 视觉原语映射，SPEC 0030）。

        自适应布局：
        - 1 图：Centered(Image)
        - 2 图：Grid(cols=2, rows=1)
        - 3-4 图：Grid(cols=2, rows=2)
        - >4 图：截断到 4 张
        """
        from pptxforge import layouts

        # 过滤出存在的图表文件
        valid_charts = [
            a for a in chart_artifacts
            if a.get("file_path") and Path(a["file_path"]).exists()
        ][:4]  # 最多 4 张

        if not valid_charts:
            return layouts.Text("[无可用图表]", role="caption")

        # 构建图片原语列表
        images = []
        for artifact in valid_charts:
            chart_name = artifact.get("name") or "未命名图表"
            images.append(
                layouts.Stack(
                    children=[
                        layouts.Image(str(artifact["file_path"]), fit="contain"),
                        layouts.Text(f"图表：{chart_name[:28]}", role="caption"),
                    ],
                        gap=0.06,
                        weights=[5.0, 0.8],
                )
            )

        count = len(images)
        if count == 1:
            return layouts.Centered(child=images[0])
        elif count == 2:
            return layouts.Grid(children=images, cols=2, rows=1, gap=0.2)
        else:
            # 3-4 图：2×2 网格
            return layouts.Grid(children=images, cols=2, rows=2, gap=0.2)

    def _extract_summary_text(self, outline_sections: list[dict]) -> str:
        """提取 SUMMARY 类型章节作为总结文本。"""
        summary_sections = [
            s for s in outline_sections
            if s.get("source_type") == "SUMMARY"
        ]
        if not summary_sections:
            return ""
        parts = [s.get("content", "") for s in summary_sections if s.get("content")]
        return "\n\n".join(parts)

    # === python-pptx 降级路径（SPEC 0024-0026） ===

    def _render_defense_with_python_pptx(
        self,
        project_name: str,
        project_topic: str,
        outline_sections: list[dict],
        chart_artifacts: list[dict],
        output_path: str,
        target_slide_count: int | None,
        primary: RGBColor,
        title_text_color: RGBColor,
    ) -> str:
        """SPEC 0043 fallback：保持答辩页主视觉优先，不退回旧双栏缩略图。"""
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        deck_plan = plan_defense_deck(outline_sections, chart_artifacts)
        defense_slides = deck_plan.slides
        if target_slide_count is not None:
            defense_slides = defense_slides[:max(1, target_slide_count - 1)]
        total_pages = 1 + len(defense_slides) + 1
        self._render_title_slide(
            prs, project_name, project_topic, primary, title_text_color,
        )
        for index, slide_plan in enumerate(defense_slides, start=1):
            self._add_defense_fallback_slide(
                prs, slide_plan, primary, title_text_color,
                page_num=index + 1, total_pages=total_pages,
                project_name=project_name,
            )
            notes = self._build_source_notes(slide_plan.chart_artifacts)
            if notes:
                prs.slides[index].notes_slide.notes_text_frame.text = notes
        self._render_summary_slide(
            prs, outline_sections, primary, title_text_color,
            page_num=total_pages, total_pages=total_pages,
            project_name=project_name,
        )
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output))
        return str(output)

    def _add_defense_fallback_slide(
        self,
        prs: Presentation,
        slide_plan,
        primary: RGBColor,
        title_text_color: RGBColor,
        page_num: int,
        total_pages: int,
        project_name: str,
    ) -> None:
        """渲染 fallback 答辩页：结果图占据主视觉，文字只承担讲述职责。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_page_title(
            slide, slide_plan.title, primary, title_text_color,
            font_size=DEFENSE_TITLE_MIN_PT,
        )
        artifacts = list(slide_plan.chart_artifacts or ())
        if artifacts:
            artifact = artifacts[0]
            file_path = str(artifact.get("file_path", ""))
            if file_path and Path(file_path).exists():
                # 10.2×5.95 英寸约占 60.7% 画布，且底部避开页脚。
                picture = slide.shapes.add_picture(
                    file_path,
                    Inches(1.566), Inches(1.02),
                    width=Inches(10.2), height=Inches(5.95),
                )
                self._style_picture(
                    picture, self._derive_color_palette(primary)[1],
                )
        else:
            text_box = slide.shapes.add_textbox(
                Inches(0.85), Inches(1.45),
                Inches(SLIDE_WIDTH - 1.7), Inches(4.9),
            )
            text_box.text_frame.word_wrap = True
            run = text_box.text_frame.paragraphs[0].add_run()
            run.text = " ".join(str(slide_plan.content or "").split())
            self._set_run_font(
                run, DEFENSE_BODY_MIN_PT, TEXT_COLOR_DARK,
            )
        self._add_footer(
            slide, project_name, page_num, total_pages,
            primary, title_text_color,
        )
    def _render_with_python_pptx(
        self,
        project_name: str,
        project_topic: str,
        outline_sections: list[dict],
        execution_artifacts: list[dict],
        output_path: str,
        target_slide_count: int | None,
        theme_color: str | None,
        include_charts: bool,
        ppt_workflow: str | None = None,
    ) -> str:
        """python-pptx 命令式渲染降级路径（SPEC 0024-0026）。

        当 pptxforge 主路径失败时使用，保留 SPEC 0024-0026 布局与视觉效果。
        """
        try:
            prs = Presentation()
            # SPEC 0024: 16:9 宽屏画布
            prs.slide_width = Inches(SLIDE_WIDTH)
            prs.slide_height = Inches(SLIDE_HEIGHT)

            # SPEC 0025: 从主题色派生三角色彩（主色/辅助色/强调色/标题文字色）
            workflow = resolve_ppt_workflow(ppt_workflow)
            fallback_color = (
                workflow.fallback_theme_color
                if workflow.fallback_theme_color
                else theme_color
            )
            theme_rgb = self._resolve_theme_color(fallback_color)
            primary, auxiliary, accent, title_text_color = (
                self._derive_color_palette(theme_rgb)
            )

            # 收集图表产物
            chart_artifacts = (
                [a for a in execution_artifacts
                 if a.get("artifact_type") == "CHART_PNG"]
                if include_charts else []
            )

            # SPEC 0043：academic/sjtu fallback 继续消费 DefenseDeckPlan。
            # 不让主路径偶发溢出后退化成旧双栏缩略图。
            if ppt_workflow in {"academic", "sjtu_academic"}:
                return self._render_defense_with_python_pptx(
                    project_name=project_name,
                    project_topic=project_topic,
                    outline_sections=outline_sections,
                    chart_artifacts=chart_artifacts,
                    output_path=output_path,
                    target_slide_count=target_slide_count,
                    primary=primary,
                    title_text_color=title_text_color,
                )
            # 构建内容页候选（按 source_type 分组，关联图表）
            content_groups = self._build_content_groups(
                outline_sections, chart_artifacts,
            )

            # 页数控制（SPEC 0011 逻辑，不改变）
            if target_slide_count is not None:
                content_groups = self._control_slide_count(
                    content_groups, target_slide_count,
                )

            # 计算剩余图表（未被内容页消耗的）
            # dict 不可哈希，用 id() 去重
            used_chart_ids = {
                id(g["chart"])
                for g in content_groups
                if g["chart"] is not None
            }
            remaining_charts = [
                a for a in chart_artifacts if id(a) not in used_chart_ids
            ]

            # 计算总页数
            total_pages = (
                1  # 封面页
                + len(content_groups)  # 内容页
                + (1 if remaining_charts else 0)  # 图表页
                + 1  # 总结页
            )

            # 1. 封面页
            self._render_title_slide(
                prs, project_name, project_topic,
                primary, title_text_color,
            )

            # 2-N. 内容页（双栏布局）
            for i, group in enumerate(content_groups):
                self._add_content_slide(
                    prs,
                    title=group["title"],
                    sections=group["sections"],
                    primary=primary,
                    auxiliary=auxiliary,
                    accent=accent,
                    title_text_color=title_text_color,
                    chart_artifact=group["chart"],
                    page_num=i + 2,
                    total_pages=total_pages,
                    project_name=project_name,
                )

            # 图表页（剩余图表）
            page_offset = 2 + len(content_groups)
            if remaining_charts:
                self._add_chart_slide(
                    prs, remaining_charts,
                    primary, title_text_color,
                    page_num=page_offset,
                    total_pages=total_pages,
                    project_name=project_name,
                )
                page_offset += 1

            # 总结页
            self._render_summary_slide(
                prs, outline_sections,
                primary, title_text_color,
                page_num=page_offset,
                total_pages=total_pages,
                project_name=project_name,
            )

            # 所有 fallback 页也持久化 [Sources]，避免主路径失败后丢失追溯。
            for slide_index, group in enumerate(content_groups, start=1):
                source_notes = self._build_source_notes(
                    [group["chart"]] if group["chart"] else []
                )
                if source_notes:
                    prs.slides[slide_index].notes_slide.notes_text_frame.text = source_notes
            if remaining_charts:
                source_notes = self._build_source_notes(remaining_charts)
                if source_notes:
                    prs.slides[1 + len(content_groups)].notes_slide.notes_text_frame.text = source_notes

            # 确保输出目录存在
            output = Path(output_path)
            output.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output))
            return str(output)
        except AppError:
            raise
        except Exception as exc:
            raise AppError(
                code="PPT_RENDER_FAILED",
                message=f"PPT 文档生成失败：{exc}",
            ) from exc

    # === 主题色处理 ===

    def _resolve_theme_color(self, theme_color: str | None) -> RGBColor:
        """解析主题色，None 时降级到默认深灰色（SPEC 0024）。

        SPEC 0011 中 None 表示默认黑色；
        SPEC 0024 改为默认深灰色 #333333，视觉更柔和。
        解析失败时记录 warning 并返回默认深灰色。
        """
        if not theme_color:
            return RGBColor.from_string(DEFAULT_THEME_COLOR)
        try:
            hex_str = theme_color.lstrip("#")
            return RGBColor.from_string(hex_str)
        except Exception as exc:
            logger.warning(
                "PPT 主题色解析失败，降级到默认深灰色：%s (value=%s)",
                exc, theme_color,
            )
            return RGBColor.from_string(DEFAULT_THEME_COLOR)

    # === SPEC 0030 主题映射 ===

    def _map_theme(self, theme_color: str | None) -> str:
        """theme_color hex → pptxforge 主题名映射（SPEC 0030 方案 B fallback）。

        当 theme_preset 为 None 时，由 theme_color 映射到 pptxforge 主题。
        映射规则基于 HLS 色相（H）和饱和度（S）：
        - 低饱和度（S < 0.20）→ 灰色系 → MONOCHROME_INK
        - 红色（H < 0.05 或 H > 0.95）→ CORAL_ENERGY
        - 橙色（0.05 ≤ H < 0.10）→ AMBER_EDITORIAL
        - 黄色（0.10 ≤ H < 0.18）→ SUNRISE_CITRUS
        - 绿色（0.18 ≤ H < 0.45）→ FOREST_MOSS
        - 青色（0.45 ≤ H < 0.55）→ PACIFIC_DEEP
        - 蓝色（0.55 ≤ H < 0.70）→ MIDNIGHT_EXECUTIVE
        - 紫色（0.70 ≤ H < 0.90）→ ROYAL_PLUM
        - 其他 → SLATE_MINIMALIST（默认）

        None 或无效色值 → SLATE_MINIMALIST（默认降级策略）。
        """
        if not theme_color:
            return "SLATE_MINIMALIST"

        try:
            hex_str = theme_color.lstrip("#")
            if len(hex_str) != 6:
                return "SLATE_MINIMALIST"
            r = int(hex_str[0:2], 16) / 255
            g = int(hex_str[2:4], 16) / 255
            b = int(hex_str[4:6], 16) / 255
            h, _l, s = colorsys.rgb_to_hls(r, g, b)
        except Exception:
            return "SLATE_MINIMALIST"

        # 低饱和度 → 灰色系
        if s < 0.20:
            return "MONOCHROME_INK"

        # 按色相映射
        if h < 0.05 or h > 0.95:
            return "CORAL_ENERGY"        # 红色
        elif h < 0.10:
            return "AMBER_EDITORIAL"     # 橙色
        elif h < 0.18:
            return "SUNRISE_CITRUS"      # 黄色
        elif h < 0.45:
            return "FOREST_MOSS"         # 绿色
        elif h < 0.55:
            return "PACIFIC_DEEP"        # 青色
        elif h < 0.70:
            return "MIDNIGHT_EXECUTIVE"  # 蓝色
        elif h < 0.90:
            return "ROYAL_PLUM"          # 紫色
        else:
            return "SLATE_MINIMALIST"    # 默认

    # === SPEC 0025 三角色彩派生 ===

    def _derive_color_palette(
        self, theme_rgb: RGBColor,
    ) -> tuple[RGBColor, RGBColor, RGBColor, RGBColor]:
        """从主题色派生三角色彩（SPEC 0025）。

        使用 Python 标准库 colorsys 进行 RGB ↔ HLS 色彩空间转换。

        返回 (primary, auxiliary, accent, title_text_color)：
        - primary: 主色（= theme_rgb），用于标题栏/页脚栏背景、封面色块
        - auxiliary: 辅助色（高亮度低饱和度浅色），用于左栏背景衬托
        - accent: 强调色（互补色相），用于圆点标记/章节标题
        - title_text_color: 标题栏文字色（白或深灰，取决于主色亮度）
        """
        # 提取 RGB 分量（0-1 范围）
        hex_str = str(theme_rgb)
        r = int(hex_str[0:2], 16) / 255
        g = int(hex_str[2:4], 16) / 255
        b = int(hex_str[4:6], 16) / 255

        # 转 HLS
        h, l, s = colorsys.rgb_to_hls(r, g, b)

        # 主色 = 原色
        primary = theme_rgb

        # 辅助色 = 高亮度（0.92）低饱和度（≤0.30），浅色背景衬托
        aux_l = 0.92
        aux_s = min(s, 0.30)
        aux_r, aux_g, aux_b = colorsys.hls_to_rgb(h, aux_l, aux_s)
        auxiliary = RGBColor(
            int(aux_r * 255), int(aux_g * 255), int(aux_b * 255),
        )

        # 强调色 = 互补色相（H+0.5），中亮度（0.45）高饱和度（0.50-0.70）
        if s < 0.20:
            # 低饱和度特殊处理：互补色区分度不足，固定使用蓝色作为强调色
            # 阈值 0.20 覆盖 #475569（饱和度约 0.193）等灰蓝色
            accent = RGBColor(0x25, 0x63, 0xEB)
        else:
            acc_h = (h + 0.5) % 1.0
            acc_l = 0.45
            acc_s = min(max(s, 0.50), 0.70)
            acc_r, acc_g, acc_b = colorsys.hls_to_rgb(acc_h, acc_l, acc_s)
            accent = RGBColor(
                int(acc_r * 255), int(acc_g * 255), int(acc_b * 255),
            )

        # 标题栏文字色：主色亮度 > 0.60 时用深灰，否则用白色（对比度保障）
        # 阈值 0.60 覆盖紫色 #7c3aed（L≈0.578），使 5 种预设深色统一用白字
        if l > 0.60:
            title_text_color = TEXT_COLOR_DARK
        else:
            title_text_color = TEXT_COLOR_WHITE

        return primary, auxiliary, accent, title_text_color

    # === SPEC 0026 视觉效果增强 ===

    def _darken_color(
        self, rgb: RGBColor, factor: float = 0.20,
    ) -> RGBColor:
        """降低颜色亮度（SPEC 0026 渐变派生）。

        在 HLS 空间将 L 降低 factor，下限 0.10 保护极暗颜色。
        用于渐变填充的结束色派生。
        """
        hex_str = str(rgb)
        r = int(hex_str[0:2], 16) / 255
        g = int(hex_str[2:4], 16) / 255
        b = int(hex_str[4:6], 16) / 255
        h, l, s = colorsys.rgb_to_hls(r, g, b)
        # 亮度降低 factor，下限 0.10 保护
        l = max(l - factor, 0.10)
        dr, dg, db = colorsys.hls_to_rgb(h, l, s)
        return RGBColor(int(dr * 255), int(dg * 255), int(db * 255))

    # === 内容页候选构建 ===

    def _build_content_groups(
        self,
        outline_sections: list[dict],
        chart_artifacts: list[dict],
    ) -> list[dict]:
        """按共享布局规划器关联章节与图表，不再按列表顺序猜配。"""

        plans = plan_section_layouts(outline_sections, chart_artifacts)
        return [
            {
                "title": plan.title,
                "sections": [plan.section],
                "chart": plan.chart_artifacts[0] if plan.chart_artifacts else None,
            }
            for plan in plans
            if plan.section.get("source_type") != "SUMMARY"
        ]
    def _control_slide_count(
        self,
        content_groups: list[dict],
        target_slide_count: int,
    ) -> list[dict]:
        """页数控制（SPEC 0011 逻辑，不改变）。

        target_slide_count 指定时，内容页数不超过 target-2（减去标题页和总结页）。
        """
        available_slots = max(0, target_slide_count - 2)
        if available_slots == 0:
            return []
        if len(content_groups) <= available_slots:
            return content_groups

        # 合并所有内容到 available_slots 个页面
        all_sections: list[dict] = []
        all_titles: list[str] = []
        all_charts: list[dict | None] = []
        for group in content_groups:
            all_titles.append(group["title"])
            all_sections.extend(group["sections"])
            if group["chart"] is not None:
                all_charts.append(group["chart"])

        merged_title = "、".join(all_titles)
        max_items = available_slots * 5
        all_sections = all_sections[:max_items]

        result: list[dict] = []
        chart_idx = 0
        for i in range(available_slots):
            start = i * 5
            end = start + 5
            chunk = all_sections[start:end]
            if chunk:
                chart = (
                    all_charts[chart_idx]
                    if chart_idx < len(all_charts) else None
                )
                if chart is not None:
                    chart_idx += 1
                result.append({
                    "title": merged_title,
                    "sections": chunk,
                    "chart": chart,
                })
        return result

    # === 封面页 ===

    def _render_title_slide(
        self,
        prs: Presentation,
        project_name: str,
        project_topic: str,
        primary: RGBColor,
        title_text_color: RGBColor,
    ) -> None:
        """渲染封面页：顶部渐变色块 + 白色大标题 + 副标题 + 底部主色窄条（三明治）。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

        # 1. 顶部渐变色块（SPEC 0026：主色 → 主色暗化 20%）
        primary_dark_20 = self._darken_color(primary, 0.20)
        self._add_gradient_block(
            slide, Inches(0), Inches(0),
            Inches(SLIDE_WIDTH), Inches(TITLE_BANNER_HEIGHT),
            primary, primary_dark_20, angle_deg=90,
        )

        # 2. 主标题（色块内白色文字，36pt Bold，居中）
        title_tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(0.7),
            Inches(SLIDE_WIDTH - 2 * MARGIN_LEFT), Inches(1.2),
        )
        tf = title_tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = project_topic or "实验报告"
        self._set_run_font(
            run, FONT_SIZE_MAIN_TITLE, title_text_color, bold=True,
        )

        # 3. 副标题（色块下方，20pt 主色）
        subtitle_tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(3.2),
            Inches(SLIDE_WIDTH - 2 * MARGIN_LEFT), Inches(1.5),
        )
        stf = subtitle_tb.text_frame
        stf.word_wrap = True
        now = datetime.now(timezone.utc)
        p1 = stf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = f"项目：{project_name}"
        self._set_run_font(r1, FONT_SIZE_SUBTITLE, primary)

        p2 = stf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = f"生成日期：{now.strftime('%Y-%m-%d')}"
        self._set_run_font(r2, FONT_SIZE_SUBTITLE, primary)

        # 4. 底部主色窄条（SPEC 0025 新增：三明治下层面包）
        self._add_color_block(
            slide,
            Inches(0), Inches(FOOTER_BAR_TOP),
            Inches(SLIDE_WIDTH), Inches(FOOTER_BAR_HEIGHT),
            primary,
        )

    # === 双栏内容页 ===

    def _add_content_slide(
        self,
        prs: Presentation,
        title: str,
        sections: list[dict],
        primary: RGBColor,
        auxiliary: RGBColor,
        accent: RGBColor,
        title_text_color: RGBColor,
        chart_artifact: dict | None = None,
        page_num: int = 0,
        total_pages: int = 0,
        project_name: str = "",
    ) -> None:
        """添加双栏内容页：深色标题栏 + 左栏文本（辅助色背景）+ 右栏图表 + 深色页脚栏。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

        # 1. 深色标题栏（SPEC 0025 三明治上层面包）
        self._add_page_title(slide, title, primary, title_text_color)

        # 2. 左栏文本（40%，辅助色背景 + 强调色圆点）
        self._add_content_left_column(slide, sections, accent, auxiliary)

        # 3. 右栏图表或补充文本（60%）
        if chart_artifact:
            self._add_content_right_chart(slide, chart_artifact, auxiliary)
        else:
            self._add_content_right_text(slide, sections)

        # 4. 深色页脚栏（SPEC 0025 三明治下层面包）
        self._add_footer(
            slide, project_name, page_num, total_pages,
            primary, title_text_color,
        )

    def _add_page_title(
        self,
        slide,
        title: str,
        primary: RGBColor,
        title_text_color: RGBColor,
        font_size: int = FONT_SIZE_PAGE_TITLE,
    ) -> None:
        """添加深色标题栏（渐变背景 + 白色/深灰标题文字，SPEC 0025/0026）。"""
        # 1. 渐变背景栏（SPEC 0026：主色 → 主色暗化 15%，全幅，高 TITLE_BAR_HEIGHT）
        primary_dark_15 = self._darken_color(primary, 0.15)
        self._add_gradient_block(
            slide,
            Inches(0), Inches(0),
            Inches(SLIDE_WIDTH), Inches(TITLE_BAR_HEIGHT),
            primary, primary_dark_15, angle_deg=90,
        )

        # 2. 标题文字（title_text_color，28pt Bold，垂直居中）
        title_tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(0.1),
            Inches(SLIDE_WIDTH - 2 * MARGIN_LEFT), Inches(TITLE_BAR_HEIGHT - 0.2),
        )
        tf = title_tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        self._set_run_font(
            run, font_size, title_text_color, bold=True,
        )

    def _add_content_left_column(
        self,
        slide,
        sections: list[dict],
        accent: RGBColor,
        auxiliary: RGBColor,
    ) -> None:
        """添加左栏文本要点（辅助色圆角背景 + 强调色圆点/标题 + 深灰说明，16pt）。"""
        # 1. 辅助色浅色圆角背景（SPEC 0026：圆角矩形替代直角矩形）
        self._add_rounded_color_block(
            slide,
            Inches(MARGIN_LEFT), Inches(CONTENT_TOP),
            Inches(CONTENT_LEFT_WIDTH), Inches(CONTENT_HEIGHT),
            auxiliary, corner_radius=0.05,
        )

        # 2. 文本要点（叠在背景上，留 0.1" 内边距）
        tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT + 0.15), Inches(CONTENT_TOP + 0.1),
            Inches(CONTENT_LEFT_WIDTH - 0.3), Inches(CONTENT_HEIGHT - 0.2),
        )
        tf = tb.text_frame
        tf.word_wrap = True

        for i, section in enumerate(sections[:5]):  # 最多 5 个要点
            content = section.get("content", "")
            short_content = content[:500] + ("…" if len(content) > 500 else "")
            section_title = section.get("title", "")

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 圆点标记（强调色，SPEC 0025）
            bullet_run = p.add_run()
            bullet_run.text = "● "
            self._set_run_font(
                bullet_run, FONT_SIZE_BODY, accent, bold=True,
            )

            # 章节标题（强调色 Bold，SPEC 0025）
            title_run = p.add_run()
            title_run.text = section_title
            self._set_run_font(
                title_run, FONT_SIZE_BODY, accent, bold=True,
            )

            # 说明文本（深灰，换行）
            if short_content:
                desc_p = tf.add_paragraph()
                desc_run = desc_p.add_run()
                desc_run.text = f"  {short_content}"
                self._set_run_font(
                    desc_run, FONT_SIZE_BODY, TEXT_COLOR_DARK,
                )

        # 超过 5 个要点加省略号
        if len(sections) > 5:
            more_p = tf.add_paragraph()
            more_run = more_p.add_run()
            more_run.text = "…"
            self._set_run_font(
                more_run, FONT_SIZE_BODY, accent,
            )

    def _add_content_right_chart(
        self, slide, artifact: dict, auxiliary: RGBColor,
    ) -> None:
        """添加右栏图表（自适应缩放 + 边框 + 阴影，SPEC 0026）。"""
        file_path = artifact.get("file_path", "")
        name = artifact.get("name", "")

        # 可用区域：top=CONTENT_TOP 到 FOOTER_BAR_TOP
        max_width = CONTENT_RIGHT_WIDTH  # 6.7"
        max_height = FOOTER_BAR_TOP - CONTENT_TOP  # 5.4"

        if file_path and Path(file_path).exists():
            w, h = self._fit_image_size(file_path, max_width, max_height)
            try:
                pic = slide.shapes.add_picture(
                    str(file_path),
                    Inches(CONTENT_RIGHT_LEFT), Inches(CONTENT_TOP),
                    width=Inches(w), height=Inches(h),
                )
                self._style_picture(pic, auxiliary)
            except Exception:
                self._add_placeholder_textbox(
                    slide, CONTENT_RIGHT_LEFT, CONTENT_TOP,
                    max_width, CONTENT_HEIGHT,
                    f"[图片无法嵌入：{name}]",
                )
        else:
            self._add_placeholder_textbox(
                slide, CONTENT_RIGHT_LEFT, CONTENT_TOP,
                max_width, CONTENT_HEIGHT,
                f"[图片文件不存在：{name}]",
            )

    def _add_content_right_text(
        self,
        slide,
        sections: list[dict],
    ) -> None:
        """无图表时右栏展示补充文本。"""
        tb = slide.shapes.add_textbox(
            Inches(CONTENT_RIGHT_LEFT), Inches(CONTENT_TOP),
            Inches(CONTENT_RIGHT_WIDTH), Inches(CONTENT_HEIGHT),
        )
        tf = tb.text_frame
        tf.word_wrap = True

        # 取所有章节的完整内容作为补充文本
        for i, section in enumerate(sections[:5]):
            content = section.get("content", "")
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = content[:500] + ("…" if len(content) > 500 else "")
            self._set_run_font(
                run, FONT_SIZE_BODY, RGBColor(0x55, 0x55, 0x55),
            )

    # === 图表自适应布局 ===

    def _add_chart_slide(
        self,
        prs: Presentation,
        chart_artifacts: list[dict],
        primary: RGBColor,
        title_text_color: RGBColor,
        page_num: int = 0,
        total_pages: int = 0,
        project_name: str = "",
    ) -> None:
        """添加关键图表页（图表自适应布局 + SPEC 0025 三明治结构）。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
        self._add_page_title(slide, "关键图表", primary, title_text_color)

        count = min(len(chart_artifacts), 4)  # 最多 4 张
        if count == 1:
            self._place_chart_centered(slide, chart_artifacts[0])
        elif count == 2:
            self._place_chart_side_by_side(slide, chart_artifacts[:2])
        elif count == 3:
            self._place_chart_three(slide, chart_artifacts[:3])
        else:
            self._place_chart_grid(slide, chart_artifacts[:4])

        # 补充图表页与正文图表页复用同一套边框/阴影组件。
        _, auxiliary, _, _ = self._derive_color_palette(primary)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._style_picture(shape, auxiliary)

        if len(chart_artifacts) > 4:
            self._add_truncation_note(slide, len(chart_artifacts))

        self._add_footer(
            slide, project_name, page_num, total_pages,
            primary, title_text_color,
        )

    def _fit_image_size(
        self, file_path: str, max_width: float, max_height: float,
    ) -> tuple[float, float]:
        """按宽高比缩放图片，确保不超过最大宽度和高度。

        返回 (width, height) 英寸。同时设置 width 和 height 不会拉伸图片，
        因为两者按原始宽高比等比计算。
        """
        try:
            from PIL import Image

            with Image.open(file_path) as img:
                orig_w, orig_h = img.size
        except Exception:
            # 无法读取尺寸时，用默认 5:3 比例
            return max_width, min(max_width * 0.6, max_height)

        ratio = orig_w / orig_h
        # 先按最大宽度缩放
        w = max_width
        h = w / ratio
        # 如果高度超出，改为按最大高度缩放
        if h > max_height:
            h = max_height
            w = h * ratio
        return round(w, 2), round(h, 2)

    def _place_chart_centered(self, slide, artifact: dict) -> None:
        """单图居中放大布局（自适应缩放到可用区域）。"""
        file_path = artifact.get("file_path", "")
        name = artifact.get("name", "")

        # 可用区域：top=1.8 到 FOOTER_BAR_TOP=7.0
        max_width = 8.0
        max_height = FOOTER_BAR_TOP - 1.8  # 5.2"

        if file_path and Path(file_path).exists():
            w, h = self._fit_image_size(file_path, max_width, max_height)
            left = (SLIDE_WIDTH - w) / 2  # 水平居中
            try:
                slide.shapes.add_picture(
                    str(file_path),
                    Inches(left), Inches(1.8),
                    width=Inches(w), height=Inches(h),
                )
            except Exception:
                self._add_placeholder_textbox(
                    slide, left, 1.8, max_width, 4,
                    f"[图片无法嵌入：{name}]",
                )
        else:
            left = (SLIDE_WIDTH - max_width) / 2
            self._add_placeholder_textbox(
                slide, left, 1.8, max_width, 4,
                f"[图片文件不存在：{name}]",
            )

    def _place_chart_side_by_side(
        self, slide, artifacts: list[dict],
    ) -> None:
        """双图左右并排布局（各自适应缩放到可用区域）。

        SPEC 0027：使用 _GridHelper 计算 1×2 网格坐标，替代硬编码 positions。
        参数选择确保与原硬编码坐标完全一致（无视觉漂移）。
        """
        max_height = FOOTER_BAR_TOP - 1.8  # 5.2"
        # SPEC 0027：用 _GridHelper 计算 1×2 网格坐标
        # 验证：cell_w=(12.3-0.7)/2=5.8, cell(0,0)=(0.5,1.8), cell(0,1)=(7.0,1.8)
        grid = self._GridHelper(
            left=Inches(0.5), top=Inches(1.8),
            width=Inches(12.3), height=Inches(max_height),
            rows=1, cols=2,
            h_gap=Inches(0.7), v_gap=0,
        )
        for i, art in enumerate(artifacts[:2]):
            file_path = art.get("file_path", "")
            name = art.get("name", "")
            cell_left, cell_top, cell_w, _ = grid.cell(0, i)
            left = Emu(cell_left).inches
            top = Emu(cell_top).inches
            max_width = Emu(cell_w).inches

            if file_path and Path(file_path).exists():
                w, h = self._fit_image_size(file_path, max_width, max_height)
                try:
                    slide.shapes.add_picture(
                        str(file_path),
                        Inches(left), Inches(top),
                        width=Inches(w), height=Inches(h),
                    )
                except Exception:
                    self._add_placeholder_textbox(
                        slide, left, top, max_width, 3.5,
                        f"[图片无法嵌入：{name}]",
                    )
            else:
                self._add_placeholder_textbox(
                    slide, left, top, max_width, 3.5,
                    f"[图片文件不存在：{name}]",
                )

    def _place_chart_three(self, slide, artifacts: list[dict]) -> None:
        """3 张图布局：上排 2 张并排 + 下排 1 张居中。

        布局约束（SPEC 0024）：
        - 上排 top=1.5, 下排 top=4.0, 行高 max 2.3"
        - 下排底部最大 6.3"，为页脚（top=7.0）预留空间

        SPEC 0027：上排使用 _GridHelper 计算 1×2 网格坐标，替代硬编码 top_positions。
        下排保持居中逻辑（动态 left，单图居中不进入网格）。
        """
        max_width_bottom = 8.0
        max_height = 2.3
        # SPEC 0027：上排用 _GridHelper 计算 1×2 网格坐标
        # 验证：cell_w=(12.3-0.7)/2=5.8, cell(0,0)=(0.5,1.5), cell(0,1)=(7.0,1.5)
        grid_top = self._GridHelper(
            left=Inches(0.5), top=Inches(1.5),
            width=Inches(12.3), height=Inches(max_height),
            rows=1, cols=2,
            h_gap=Inches(0.7), v_gap=0,
        )

        # 上排 2 张
        for i, art in enumerate(artifacts[:2]):
            file_path = art.get("file_path", "")
            name = art.get("name", "")
            cell_left, cell_top, cell_w, _ = grid_top.cell(0, i)
            left = Emu(cell_left).inches
            top = Emu(cell_top).inches
            max_width = Emu(cell_w).inches
            if file_path and Path(file_path).exists():
                w, h = self._fit_image_size(file_path, max_width, max_height)
                try:
                    slide.shapes.add_picture(
                        str(file_path),
                        Inches(left), Inches(top),
                        width=Inches(w), height=Inches(h),
                    )
                except Exception:
                    self._add_placeholder_textbox(
                        slide, left, top, max_width, max_height,
                        f"[图片无法嵌入：{name}]",
                    )
            else:
                self._add_placeholder_textbox(
                    slide, left, top, max_width, max_height,
                    f"[图片文件不存在：{name}]",
                )

        # 下排 1 张居中
        if len(artifacts) >= 3:
            art = artifacts[2]
            file_path = art.get("file_path", "")
            name = art.get("name", "")
            top = 4.0
            if file_path and Path(file_path).exists():
                w, h = self._fit_image_size(file_path, max_width_bottom, max_height)
                left = (SLIDE_WIDTH - w) / 2  # 水平居中
                try:
                    slide.shapes.add_picture(
                        str(file_path),
                        Inches(left), Inches(top),
                        width=Inches(w), height=Inches(h),
                    )
                except Exception:
                    self._add_placeholder_textbox(
                        slide, left, top, max_width_bottom, max_height,
                        f"[图片无法嵌入：{name}]",
                    )
            else:
                left = (SLIDE_WIDTH - max_width_bottom) / 2
                self._add_placeholder_textbox(
                    slide, left, top, max_width_bottom, max_height,
                    f"[图片文件不存在：{name}]",
                )

    def _place_chart_grid(self, slide, artifacts: list[dict]) -> None:
        """多图 2×2 网格布局（各自适应缩放到网格单元）。

        布局约束（SPEC 0024）：
        - 上排 top=1.5, 下排 top=4.0, 行高 max 2.3"
        - 下排底部最大 6.3"，为截断注释（top=6.5）和页脚（top=7.0）预留空间
        - 避免图片与截断注释/页脚重叠

        SPEC 0027：使用 _GridHelper 计算 2×2 网格坐标，替代硬编码 positions。
        参数选择确保与原硬编码坐标完全一致（无视觉漂移）。
        """
        # SPEC 0027：用 _GridHelper 计算 2×2 网格坐标
        # 验证：cell_w=(9.9-2.3)/2=3.8, cell_h=(4.8-0.2)/2=2.3
        # cell(0,0)=(0.7,1.5), cell(0,1)=(6.8,1.5), cell(1,0)=(0.7,4.0), cell(1,1)=(6.8,4.0)
        grid = self._GridHelper(
            left=Inches(0.7), top=Inches(1.5),
            width=Inches(9.9), height=Inches(4.8),
            rows=2, cols=2,
            h_gap=Inches(2.3), v_gap=Inches(0.2),
        )
        for i, art in enumerate(artifacts[:4]):
            file_path = art.get("file_path", "")
            name = art.get("name", "")
            row, col = divmod(i, 2)
            cell_left, cell_top, cell_w, cell_h = grid.cell(row, col)
            left = Emu(cell_left).inches
            top = Emu(cell_top).inches
            max_width = Emu(cell_w).inches
            max_height = Emu(cell_h).inches

            if file_path and Path(file_path).exists():
                w, h = self._fit_image_size(file_path, max_width, max_height)
                try:
                    slide.shapes.add_picture(
                        str(file_path),
                        Inches(left), Inches(top),
                        width=Inches(w), height=Inches(h),
                    )
                except Exception:
                    self._add_placeholder_textbox(
                        slide, left, top, max_width, max_height,
                        f"[图片无法嵌入：{name}]",
                    )
            else:
                self._add_placeholder_textbox(
                    slide, left, top, max_width, max_height,
                    f"[图片文件不存在：{name}]",
                )

    def _add_truncation_note(self, slide, total_count: int) -> None:
        """添加截断注释（超过 4 张图表时）。"""
        tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(6.5),
            Inches(SLIDE_WIDTH - 2 * MARGIN_LEFT), Inches(0.4),
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"共 {total_count} 张图表，已展示前 4 张"
        self._set_run_font(
            run, FONT_SIZE_CAPTION, RGBColor(0x88, 0x88, 0x88),
        )

    # === 总结页 ===

    def _render_summary_slide(
        self,
        prs: Presentation,
        outline_sections: list[dict],
        primary: RGBColor,
        title_text_color: RGBColor,
        page_num: int = 0,
        total_pages: int = 0,
        project_name: str = "",
    ) -> None:
        """渲染总结页：深色标题栏 + 居中正文 + 深色页脚栏（SPEC 0025 三明治结构）。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

        # 1. 深色标题栏（SPEC 0025 三明治上层面包）
        self._add_page_title(slide, "总结", primary, title_text_color)

        # 2. 提取 SUMMARY 类型章节作为总结
        summary_sections = [
            s for s in outline_sections
            if s.get("source_type") == "SUMMARY"
        ]

        # 3. 总结正文（居中，20pt 深灰）
        tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(2.5),
            Inches(SLIDE_WIDTH - 2 * MARGIN_LEFT), Inches(3.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        if summary_sections:
            for i, section in enumerate(summary_sections):
                content = section.get("content", "")
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = content
                self._set_run_font(
                    run, FONT_SIZE_SUBTITLE, TEXT_COLOR_DARK,
                )
        else:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "本实验已按既定方案完成数据分析与可视化。"
            self._set_run_font(
                run, FONT_SIZE_SUBTITLE, TEXT_COLOR_DARK,
            )

        # 4. 深色页脚栏（SPEC 0025 三明治下层面包，取代旧装饰线）
        self._add_footer(
            slide, project_name, page_num, total_pages,
            primary, title_text_color,
        )

    # === 辅助方法：样式 ===

    def _set_run_font(
        self,
        run,
        size: int,
        color_rgb: RGBColor | None = None,
        bold: bool = False,
        font_name: str = FONT_NAME_CN,
    ) -> None:
        """统一设置 run 的字体样式（含东亚字体）。"""
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        if color_rgb is not None:
            run.font.color.rgb = color_rgb
        # 设置东亚字体（确保中文正确渲染）
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", font_name)

    def _add_color_block(
        self,
        slide,
        left,
        top,
        width,
        height,
        color_rgb: RGBColor,
    ) -> None:
        """添加纯色色块（无边框矩形）。"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_rgb
        shape.line.fill.background()  # 无边框

    def _add_gradient_block(
        self,
        slide,
        left,
        top,
        width,
        height,
        color_start: RGBColor,
        color_end: RGBColor,
        angle_deg: float = 90,
    ) -> None:
        """添加渐变色块（线性渐变，SPEC 0026）。

        使用 python-pptx 原生 fill.gradient() API。
        参数：
        - color_start: 起始色（position=0.0）
        - color_end: 结束色（position=1.0）
        - angle_deg: 渐变角度（90=上→下）
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = angle_deg
        # 设置两端颜色（默认 gradient() 生成两停止点）
        stops = fill.gradient_stops
        stops[0].color.rgb = color_start
        stops[0].position = 0.0
        stops[1].color.rgb = color_end
        stops[1].position = 1.0
        shape.line.fill.background()  # 无边框

    def _add_rounded_color_block(
        self,
        slide,
        left,
        top,
        width,
        height,
        color_rgb: RGBColor,
        corner_radius: float = 0.05,
    ) -> None:
        """添加圆角色块（圆角矩形 + 纯色填充，SPEC 0026）。

        参数：
        - corner_radius: 圆角半径（0.0-1.0，相对短边比例）
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.adjustments[0] = corner_radius
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_rgb
        shape.line.fill.background()  # 无边框

    def _style_picture(self, picture, auxiliary: RGBColor) -> None:
        # 复用 SPEC 0026 的细边框和外阴影科研图表组件。
        picture.line.color.rgb = auxiliary
        picture.line.width = Pt(1)
        self._add_picture_shadow(picture)

    def _add_picture_shadow(
        self,
        picture,
        blur_radius_pt: float = 8,
        distance_pt: float = 4,
        direction_deg: float = 315,
        alpha_pct: int = 30,
    ) -> None:
        """为图片添加外阴影效果（oxml 操作 a:effectLst，SPEC 0026）。

        python-pptx 未暴露 effect_format API，需直接操作 oxml。
        参数：
        - blur_radius_pt: 模糊半径（Pt）
        - distance_pt: 阴影距离（Pt）
        - direction_deg: 阴影方向（度，0=右，90=下，180=左，270=上）
        - alpha_pct: 不透明度百分比（0-100）
        """
        # EMU 转换常量：1 Pt = 12700 EMU；1 度 = 60000（1/60000 度）
        blur_emu = str(int(blur_radius_pt * 12700))
        dist_emu = str(int(distance_pt * 12700))
        dir_emu = str(int(direction_deg * 60000))
        # alpha 百分比 → 千分比（100% = 100000）
        alpha_val = str(int(alpha_pct * 1000))

        spPr = picture._element.spPr
        # 移除已有 effectLst（避免重复）
        existing = spPr.find(qn('a:effectLst'))
        if existing is not None:
            spPr.remove(existing)

        effect_lst = spPr.makeelement(qn('a:effectLst'), {})
        outer_shdw = effect_lst.makeelement(qn('a:outerShdw'), {
            'blurRad': blur_emu,
            'dist': dist_emu,
            'dir': dir_emu,
            'rotWithShape': '0',
        })
        clr = outer_shdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': alpha_val})
        clr.append(alpha)
        outer_shdw.append(clr)
        effect_lst.append(outer_shdw)
        spPr.append(effect_lst)

    def _add_divider(
        self,
        slide,
        left,
        top,
        width,
        color_rgb: RGBColor,
        height: float = DIVIDER_HEIGHT,
    ) -> None:
        """添加分隔线（细矩形）。"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_rgb
        shape.line.fill.background()

    def _add_footer(
        self,
        slide,
        project_name: str,
        page_num: int,
        total_pages: int,
        primary: RGBColor,
        title_text_color: RGBColor,
    ) -> None:
        """添加深色页脚栏（渐变背景 + 白色项目名和页码，SPEC 0025/0026）。"""
        # 1. 渐变背景栏（SPEC 0026：主色暗化 15% → 主色，全幅，高 FOOTER_BAR_HEIGHT）
        primary_dark_15 = self._darken_color(primary, 0.15)
        self._add_gradient_block(
            slide,
            Inches(0), Inches(FOOTER_BAR_TOP),
            Inches(SLIDE_WIDTH), Inches(FOOTER_BAR_HEIGHT),
            primary_dark_15, primary, angle_deg=90,
        )

        # 2. 项目名（左，title_text_color）
        left_tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(FOOTER_BAR_TOP + 0.1),
            Inches(6), Inches(0.3),
        )
        lt = left_tb.text_frame
        lr = lt.paragraphs[0].add_run()
        lr.text = project_name
        self._set_run_font(
            lr, FONT_SIZE_CAPTION, title_text_color,
        )

        # 3. 页码（右，title_text_color）
        right_tb = slide.shapes.add_textbox(
            Inches(SLIDE_WIDTH - 3), Inches(FOOTER_BAR_TOP + 0.1),
            Inches(2.5), Inches(0.3),
        )
        rt = right_tb.text_frame
        rt.paragraphs[0].alignment = PP_ALIGN.RIGHT
        rr = rt.paragraphs[0].add_run()
        rr.text = f"第 {page_num} / {total_pages} 页"
        self._set_run_font(
            rr, FONT_SIZE_CAPTION, title_text_color,
        )

    def _add_placeholder_textbox(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
    ) -> None:
        """添加占位文本框（图片嵌入失败时使用）。"""
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        self._set_run_font(
            run, FONT_SIZE_BODY, RGBColor(0x88, 0x88, 0x88),
        )
